"""
common.py — Shared module for the SEGA Shooter Market Intelligence app.

This module holds everything every page needs: constants, data fetchers,
prompt builders, PPTX/HTML/PDF export, translations, auth, and shared UI
chrome (topbar/footer/CSS). It is imported by the main entry script and by
every file under pages/.

Streamlit's multipage model runs each page as an independent script — there
is no guarantee the main file executes before a page does. Every page must
therefore call, in this order, at the very top of its script:

    st.set_page_config(...)      # must be the first Streamlit call, per-page
    inject_css()
    init_session_defaults()
    require_auth()                # st.stop()s here if not authenticated
    render_topbar()

Optional extras (snapshot PowerPoint export):
    pip install playwright pillow kaleido
    playwright install chromium
"""

import sys
import time
import re
import io
import json
import random
import base64
import hmac
import hashlib
import concurrent.futures
import threading
from datetime import datetime, timedelta
from collections import Counter
from pathlib import Path
import os

import requests
import pandas as pd
import plotly.graph_objects as go
import streamlit as st
import streamlit.components.v1 as _st_components

try:
    import markdown as _md_lib
    MARKDOWN_AVAILABLE = True
except ImportError:
    MARKDOWN_AVAILABLE = False

try:
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import cm
    from reportlab.lib import colors as _rl_colors
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, HRFlowable, Preformatted
    import io as _rl_io
    _REPORTLAB_AVAILABLE = True
except ImportError:
    _REPORTLAB_AVAILABLE = False

try:
    import anthropic as _anthropic
    ANTHROPIC_AVAILABLE = True
except ImportError:
    ANTHROPIC_AVAILABLE = False

# ─────────────────────────────────────────────────────────────
# HTML TABLE HELPER
# ─────────────────────────────────────────────────────────────

def parse_md_table(block: str) -> tuple[list[str], list[dict]] | None:
    """Parse a markdown table block into (headers, rows).
    Returns None if the block isn't a valid markdown table."""
    lines = [l.strip() for l in block.strip().splitlines() if l.strip()]
    if len(lines) < 3:
        return None
    # Must have a separator line (---|---|---)
    sep_idx = next((i for i, l in enumerate(lines)
                    if re.match(r"^\|?[\s\-:]+(\|[\s\-:]+)+\|?$", l)), None)
    if sep_idx is None or sep_idx == 0:
        return None
    def clean_cell(text: str) -> str:
        """Strip markdown bold/italic markers from cell text."""
        t = text.strip()
        t = re.sub(r"\*\*(.+?)\*\*", r"\1", t)
        t = re.sub(r"\*(.+?)\*",     r"\1", t)
        t = re.sub(r"__(.+?)__",       r"\1", t)
        t = re.sub(r"_(.+?)_",         r"\1", t)
        return t.strip()

    def split_row(line):
        return [clean_cell(c) for c in line.strip().strip("|").split("|")]
    headers = split_row(lines[sep_idx - 1])
    rows = []
    for line in lines[sep_idx + 1:]:
        if not line.startswith("|") and "|" not in line:
            break
        cells = split_row(line)
        # Pad or trim to match header count
        while len(cells) < len(headers):
            cells.append("")
        rows.append(dict(zip(headers, cells[:len(headers)])))
    return (headers, rows) if rows else None


def render_report_with_tables(report_md: str) -> None:
    """Render an AI markdown report, replacing any markdown tables
    with sortable HTML tables via render_table()."""
    # Split report into chunks: text blocks and table blocks
    # A table block is a sequence of lines containing |
    lines = report_md.splitlines(keepends=True)
    chunks = []
    i = 0
    while i < len(lines):
        line = lines[i]
        # Detect start of a markdown table (line with | characters)
        if "|" in line and line.strip().startswith("|"):
            table_lines = []
            while i < len(lines) and ("|" in lines[i] or lines[i].strip() == ""):
                if "|" in lines[i]:
                    table_lines.append(lines[i])
                    i += 1
                else:
                    break
            block = "".join(table_lines)
            parsed = parse_md_table(block)
            if parsed:
                chunks.append(("table", parsed))
            else:
                chunks.append(("text", block))
        else:
            # Accumulate text until next table
            text = ""
            while i < len(lines) and not ("|" in lines[i] and lines[i].strip().startswith("|")):
                text += lines[i]
                i += 1
            if text:
                chunks.append(("text", text))

    for kind, content in chunks:
        if kind == "text":
            if content.strip():
                st.markdown(content)
        else:
            headers, rows = content
            render_table(rows, headers)


def render_table(rows: list[dict], col_order: list[str] = None,
                 green_cols: set = None, red_cols: set = None,
                 height: int = None) -> None:
    """Render html_table inside a components iframe so <script> executes."""
    html = html_table(rows, col_order, green_cols, red_cols)
    if not html:
        return
    # Wrap in a full dark-themed page so the iframe matches the app
    full = f"""<!DOCTYPE html><html>
<head><style>
  * {{ margin:0;padding:0;box-sizing:border-box;font-family:Poppins,sans-serif; }}
  body {{ background:#0a0c1a;color:#eef0fa; }}
  table {{ width:100%;border-collapse:collapse; }}
  th {{ cursor:pointer; user-select:none; }}
  th:hover {{ opacity:0.8; }}
</style></head>
<body>{html}</body></html>"""
    # Auto-height: ~36px per row + 50px header + padding
    if height is None:
        height = len(rows) * 36 + 60
    _st_components.html(full, height=height, scrolling=False)

# Columns where a + prefix (even without %) means green, - means red
# Both English and Japanese column names are listed so conditional
# colouring keeps working when the UI language is switched to JP.
_DELTA_COLS = {
    "YoY", "MoM", "Annual Change", "Month Change", "Weekly Change",
    "Change (CCU)", "Review",
    "年間変化", "月間変化", "週次変化", "変化 (CCU)", "レビュー",
}

def html_table(rows: list[dict], col_order: list[str] = None,
               green_cols: set[str] = None, red_cols: set[str] = None) -> str:
    """Render a list of dicts as a fully styled dark HTML table.
    Conditional colouring rules (applied in order):
      1. Cells containing a value that starts with + → green
      2. Cells containing a value that starts with - → red
      3. Columns in green_cols override → always green
      4. Columns in red_cols override → always red
    """
    if not rows:
        return ""
    cols = col_order or list(rows[0].keys())

    BG1   = "#0f1120"
    BG2   = "#141728"
    HDR   = "#0a0c1a"
    BORD  = "#232640"
    TXT   = "#eef0fa"
    POS   = "#20c65a"
    POS_BG= "rgba(32,198,90,0.08)"
    NEG   = "#ff4d4d"
    NEG_BG= "rgba(255,77,77,0.08)"
    MUTED = "#5a5f82"

    def cell_style(col: str, val) -> tuple[str, str]:
        """Return (text_color, bg_tint) for a cell."""
        s = str(val).strip()
        # explicit column overrides
        if green_cols and col in green_cols: return POS, POS_BG
        if red_cols   and col in red_cols:   return NEG, NEG_BG
        # value-based: any delta-style column or any value starting with +/-
        is_delta = col in _DELTA_COLS
        if s.startswith("+") and s not in ("", "+"): return POS, POS_BG if is_delta else ""
        if s.startswith("-") and s not in ("", "-"): return NEG, NEG_BG if is_delta else ""
        return TXT, ""

    th = (f"padding:9px 14px;text-align:left;font-size:0.7rem;font-weight:700;"
          f"letter-spacing:0.1em;text-transform:uppercase;color:{MUTED};"
          f"background:{HDR};border-bottom:2px solid {BORD};white-space:nowrap;")
    td_base = "padding:7px 14px;font-size:0.83rem;border-bottom:1px solid {bord};white-space:nowrap;"

    head = "".join(f"<th style='{th}'>{c}</th>" for c in cols)
    body_rows = []
    for i, row in enumerate(rows):
        row_bg = BG1 if i % 2 == 0 else BG2
        cells = []
        for c in cols:
            val = row.get(c, "")
            txt_col, bg_tint = cell_style(c, val)
            bg = bg_tint if bg_tint else row_bg
            fw = "font-weight:600;" if txt_col in (POS, NEG) else ""
            td_style = td_base.format(bord=BORD) + f"background:{bg};color:{txt_col};{fw}"
            cells.append(f"<td style='{td_style}'>{val}</td>")
        body_rows.append(f"<tr>{''.join(cells)}</tr>")

    # Give each table a unique id for the sort script
    import hashlib as _hl
    tbl_id = "ht_" + _hl.md5(str(cols).encode()).hexdigest()[:8]

    sort_js = f"""<script>
(function(){{
  var tbl = document.getElementById('{tbl_id}');
  if(!tbl) return;
  var headers = tbl.querySelectorAll('th');
  var asc = {{}};
  headers.forEach(function(th, ci){{
    th.style.cursor='pointer';
    th.title='Click to sort';
    th.addEventListener('click', function(){{
      var tbody = tbl.querySelector('tbody');
      var rows  = Array.from(tbody.querySelectorAll('tr'));
      asc[ci]   = !asc[ci];
      rows.sort(function(a,b){{
        var av = a.cells[ci].innerText.trim();
        var bv = b.cells[ci].innerText.trim();
        // strip non-numeric chars for numeric sort
        var an = parseFloat(av.replace(/[^0-9.-]/g,''));
        var bn = parseFloat(bv.replace(/[^0-9.-]/g,''));
        if(!isNaN(an) && !isNaN(bn)) return asc[ci] ? an-bn : bn-an;
        return asc[ci] ? av.localeCompare(bv) : bv.localeCompare(av);
      }});
      rows.forEach(function(r){{ tbody.appendChild(r); }});
      // re-stripe rows
      rows.forEach(function(r,i){{
        r.querySelectorAll('td').forEach(function(td){{
          td.style.background = td.style.background; // keep cell bg
        }});
      }});
    }});
  }});
}})();
</script>"""

    return (f"<div style='overflow-x:auto;border:1px solid {BORD};"
            f"border-radius:8px;margin-bottom:0.5rem;'>"
            f"<table id='{tbl_id}' style='width:100%;border-collapse:collapse;"
            f"font-family:Poppins,sans-serif;'>"
            f"<thead><tr>{head}</tr></thead>"
            f"<tbody>{''.join(body_rows)}</tbody>"
            f"</table></div>{sort_js}")

# ─────────────────────────────────────────────────────────────
# SEGA BRAND STYLES
# ─────────────────────────────────────────────────────────────

def inject_css() -> None:
    """Inject the SEGA dark-theme stylesheet, including mobile breakpoints.
    Call once per page, immediately after st.set_page_config()."""
    st.markdown("""
<style>
@import url('https://fonts.googleapis.com/css2?family=Inter+Tight:wght@400;700;800;900&family=Poppins:wght@300;400;500;600&display=swap');

/* ── FORCE DARK MODE ── */
html, body, #root, #root > div,
[data-testid="stAppViewContainer"],
[data-testid="stAppViewContainer"] > section,
[data-testid="stMain"], [data-testid="stHeader"],
[data-testid="stSidebar"], [data-testid="stSidebarContent"],
[data-testid="stBottom"], .stApp, .main {
    background-color: #0a0c1a !important;
    color: #eef0fa !important;
    color-scheme: dark !important;
}
/* Kill any light-mode class Streamlit injects */
[data-testid="stAppViewContainer"][class*="light"],
.stApp[class*="light"],
html[data-theme="light"] .stApp,
html[data-theme="light"] body {
    background-color: #0a0c1a !important;
    color: #eef0fa !important;
    color-scheme: dark !important;
}
/* Dataframe dark internals */
[data-testid="stDataFrame"] > div { background: #0f1120 !important; color-scheme: dark !important; }
.dvn-scroller { background: #0f1120 !important; }
.col-header-title { color: #b8bcd4 !important; background: #141728 !important; font-weight: 600 !important; font-size: 0.75rem !important; letter-spacing: 0.05em !important; }
[data-testid="glideDataEditor"] { background: #0f1120 !important; }

/* Expander summary styling */
[data-testid="stExpander"] summary {
    font-size: 0.84rem !important; font-weight: 600 !important;
    padding: 0.7rem 1rem !important;
    background: #0f1120 !important;
}
[data-testid="stExpander"] summary:hover { background: #141728 !important; }

:root,
html[data-theme="light"],
html[data-theme="dark"],
[data-theme="light"],
[data-theme="dark"] {
    color-scheme: dark !important;
    --bg:           #0a0c1a;
    --surface:      #0f1120;
    --surface2:     #141728;
    --surface3:     #1a1e30;
    --border:       #232640;
    --border-hi:    #323760;
    --blue:         #4080ff;
    --blue-lo:      #1a3acc;
    --blue-glow:    rgba(64,128,255,0.16);
    --blue-glow-hi: rgba(64,128,255,0.32);
    --text:         #eef0fa;
    --text-dim:     #b8bcd4;
    --muted:        #5a5f82;
    --pos:          #20c65a;
    --pos-dim:      rgba(32,198,90,0.14);
    --neg:          #ff3d52;
    --neg-dim:      rgba(255,61,82,0.14);
    --amber:        #ffb938;
    --amber-dim:    rgba(255,185,56,0.14);
    --purple:       #a855f7;
    --purple-dim:   rgba(168,85,247,0.14);
}

html, body {
    background: var(--bg) !important;
    color: var(--text) !important;
    color-scheme: dark !important;
}
.stApp,
.stApp > div,
section[data-testid="stAppViewContainer"],
section[data-testid="stAppViewContainer"] > div,
div[data-testid="stMain"],
div[data-testid="stVerticalBlock"],
div[data-testid="stHorizontalBlock"],
.main .block-container,
.block-container {
    background-color: var(--bg) !important;
    color: var(--text) !important;
}

*, *::before, *::after { font-family: 'Poppins', sans-serif; box-sizing: border-box; }

p, span, div, li, td, th, label,
h1, h2, h3, h4, h5, h6,
.stMarkdown, .stMarkdown p, .stMarkdown span,
[data-testid="stText"],
[data-testid="stMarkdownContainer"],
[data-testid="stMarkdownContainer"] p,
[data-testid="stMarkdownContainer"] span,
[data-testid="stMarkdownContainer"] li,
[data-testid="stMarkdownContainer"] strong,
[data-testid="stMarkdownContainer"] em,
[class*="css"] { color: var(--text) !important; }

.stCaption, [data-testid="stCaptionContainer"],
[data-testid="stCaptionContainer"] p { color: var(--muted) !important; }

code { background: var(--surface3) !important; color: var(--blue) !important; padding: 0.1em 0.4em; border-radius: 3px; }

#MainMenu, footer, header { visibility: hidden; }
.block-container { padding: 0 2.5rem 4rem !important; max-width: 1440px !important; }

::-webkit-scrollbar { width: 5px; height: 5px; }
::-webkit-scrollbar-track { background: transparent; }
::-webkit-scrollbar-thumb { background: var(--border-hi); border-radius: 4px; }
::-webkit-scrollbar-thumb:hover { background: var(--muted); }

/*  TOP NAV  */
.topbar {
    background: var(--surface);
    border-bottom: 1px solid var(--border);
    padding: 0.8rem 2.5rem;
    margin: 0 -2.5rem 1.75rem;
    display: flex;
    align-items: center;
    gap: 1.25rem;
    position: relative;
}
.topbar::after {
    content: '';
    position: absolute;
    bottom: -1px; left: 0; right: 0; height: 1px;
    background: linear-gradient(90deg, var(--blue) 0%, rgba(64,128,255,0) 55%);
}
.topbar-logo { font-family: 'Inter Tight', sans-serif; font-size: 0.95rem; font-weight: 900; color: var(--text) !important; letter-spacing: 0.12em; text-transform: uppercase; }
.topbar-logo .seg { color: var(--blue); }
.topbar-divider { width: 1px; height: 18px; background: var(--border-hi); flex-shrink: 0; }
.topbar-label { font-size: 0.6rem; font-weight: 600; color: var(--muted) !important; letter-spacing: 0.2em; text-transform: uppercase; }
.topbar-pill { margin-left: auto; background: var(--blue-glow); border: 1px solid rgba(64,128,255,0.28); border-radius: 20px; padding: 0.18rem 0.7rem; font-size: 0.58rem; font-weight: 700; letter-spacing: 0.14em; text-transform: uppercase; color: var(--blue) !important; }

/*  NAV TABS  (cross-page navigation, since st.tabs() can't navigate
    between separate multipage-app pages — these are styled page_links) */
.nav-tabs {
    display: flex;
    gap: 0.25rem;
    border-bottom: 1px solid var(--border);
    margin: 0 0 1.5rem;
    padding-bottom: 0;
}
.nav-tabs [data-testid="stPageLink"],
.nav-tabs-active {
    flex: 1;
}
.nav-tabs [data-testid="stPageLink"] {
    border-radius: 6px 6px 0 0 !important;
}
.nav-tabs [data-testid="stPageLink"] p {
    font-family: 'Inter Tight', sans-serif !important;
    font-size: 0.72rem !important;
    font-weight: 700 !important;
    letter-spacing: 0.08em !important;
    text-transform: uppercase !important;
    color: var(--muted) !important;
}
.nav-tabs [data-testid="stPageLink"]:hover p { color: var(--text-dim) !important; }
.nav-tab-active {
    font-family: 'Inter Tight', sans-serif;
    font-size: 0.72rem;
    font-weight: 700;
    letter-spacing: 0.08em;
    text-transform: uppercase;
    color: var(--blue) !important;
    text-align: center;
    padding: 0.5rem 0.75rem;
    border-bottom: 2px solid var(--blue);
    margin-bottom: -1px;
    white-space: nowrap;
}

/*  HERO  */
.hero { padding: 1.5rem 0 0.75rem; }
.hero-title { font-family: 'Inter Tight', sans-serif; font-size: 2.4rem; font-weight: 900; line-height: 1.05; color: var(--text) !important; letter-spacing: -0.03em; margin-bottom: 0.5rem; }
.hero-title .accent { color: var(--blue); }
.hero-sub { font-size: 0.87rem; font-weight: 300; color: var(--muted) !important; max-width: 580px; line-height: 1.65; }

/*  QUERY BLOCK  */
.query-block {
    background: var(--surface);
    border: 1px solid var(--border);
    border-top: 2px solid var(--blue);
    border-radius: 0 0 10px 10px;
    padding: 1.4rem 1.75rem 1.25rem;
    margin: 1.25rem 0 0;
}
.field-label { font-size: 0.58rem; font-weight: 700; letter-spacing: 0.22em; text-transform: uppercase; color: var(--muted) !important; margin-bottom: 0.3rem; }

/*  FORM CONTROLS  */
.stTextInput > div > div > input,
.stTextArea > div > div > textarea {
    background: var(--bg) !important;
    border: 1px solid var(--border) !important;
    border-radius: 6px !important;
    color: var(--text) !important;
    font-family: 'Poppins', sans-serif !important;
    font-size: 0.88rem !important;
    caret-color: var(--blue) !important;
}
.stTextInput > div > div > input:focus,
.stTextArea > div > div > textarea:focus {
    border-color: var(--blue) !important;
    box-shadow: 0 0 0 3px var(--blue-glow) !important;
}
input::placeholder, textarea::placeholder { color: var(--muted) !important; opacity: 0.6 !important; }

div[data-baseweb="select"] > div,
div[data-baseweb="select"] > div > div {
    background: var(--bg) !important;
    border-color: var(--border) !important;
    color: var(--text) !important;
}
div[data-baseweb="select"] svg { fill: var(--muted) !important; }
div[data-baseweb="select"] span,
div[data-baseweb="select"] input { color: var(--text) !important; }
div[data-baseweb="menu"],
div[data-baseweb="popover"] {
    background: var(--surface2) !important;
    border: 1px solid var(--border-hi) !important;
    box-shadow: 0 8px 32px rgba(0,0,0,0.5) !important;
}
div[data-baseweb="menu"] li,
div[data-baseweb="menu"] [role="option"] { color: var(--text) !important; background: transparent !important; }
div[data-baseweb="menu"] li:hover,
div[data-baseweb="menu"] [aria-selected="true"] { background: var(--surface3) !important; color: var(--text) !important; }

.stCheckbox > label,
.stCheckbox > label > span,
.stCheckbox label p,
[data-testid="stCheckbox"] span,
[data-testid="stCheckbox"] p { color: var(--text) !important; font-size: 0.84rem !important; }

/*  BUTTONS  */
.stButton > button {
    background: var(--blue) !important;
    color: #fff !important;
    border: none !important;
    border-radius: 6px !important;
    font-family: 'Inter Tight', sans-serif !important;
    font-size: 0.78rem !important;
    font-weight: 800 !important;
    letter-spacing: 0.12em !important;
    text-transform: uppercase !important;
    padding: 0.5rem 1.5rem !important;
    transition: background 0.15s, box-shadow 0.15s, transform 0.1s !important;
    box-shadow: 0 2px 10px rgba(64,128,255,0.3) !important;
}
.stButton > button:hover {
    background: #2d6aee !important;
    box-shadow: 0 4px 18px rgba(64,128,255,0.45) !important;
    transform: translateY(-1px) !important;
}
.stButton > button:active { transform: translateY(0px) !important; }

.stDownloadButton > button {
    background: transparent !important;
    color: var(--blue) !important;
    border: 1px solid rgba(64,128,255,0.35) !important;
    border-radius: 6px !important;
    font-family: 'Inter Tight', sans-serif !important;
    font-size: 0.72rem !important;
    font-weight: 700 !important;
    letter-spacing: 0.1em !important;
    text-transform: uppercase !important;
    transition: all 0.15s !important;
    box-shadow: none !important;
}
.stDownloadButton > button:hover {
    background: var(--blue-glow) !important;
    border-color: var(--blue) !important;
    transform: none !important;
}

/*  METRIC CARDS  */
.metric-card {
    background: var(--surface);
    border: 1px solid var(--border);
    border-radius: 8px;
    padding: 1.2rem 1.4rem;
    position: relative;
    overflow: hidden;
    transition: border-color 0.2s, box-shadow 0.2s;
    height: 100%;
}
.metric-card::after {
    content: '';
    position: absolute;
    top: 0; left: 0; right: 0;
    height: 80px;
    background: linear-gradient(180deg, rgba(64,128,255,0.04) 0%, transparent 100%);
    pointer-events: none;
}
.metric-card.blue-top  { border-top: 2px solid var(--blue); }
.metric-card.pos-top   { border-top: 2px solid var(--pos); }
.metric-card.amber-top { border-top: 2px solid var(--amber); }
.metric-card.purple-top{ border-top: 2px solid var(--purple); }
.metric-card:hover { border-color: var(--border-hi); box-shadow: 0 4px 24px rgba(0,0,0,0.3); }
.metric-label { font-size: 0.58rem; font-weight: 700; letter-spacing: 0.22em; text-transform: uppercase; color: var(--muted) !important; margin-bottom: 0.45rem; }
.metric-value { font-family: 'Inter Tight', sans-serif; font-size: 2.1rem; font-weight: 900; color: var(--text) !important; line-height: 1; margin-bottom: 0.25rem; letter-spacing: -0.025em; }
.metric-sub { font-size: 0.69rem; color: var(--muted) !important; font-weight: 300; }

/*  SECTION HEADER  */
.section-header {
    font-family: 'Inter Tight', sans-serif;
    font-size: 0.68rem;
    font-weight: 800;
    letter-spacing: 0.24em;
    text-transform: uppercase;
    color: var(--text-dim) !important;
    margin: 1.75rem 0 0.9rem;
    padding-bottom: 0.55rem;
    border-bottom: 1px solid var(--border);
    display: flex;
    align-items: center;
    gap: 0.55rem;
}
.section-header .dot { width: 5px; height: 5px; background: var(--blue); border-radius: 1px; display: inline-block; flex-shrink: 0; box-shadow: 0 0 5px var(--blue); }

/*  PROGRESS BARS  */
.stProgress > div > div > div > div {
    background: linear-gradient(90deg, var(--blue) 0%, #7ab0ff 100%) !important;
    border-radius: 4px !important;
}

/*  TABS  */
.stTabs [data-baseweb="tab-list"] {
    gap: 0 !important;
    border-bottom: 1px solid var(--border) !important;
    background: transparent !important;
    margin-bottom: 0.25rem !important;
}
.stTabs [data-baseweb="tab"] {
    background: transparent !important;
    color: var(--muted) !important;
    font-family: 'Inter Tight', sans-serif !important;
    font-weight: 700 !important;
    font-size: 0.68rem !important;
    letter-spacing: 0.16em !important;
    text-transform: uppercase !important;
    padding: 0.6rem 1.1rem !important;
    border-bottom: 2px solid transparent !important;
    transition: color 0.15s !important;
}
.stTabs [data-baseweb="tab"]:hover { color: var(--text-dim) !important; }
.stTabs [aria-selected="true"] { color: var(--text) !important; border-bottom-color: var(--blue) !important; }
.stTabs [data-baseweb="tab-panel"] { padding-top: 0.5rem !important; }

/*  EXPANDERS  */
[data-testid="stExpander"],
details[data-testid="stExpander"] {
    background: var(--surface) !important;
    border: 1px solid var(--border) !important;
    border-radius: 8px !important;
    overflow: hidden;
}
[data-testid="stExpander"] summary,
[data-testid="stExpander"] summary span,
[data-testid="stExpander"] summary p,
[data-testid="stExpander"] summary div { color: var(--text) !important; background: var(--surface) !important; }
[data-testid="stExpanderDetails"],
[data-testid="stExpanderDetails"] > div { background: var(--surface) !important; color: var(--text) !important; }

/*  DATA TABLE  */
[data-testid="stDataFrame"] {
    background: var(--surface) !important;
    border: 1px solid var(--border) !important;
    border-radius: 8px !important;
    overflow: hidden;
}

/*  ALERTS  */
[data-testid="stAlert"],
div[data-baseweb="notification"] {
    background: var(--surface2) !important;
    border: 1px solid var(--border-hi) !important;
    border-radius: 6px !important;
    color: var(--text) !important;
}
[data-testid="stAlert"] p,
[data-testid="stAlert"] span { color: var(--text) !important; }

/*  SPINNER  */
[data-testid="stSpinner"] p { color: var(--text) !important; }

/*  INSIGHT CARD (query preset)  */
.insight-card {
    background: var(--surface);
    border: 1px solid var(--border);
    border-left: 3px solid var(--blue);
    border-radius: 0 8px 8px 0;
    padding: 1rem 1.2rem;
    margin-bottom: 0.75rem;
    cursor: pointer;
    transition: border-color 0.15s, box-shadow 0.15s;
}
.insight-card:hover { border-left-color: #7ab0ff; box-shadow: 0 4px 20px rgba(64,128,255,0.15); }
.insight-card-title { font-family: 'Inter Tight', sans-serif; font-weight: 800; font-size: 0.82rem; color: var(--text) !important; letter-spacing: 0.04em; margin-bottom: 0.3rem; }
.insight-card-desc { font-size: 0.74rem; color: var(--muted) !important; line-height: 1.55; }

/*  TAG BADGE  */
.tag { display: inline-block; background: var(--blue-glow); border: 1px solid rgba(64,128,255,0.28); border-radius: 4px; padding: 0.1rem 0.5rem; font-size: 0.6rem; font-weight: 700; letter-spacing: 0.12em; text-transform: uppercase; color: var(--blue) !important; margin-right: 0.35rem; }
.tag.amber { background: var(--amber-dim); border-color: rgba(255,185,56,0.28); color: var(--amber) !important; }
.tag.pos   { background: var(--pos-dim);   border-color: rgba(32,198,90,0.28);  color: var(--pos)   !important; }
.tag.purple{ background: var(--purple-dim);border-color: rgba(168,85,247,0.28); color: var(--purple)!important; }

/*  EMPTY STATE  */
.empty-state {
    margin-top: 3.5rem;
    text-align: center;
    padding: 4rem 2rem;
    border: 1px dashed var(--border-hi);
    border-radius: 12px;
    background: radial-gradient(ellipse at 50% 0%, rgba(64,128,255,0.05) 0%, transparent 65%);
}
.empty-title { font-family: 'Inter Tight', sans-serif; font-size: 2rem; font-weight: 900; color: var(--border-hi) !important; letter-spacing: -0.02em; margin-bottom: 0.7rem; }
.empty-sub { font-size: 0.86rem; color: var(--muted) !important; max-width: 380px; margin: 0 auto; line-height: 1.75; }

/*  FOOTER  */
.footer {
    margin-top: 4rem;
    padding-top: 1.25rem;
    border-top: 1px solid var(--border);
    display: flex;
    justify-content: space-between;
    align-items: center;
}
.footer-brand { font-family: 'Inter Tight', sans-serif; font-weight: 900; font-size: 0.7rem; color: var(--border-hi) !important; letter-spacing: 0.14em; text-transform: uppercase; }
.footer-note { font-size: 0.63rem; color: var(--muted) !important; }

/*  CHAT  */
[data-testid="stChatMessage"] { background: var(--surface) !important; border: 1px solid var(--border) !important; border-radius: 8px !important; }
[data-testid="stChatInput"] > div { background: var(--surface2) !important; border-color: var(--border-hi) !important; }
[data-testid="stChatInput"] textarea { color: var(--text) !important; }

/* ─────────────────────────────────────────────────────────────
   MOBILE RESPONSIVENESS
   Below 768px: stack multi-column layouts, enlarge tap targets,
   shrink the oversized hero/metric typography that otherwise
   overflows a phone viewport.
   ───────────────────────────────────────────────────────────── */
@media (max-width: 768px) {
    .block-container { padding: 0 1rem 3rem !important; }

    .topbar { padding: 0.7rem 1rem; margin: 0 -1rem 1.25rem; flex-wrap: wrap; gap: 0.6rem; }
    .topbar-logo { font-size: 0.78rem; }
    .topbar-label { display: none; }   /* redundant subtitle on narrow screens */
    .topbar-pill { margin-left: 0; }

    .hero-title { font-size: 1.6rem; }
    .hero-sub   { font-size: 0.8rem; max-width: 100%; }

    .query-block { padding: 1rem 1.1rem 0.9rem; }

    /* Streamlit's horizontal blocks (st.columns) collapse to single column */
    div[data-testid="stHorizontalBlock"] {
        flex-direction: column !important;
    }
    div[data-testid="stHorizontalBlock"] > div[data-testid="column"] {
        width: 100% !important;
        flex: 1 1 100% !important;
        min-width: 100% !important;
        margin-bottom: 0.6rem;
    }

    .metric-card { padding: 1rem 1.1rem; }
    .metric-value { font-size: 1.6rem; }

    .section-header { font-size: 0.6rem; margin: 1.25rem 0 0.7rem; }

    /* Buttons and download buttons need a larger tap target on touch screens */
    .stButton > button, .stDownloadButton > button {
        padding: 0.65rem 1.1rem !important;
        font-size: 0.74rem !important;
        width: 100%;
    }

    .empty-state { padding: 2.5rem 1.25rem; margin-top: 2rem; }
    .empty-title { font-size: 1.4rem; }

    .footer { flex-direction: column; align-items: flex-start; gap: 0.4rem; }
}

@media (max-width: 480px) {
    .hero-title { font-size: 1.35rem; }
    .metric-value { font-size: 1.4rem; }
}
</style>

""", unsafe_allow_html=True)


# ─────────────────────────────────────────────────────────────
# CONSTANTS & DATA
# ─────────────────────────────────────────────────────────────

PLOTLY_BASE = dict(
    paper_bgcolor="rgba(0,0,0,0)",
    plot_bgcolor="rgba(0,0,0,0)",
    font=dict(family="Poppins, sans-serif", color="#eef0fa"),
    margin=dict(l=10, r=10, t=30, b=10),
)

# Top 10 shooters on Steam (source: SteamDB live charts, March 2026)
#  Roster data — all unique titles across both lists
# Sub-genre and publisher are looked up by app_id; the two ranked lists
# below reference these by ID only to avoid duplication.

GAME_CATALOG = {
    730:     {"name": "Counter-Strike 2",             "sub": "Tactical / Competitive",   "publisher": "Valve",                  "f2p": True, "year": 2023},
    578080:  {"name": "PUBG: Battlegrounds",          "sub": "Battle Royale",            "publisher": "Krafton",                "f2p": True, "year": 2022},
    252490:  {"name": "Rust",                         "sub": "Open World / Survival",    "publisher": "Facepunch",              "f2p": False, "year": 2013},
    1172470: {"name": "Apex Legends",                 "sub": "Battle Royale / Hero",     "publisher": "EA / Respawn",           "f2p": True, "year": 2022},
    3764200: {"name": "Monster Hunter Wilds",         "sub": "Action / Co-op",           "publisher": "Capcom",                 "f2p": False, "year": 2025},
    2357570: {"name": "Overwatch 2",                  "sub": "Hero Shooter",             "publisher": "Blizzard",               "f2p": True, "year": 2022},
    2507950: {"name": "Delta Force",                  "sub": "Military FPS / Extraction","publisher": "Team Jade",              "f2p": True, "year": 2024},
    359550:  {"name": "Rainbow Six Siege",            "sub": "Tactical / Competitive",   "publisher": "Ubisoft",                "f2p": False, "year": 2015},
    440:     {"name": "Team Fortress 2",              "sub": "Arena / Class FPS",        "publisher": "Valve",                  "f2p": True, "year": 2007},
    221100:  {"name": "DayZ",                         "sub": "Survival / Open World",    "publisher": "Bohemia Interactive",    "f2p": False, "year": 2018},
    2767030: {"name": "Marvel Rivals",                "sub": "Hero Shooter",             "publisher": "NetEase Games",          "f2p": True, "year": 2024},
    1366800: {"name": "Crosshair X",                  "sub": "Aim Trainer / Utility",    "publisher": "Vanguard",               "f2p": False, "year": 2021},
    2807960: {"name": "Battlefield 6",                "sub": "Military FPS",             "publisher": "EA / DICE",              "f2p": False, "year": 2025},
    3065800: {"name": "Marathon",                     "sub": "Extraction Shooter",       "publisher": "Bungie / Sony",          "f2p": False, "year": 2025},
    4465480: {"name": "CS:GO",                        "sub": "Tactical / Competitive",   "publisher": "Valve",                  "f2p": True, "year": 2012},
    1938090: {"name": "Call of Duty",                 "sub": "Military FPS",             "publisher": "Activision",             "f2p": True, "year": 2023},
    1174180: {"name": "Red Dead Redemption 2",        "sub": "Open World / TPS",         "publisher": "Rockstar Games",         "f2p": False, "year": 2018},
    4000:    {"name": "Garry's Mod",                  "sub": "Sandbox / Shooter",        "publisher": "Facepunch",              "f2p": False, "year": 2004},
    1091500: {"name": "Cyberpunk 2077",               "sub": "FPS / RPG",                "publisher": "CD Projekt Red",         "f2p": False, "year": 2020},
    2073620: {"name": "Arena Breakout: Infinite",     "sub": "Extraction Shooter",       "publisher": "Level Infinite",         "f2p": True, "year": 2024},
    251570:  {"name": "7 Days to Die",                "sub": "Survival / FPS",           "publisher": "The Fun Pimps",          "f2p": False, "year": 2013},
    1818450: {"name": "The Finals",                   "sub": "Arena / Team FPS",         "publisher": "Embark Studios",         "f2p": True, "year": 2023},
    377160:  {"name": "Fallout 4",                    "sub": "Open World FPS / RPG",     "publisher": "Bethesda",               "f2p": False, "year": 2015},
    550:     {"name": "Left 4 Dead 2",                "sub": "Co-op / Survival FPS",     "publisher": "Valve",                  "f2p": False, "year": 2009},
    1151340: {"name": "Fallout 76",                   "sub": "Online FPS / RPG",         "publisher": "Bethesda",               "f2p": False, "year": 2018},
    1808500: {"name": "ARC Raiders",                  "sub": "Extraction Shooter",       "publisher": "Embark Studios",         "f2p": False, "year": 2025},
    271590:  {"name": "GTA V Legacy",                 "sub": "Open World / Action",      "publisher": "Rockstar Games",         "f2p": False, "year": 2015},
    236390:  {"name": "War Thunder",                  "sub": "Vehicle Combat / MMO",     "publisher": "Gaijin",                 "f2p": True, "year": 2013},
    1422450: {"name": "Deadlock",                     "sub": "Hero Shooter / MOBA",      "publisher": "Valve",                  "f2p": True, "year": 2024},
    230410:  {"name": "Warframe",                     "sub": "Looter Shooter / Co-op",   "publisher": "Digital Extremes",       "f2p": True, "year": 2013},
    3240220: {"name": "GTA V Enhanced",               "sub": "Open World / Action",      "publisher": "Rockstar Games",         "f2p": False, "year": 2025},
    553850:  {"name": "Helldivers 2",                 "sub": "Co-op / Third-Person",     "publisher": "PlayStation Studios",    "f2p": False, "year": 2024},
    1623730: {"name": "Palworld",                     "sub": "Survival / Action",        "publisher": "Pocketpair",             "f2p": False, "year": 2024},
    2050650: {"name": "Resident Evil 4 Remake",       "sub": "Survival Horror / TPS",    "publisher": "Capcom",                 "f2p": False, "year": 2023},
    2221490: {"name": "Tom Clancy's The Division 2",  "sub": "Cover Shooter / MMO",      "publisher": "Ubisoft",                "f2p": False, "year": 2019},
    2183900: {"name": "Warhammer 40K: Space Marine 2","sub": "TPS / Co-op Action",       "publisher": "Saber Interactive",      "f2p": False, "year": 2024},
    107410:  {"name": "Arma 3",                       "sub": "Tactical / Mil-Sim",       "publisher": "Bohemia Interactive",    "f2p": False, "year": 2013},
    1407200: {"name": "World of Tanks",               "sub": "Vehicle Combat / MMO",     "publisher": "Wargaming",              "f2p": True, "year": 2011},
    3405340: {"name": "Split Fiction",                "sub": "Co-op / Action",           "publisher": "Hazelight / EA",         "f2p": False, "year": 2025},
    3659280: {"name": "Atomfall",                     "sub": "Open World FPS / RPG",     "publisher": "Rebellion",              "f2p": False, "year": 2025},
    552990:  {"name": "Borderlands 3",                "sub": "Looter Shooter",           "publisher": "2K / Gearbox",           "f2p": False, "year": 2019},
    240:     {"name": "Counter-Strike: Source",        "sub": "Tactical / Competitive",   "publisher": "Valve",                  "f2p": False, "year": 2004},
    1659040: {"name": "HITMAN World of Assassination",  "sub": "Stealth / TPS",            "publisher": "IO Interactive",         "f2p": False, "year": 2021},
}

# Ranked list of top 25 FPS titles (by Steam CCU, Mar 2026)
# Locked to the 42 SteamDB CSVs in /data — 1818450 removed (no CSV), 240 added
FPS_ROSTER_IDS = [
    730, 578080, 252490, 1172470, 3764200, 2357570, 2507950, 359550, 440,
    221100, 2767030, 2807960, 3065800, 4465480, 1938090, 1174180,
    4000, 1091500, 2073620, 251570, 377160, 550, 1151340, 240,
]

# Ranked list of top 25 TPS titles (by Steam CCU, Mar 2026)
# Locked to the 42 SteamDB CSVs in /data — 1329410 removed (no CSV), 1659040 added
TPS_ROSTER_IDS = [
    578080, 1808500, 271590, 3764200, 2357570, 236390, 1422450, 230410,
    2767030, 3240220, 553850, 1938090, 1623730, 1174180, 2050650,
    377160, 2221490, 2183900, 107410, 1407200, 3405340, 3659280, 552990,
    1659040,
]

# ─────────────────────────────────────────────────────────────
# NOTABLE EVENTS PER GAME  (date → label for chart annotations)
# ─────────────────────────────────────────────────────────────

# Steam News API — replaces hardcoded GAME_EVENTS
STEAM_NEWS_URL = "https://api.steampowered.com/ISteamNews/GetNewsForApp/v2/"

# Keywords that suggest a meaningful update (patch notes, seasons, DLC, major releases)
_NEWS_KEYWORDS = [
    "update", "patch", "season", "dlc", "expansion", "launch", "release",
    "major", "anniversary", "free", "wipe", "reset", "event", "content",
    "chapter", "operation", "overhaul", "rework", "battle pass", "new map",
    "new mode", "hotfix", "balance", "2.0", "3.0",
]

@st.cache_data(ttl=3600, show_spinner=False)
def fetch_steam_news(app_id: int, count: int = 100) -> list[tuple[str, str]]:
    """
    Fetch latest news for a game from the Steam News API.
    Returns list of (YYYY-MM, title) tuples filtered to meaningful update events,
    deduplicated to one event per month (the most prominent).
    """
    try:
        r = requests.get(
            STEAM_NEWS_URL,
            params={
                "appid":     app_id,
                "count":     count,
                "maxlength": 0,       # titles only, no body content
                "format":    "json",
            },
            timeout=10,
        )
        if not r.ok:
            return []
        items = r.json().get("appnews", {}).get("newsitems", [])
    except Exception:
        return []

    events: dict[str, str] = {}  # month -> best label (dedup to 1 per month)
    for item in items:
        title = item.get("title", "").strip()
        ts    = item.get("date", 0)
        if not title or not ts:
            continue
        title_lower = title.lower()
        if not any(kw in title_lower for kw in _NEWS_KEYWORDS):
            continue
        from datetime import datetime, timezone
        dt    = datetime.fromtimestamp(ts, tz=timezone.utc)
        month = dt.strftime("%Y-%m")
        # Truncate long titles for annotation legibility
        label = title[:40].rstrip() + ("…" if len(title) > 40 else "")
        # Keep shortest/clearest label per month (prefer patch notes)
        if month not in events or len(label) < len(events[month]):
            events[month] = label

    # Return sorted oldest → newest
    return sorted(events.items())


def get_game_events(app_id: int) -> list[tuple[str, str]]:
    """Wrapper — fetches from Steam News API with caching."""
    return fetch_steam_news(app_id)

def get_roster(genre: str = "FPS") -> list[dict]:
    """Return roster list for given genre, merging catalog metadata."""
    ids = FPS_ROSTER_IDS if genre == "FPS" else TPS_ROSTER_IDS
    result = []
    for app_id in ids:
        if app_id in GAME_CATALOG:
            result.append({"app_id": app_id, **GAME_CATALOG[app_id]})
    return result

# Default active roster (overridden by sidebar filter at runtime)
SHOOTER_ROSTER = get_roster("FPS")

# Folder containing SteamDB CSVs — checks several locations in priority order
def _find_data_dir() -> Path | None:
    candidates = [
        Path(__file__).parent / "data",          # local dev & Streamlit Cloud repo
        Path(__file__).parent,                    # CSVs alongside the .py file
        Path("/mount/src") / Path(__file__).stem / "data",  # Streamlit Cloud alt mount
        Path("/data"),                            # Docker / custom deploy
    ]
    for p in candidates:
        if p.exists() and any(p.glob("steamdb_chart_*.csv")):
            return p
    # Return the default even if empty — loaders will handle it gracefully
    return Path(__file__).parent / "data"

DATA_DIR = _find_data_dir()

# Steam CCU endpoint
CCU_URL = "https://api.steampowered.com/ISteamUserStats/GetNumberOfCurrentPlayers/v1/"

# Max concurrent requests per upstream API. The original sequential fetch
# loop naturally spaced 41 requests out over ~16+ seconds (0.4s sleep between
# each); the parallel fetch can blast a burst of simultaneous requests at
# Steam/SteamSpy from a single IP instead, which is a much more likely
# trigger for rate-limiting/throttling than the original code ever hit.
# These semaphores cap concurrency per-API regardless of how many worker
# threads are running, and _http_get_with_retry() below adds a couple of
# retries with backoff so a single transient blip doesn't permanently zero
# out that title's CCU for the rest of the 5-minute cache window.
_STEAM_CCU_SEM = threading.Semaphore(6)
_STEAMSPY_SEM  = threading.Semaphore(4)


def _http_get_with_retry(url: str, params: dict, timeout: float,
                         max_retries: int = 3, base_delay: float = 0.6) -> requests.Response | None:
    """GET with retry-with-backoff. Returns the response on success (status
    200-299), or None if every attempt failed (timeout, connection error, or
    non-2xx status). Used by fetch_ccu / fetch_steamspy / fetch_steam_reviews
    so a single rate-limited or transiently-failed request doesn't silently
    zero out that title for the rest of the cache TTL."""
    for attempt in range(max_retries):
        try:
            r = requests.get(url, params=params, timeout=timeout)
            if r.ok:
                return r
        except Exception:
            pass
        if attempt < max_retries - 1:
            time.sleep(base_delay * (2 ** attempt) + random.uniform(0, 0.25))
    return None

# ── Twitch Helix API ─────────────────────────────────────────
TWITCH_TOKEN_URL   = "https://id.twitch.tv/oauth2/token"
TWITCH_GAMES_URL   = "https://api.twitch.tv/helix/games"
TWITCH_STREAMS_URL = "https://api.twitch.tv/helix/streams"

# Process-lifetime token cache shared across all Streamlit reruns and worker threads.
# The lock prevents concurrent token-refresh races when the pool starts.
_TWITCH_TOKEN: dict[str, str | float] = {"value": "", "expires_at": 0.0}
_TWITCH_TOKEN_LOCK = threading.Lock()

# Game-name → Twitch game_id, cached for the process lifetime (avoids repeated search calls).
_TWITCH_GAME_ID_CACHE: dict[str, str | None] = {}

# Games whose name on Twitch differs from GAME_CATALOG (case-sensitive exact match required).
TWITCH_NAME_MAP: dict[int, str] = {
    578080:  "PLAYERUNKNOWN'S BATTLEGROUNDS",
    1174180: "Red Dead Redemption 2",
    271590:  "Grand Theft Auto V",
    3240220: "Grand Theft Auto V",          # GTA V Enhanced shares the same Twitch category
    4465480: "Counter-Strike: Global Offensive",
    230410:  "Warframe",
    1407200: "World of Tanks",
    1659040: "HITMAN World of Assassination",
    2183900: "Warhammer 40,000: Space Marine 2",
}


def _get_twitch_token() -> str:
    """Return a valid Twitch OAuth app-access token, refreshing automatically.
    Thread-safe: _TWITCH_TOKEN_LOCK prevents concurrent refresh requests."""
    with _TWITCH_TOKEN_LOCK:
        now = time.time()
        if _TWITCH_TOKEN["value"] and now < float(_TWITCH_TOKEN["expires_at"]) - 60:
            return str(_TWITCH_TOKEN["value"])
        try:
            r = requests.post(
                TWITCH_TOKEN_URL,
                params={
                    "client_id":     st.secrets.get("TWITCH_CLIENT_ID", ""),
                    "client_secret": st.secrets.get("TWITCH_CLIENT_SECRET", ""),
                    "grant_type":    "client_credentials",
                },
                timeout=8,
            )
            if r.ok:
                d = r.json()
                _TWITCH_TOKEN["value"]      = d["access_token"]
                _TWITCH_TOKEN["expires_at"] = now + d.get("expires_in", 3600)
                return str(_TWITCH_TOKEN["value"])
        except Exception:
            pass
    return ""


def _twitch_game_id(game_name: str) -> str | None:
    """Resolve a Twitch game_id from a display name.
    Results are cached per process — avoids repeated search calls for the same title."""
    if game_name in _TWITCH_GAME_ID_CACHE:
        return _TWITCH_GAME_ID_CACHE[game_name]
    token = _get_twitch_token()
    if not token:
        _TWITCH_GAME_ID_CACHE[game_name] = None
        return None
    try:
        r = requests.get(
            TWITCH_GAMES_URL,
            params={"name": game_name},
            headers={
                "Client-ID":     st.secrets.get("TWITCH_CLIENT_ID", ""),
                "Authorization": f"Bearer {token}",
            },
            timeout=8,
        )
        if r.ok:
            data = r.json().get("data", [])
            gid  = data[0]["id"] if data else None
            _TWITCH_GAME_ID_CACHE[game_name] = gid
            return gid
    except Exception:
        pass
    _TWITCH_GAME_ID_CACHE[game_name] = None
    return None


@st.cache_data(ttl=300, show_spinner=False)
def fetch_twitch_viewers(app_id: int, game_name: str) -> int | None:
    """
    Return total live viewer count on Twitch (sum across top-100 streams).
    Returns None when TWITCH_CLIENT_ID is not configured or the game isn't found.

    Add to .streamlit/secrets.toml to enable:
        TWITCH_CLIENT_ID     = "..."
        TWITCH_CLIENT_SECRET = "..."
    """
    if not st.secrets.get("TWITCH_CLIENT_ID"):
        return None
    twitch_name = TWITCH_NAME_MAP.get(app_id, game_name)
    game_id     = _twitch_game_id(twitch_name)
    if not game_id:
        return None
    token = _get_twitch_token()
    if not token:
        return None
    try:
        r = requests.get(
            TWITCH_STREAMS_URL,
            params={"game_id": game_id, "first": 100},
            headers={
                "Client-ID":     st.secrets.get("TWITCH_CLIENT_ID", ""),
                "Authorization": f"Bearer {token}",
            },
            timeout=10,
        )
        if r.ok:
            return sum(s.get("viewer_count", 0) for s in r.json().get("data", []))
    except Exception:
        pass
    return None

# ─────────────────────────────────────────────────────────────
# STEAMDB HISTORICAL CSV LOADER
# ─────────────────────────────────────────────────────────────

@st.cache_data(show_spinner=False)
def _parse_steamdb_bytes(raw_bytes: bytes, app_id: int) -> pd.DataFrame | None:
    """Parse raw SteamDB CSV bytes into a monthly peak/avg DataFrame."""
    try:
        import io
        df = pd.read_csv(io.BytesIO(raw_bytes), encoding="utf-8-sig")
        df.columns = [c.strip().strip('"') for c in df.columns]
        df["DateTime"] = pd.to_datetime(df["DateTime"], errors="coerce")
        df = df.dropna(subset=["DateTime"])
        df["Players"] = pd.to_numeric(df["Players"], errors="coerce")
        df["Average Players"] = pd.to_numeric(df.get("Average Players", pd.Series()), errors="coerce")
        df["month"] = df["DateTime"].dt.to_period("M")
        monthly = (
            df.groupby("month")
            .agg(peak_ccu=("Players", "max"), avg_ccu=("Average Players", "mean"))
            .reset_index()
        )
        return monthly.sort_values("month")
    except Exception:
        return None


def _parse_steamdb_bytes_raw(raw_bytes: bytes) -> pd.DataFrame | None:
    """Parse raw SteamDB CSV bytes into a raw (10-min interval) DataFrame."""
    try:
        import io
        df = pd.read_csv(io.BytesIO(raw_bytes), encoding="utf-8-sig")
        df.columns = [c.strip().strip('"') for c in df.columns]
        df["DateTime"] = pd.to_datetime(df["DateTime"], errors="coerce", utc=True)
        df = df.dropna(subset=["DateTime"])
        df["Players"] = pd.to_numeric(df["Players"], errors="coerce")
        df = df.dropna(subset=["Players"])
        return df[["DateTime", "Players"]].sort_values("DateTime").reset_index(drop=True)
    except Exception:
        return None


@st.cache_data(show_spinner=False)
def load_all_historical(roster_ids: frozenset[int] | None = None) -> dict[int, pd.DataFrame]:
    """
    Loads SteamDB CSVs — from /data folder AND from sidebar-uploaded files.
    Returns a dict of {app_id: monthly_df} with columns: month (Period), peak_ccu, avg_ccu

    roster_ids: when provided, only parse CSVs whose app_id is in this set.
    Pass None (default) to load every CSV on disk — used by the background scheduler.
    """
    historical: dict[int, pd.DataFrame] = {}

    # Priority 1: sidebar-uploaded files (in session state)
    for app_id, raw_bytes in st.session_state.get("uploaded_csvs", {}).items():
        if roster_ids is not None and app_id not in roster_ids:
            continue
        mdf = _parse_steamdb_bytes(raw_bytes, app_id)
        if mdf is not None and not mdf.empty:
            historical[app_id] = mdf

    # Priority 2: files on disk — skip CSVs outside the active roster
    if not DATA_DIR.exists():
        return historical

    for csv_path in sorted(DATA_DIR.glob("steamdb_chart_*.csv")):
        try:
            app_id = int(csv_path.stem.replace("steamdb_chart_", ""))
        except ValueError:
            continue
        if roster_ids is not None and app_id not in roster_ids:
            continue
        try:
            df = pd.read_csv(csv_path, encoding="utf-8-sig")
            df.columns = [c.strip().strip('"') for c in df.columns]
            df["DateTime"] = pd.to_datetime(df["DateTime"], errors="coerce")
            df = df.dropna(subset=["DateTime"])
            df["Players"] = pd.to_numeric(df["Players"], errors="coerce")
            df["Average Players"] = pd.to_numeric(df["Average Players"], errors="coerce")
            df["month"] = df["DateTime"].dt.to_period("M")

            monthly = (
                df.groupby("month")
                .agg(peak_ccu=("Players", "max"), avg_ccu=("Average Players", "mean"))
                .reset_index()
            )
            monthly = monthly.sort_values("month")
            historical[app_id] = monthly
        except Exception:
            pass  # Skip malformed CSVs silently

    return historical


def compute_yoy(monthly_df: pd.DataFrame) -> tuple[str, float]:
    """
    Compute YoY change in average CCU: compare most recent complete month
    to the same month one year prior. Returns (display_str, pct_float).
    """
    if monthly_df is None or len(monthly_df) < 2:
        return "N/A", 0.0

    now = pd.Period.now("M")
    # Use last complete month (not current partial month)
    last_complete = now - 1
    year_ago = last_complete - 12

    row_now  = monthly_df[monthly_df["month"] == last_complete]
    row_prev = monthly_df[monthly_df["month"] == year_ago]

    if row_now.empty or row_prev.empty:
        # Fall back to most recent vs same-month last year using available data
        latest = monthly_df.iloc[-1]
        target_month = latest["month"] - 12
        row_prev2 = monthly_df[monthly_df["month"] == target_month]
        if row_prev2.empty:
            return "N/A", 0.0
        val_now  = latest["avg_ccu"] if not pd.isna(latest["avg_ccu"]) else latest["peak_ccu"]
        val_prev = row_prev2.iloc[0]["avg_ccu"]
        if pd.isna(val_prev) or val_prev == 0:
            val_prev = row_prev2.iloc[0]["peak_ccu"]
    else:
        val_now  = row_now.iloc[0]["avg_ccu"]
        val_prev = row_prev.iloc[0]["avg_ccu"]
        if pd.isna(val_now):  val_now  = row_now.iloc[0]["peak_ccu"]
        if pd.isna(val_prev): val_prev = row_prev.iloc[0]["peak_ccu"]

    if pd.isna(val_now) or pd.isna(val_prev) or val_prev == 0:
        return "N/A", 0.0

    pct = (val_now - val_prev) / val_prev * 100
    pct_capped = max(-999.0, min(999.0, pct))
    sign = "+" if pct_capped >= 0 else ""
    return f"{sign}{round(pct_capped)}%", pct_capped


def compute_yoy_from_snapshots(snapshots: list[dict], app_id: int) -> tuple[str, float] | None:
    """
    Compute real YoY CCU change from accumulated app snapshots.

    Needs at least one saved snapshot within ±30 days of exactly 365 days ago.
    Returns None when history is too short (caller falls back to CSV or SteamSpy).
    Becomes accurate after ~13 months of daily use; silent no-op before that.
    """
    if not snapshots or len(snapshots) < 2:
        return None

    now    = datetime.utcnow()
    target = now - timedelta(days=365)
    tol    = timedelta(days=30)

    # Most-recent snapshot → current CCU baseline
    curr_snap = max(snapshots, key=lambda s: s["ts"])
    curr_ccu  = curr_snap["data"].get(str(app_id))
    if not curr_ccu:
        return None
    curr_ccu = int(curr_ccu)

    # Snapshot closest to 365 days ago
    best_val = None
    best_gap: timedelta | None = None
    for s in snapshots:
        try:
            ts  = datetime.fromisoformat(s["ts"])
            gap = abs(ts - target)
            v   = s["data"].get(str(app_id))
            if v is None:
                continue
            if best_gap is None or gap < best_gap:
                best_gap = gap
                best_val = int(v)
        except Exception:
            continue

    if best_val is None or best_gap is None or best_gap > tol or best_val == 0:
        return None

    pct       = (curr_ccu - best_val) / best_val * 100
    pct_cap   = max(-999.0, min(999.0, pct))
    sign      = "+" if pct_cap >= 0 else ""
    return f"{sign}{round(pct_cap)}%", pct_cap



def get_historical_summary(monthly_df: pd.DataFrame) -> dict:
    """Return a summary dict for the AI prompt from historical data."""
    if monthly_df is None or monthly_df.empty:
        return {}
    last_12 = monthly_df.tail(12)
    peak_ever = monthly_df["peak_ccu"].max()
    peak_12m  = last_12["peak_ccu"].max()
    avg_12m   = last_12["avg_ccu"].mean()
    # Month-over-month trend (slope sign over last 3 months)
    last_3 = monthly_df.tail(3)
    if len(last_3) >= 2:
        vals = last_3["avg_ccu"].fillna(last_3["peak_ccu"]).dropna().tolist()
        mom_trend = "↑" if len(vals) >= 2 and vals[-1] > vals[0] else "↓"
    else:
        mom_trend = "—"
    # Real MoM percentage: last month vs month before
    last_2 = monthly_df.tail(2)
    mom_pct = None
    if len(last_2) == 2:
        v1 = last_2.iloc[-2]["avg_ccu"] or last_2.iloc[-2]["peak_ccu"]
        v2 = last_2.iloc[-1]["avg_ccu"] or last_2.iloc[-1]["peak_ccu"]
        if v1 and v1 > 0 and not pd.isna(v1) and not pd.isna(v2):
            mom_pct = (v2 - v1) / v1 * 100
            sign = "+" if mom_pct >= 0 else ""
            mom_trend = f"{sign}{round(mom_pct)}%"
    # 1-year-ago peak CCU (for live vs 1yr comparison)
    yoy_ccu = None
    if len(monthly_df) >= 12:
        yr_ago_row = monthly_df.iloc[-12]
        v = yr_ago_row.get("peak_ccu") or yr_ago_row.get("avg_ccu")
        if v and not pd.isna(v):
            yoy_ccu = int(v)
    return {
        "peak_ever":  int(peak_ever) if not pd.isna(peak_ever) else None,
        "peak_12m":   int(peak_12m)  if not pd.isna(peak_12m)  else None,
        "avg_12m":    int(avg_12m)   if not pd.isna(avg_12m)   else None,
        "mom_trend":  mom_trend,
        "mom_pct":    mom_pct,
        "months_data": len(monthly_df),
        "yoy_ccu":    yoy_ccu,
    }


def compute_review_velocity(snapshots: list[dict], app_id: int,
                             days: int = 7) -> float | None:
    """
    Compute the rate of new Steam reviews (reviews/day) over the last `days`.

    Review counts are stored alongside CCU in each snapshot entry under "reviews".
    Returns None when there isn't a snapshot from ~`days` ago (new installs, or
    when the snapshot file pre-dates the reviews field being added).
    """
    if not snapshots:
        return None

    # Most-recent snapshot with review data → current total
    curr_total: int | None = None
    for s in sorted(snapshots, key=lambda x: x["ts"], reverse=True):
        v = s.get("reviews", {}).get(str(app_id))
        if v is not None:
            curr_total = int(v)
            break
    if curr_total is None:
        return None

    # Snapshot closest to `days` ago
    target   = datetime.utcnow() - timedelta(days=days)
    tol      = timedelta(hours=96)
    prev_total: int | None = None
    best_gap: timedelta | None = None
    for s in snapshots:
        try:
            ts  = datetime.fromisoformat(s["ts"])
            gap = abs(ts - target)
            v   = s.get("reviews", {}).get(str(app_id))
            if v is None:
                continue
            if best_gap is None or gap < best_gap:
                best_gap  = gap
                prev_total = int(v)
        except Exception:
            continue

    if prev_total is None or best_gap is None or best_gap > tol:
        return None

    delta = curr_total - prev_total
    if delta < 0:
        return None   # data anomaly — reviews shouldn't decrease
    return round(delta / days, 1)


@st.cache_data(ttl=600, show_spinner=False)
def load_all_raw(roster_ids: frozenset[int] | None = None) -> dict[int, pd.DataFrame]:
    """Load raw 10-minute interval CSV data (no aggregation) for WoW diff.
    Also reads from sidebar-uploaded session-state files.

    roster_ids: when provided, only load CSVs for games in this set.
    """
    raw: dict[int, pd.DataFrame] = {}

    # Priority 1: sidebar-uploaded files
    for app_id, rb in st.session_state.get("uploaded_csvs", {}).items():
        if roster_ids is not None and app_id not in roster_ids:
            continue
        rdf = _parse_steamdb_bytes_raw(rb)
        if rdf is not None and not rdf.empty:
            raw[app_id] = rdf

    # Priority 2: disk — skip CSVs outside the active roster
    if not DATA_DIR.exists():
        return raw
    for csv_path in sorted(DATA_DIR.glob("steamdb_chart_*.csv")):
        try:
            app_id = int(csv_path.stem.replace("steamdb_chart_", ""))
        except ValueError:
            continue
        if roster_ids is not None and app_id not in roster_ids:
            continue
        try:
            df = pd.read_csv(csv_path, encoding="utf-8-sig")
            df.columns = [c.strip().strip('"') for c in df.columns]
            df["DateTime"] = pd.to_datetime(df["DateTime"], errors="coerce", utc=True)
            df = df.dropna(subset=["DateTime"])
            df["Players"] = pd.to_numeric(df["Players"], errors="coerce")
            df = df.dropna(subset=["Players"])
            df = df.sort_values("DateTime")
            raw[app_id] = df[["DateTime", "Players"]].reset_index(drop=True)
        except Exception:
            pass
    return raw


def compute_period_diff(
    raw_data: dict[int, pd.DataFrame],
    live_ccu: dict[int, int],
    days: int = 7,
    tolerance_hours: int = 96,
) -> dict[int, dict]:
    """
    Compare live CCU against a reference point for `days` ago.

    Reference priority:
      1. Saved CCU snapshot from ~`days` ago (ccu_snapshots.json)
      2. Raw CSV row closest to `days` ago (within tolerance)
      3. Last row in the CSV (fallback when no recent data exists)

    Returns {app_id: {prev_ccu, delta, delta_pct, period_label, source}}
    """
    now      = pd.Timestamp.utcnow()
    target   = now - pd.Timedelta(days=days)
    tol      = pd.Timedelta(hours=tolerance_hours)
    snapshots = load_ccu_snapshots()
    result   = {}

    all_app_ids = set(raw_data.keys()) | set(live_ccu.keys())

    for app_id in all_app_ids:
        if app_id not in live_ccu:
            continue
        curr_ccu = live_ccu[app_id]
        prev_ccu = None
        source   = ""
        ref_label = ""

        # Priority 1: saved JSON snapshot near target date
        snap_val = find_snapshot_near(snapshots, app_id, days, tolerance_hours)
        if snap_val is not None:
            prev_ccu  = snap_val
            source    = "snapshot"
            ref_label = f"~{days}d ago (saved)"

        # Priority 2: CSV row closest to target date
        if prev_ccu is None and app_id in raw_data:
            df = raw_data[app_id]
            if not df.empty:
                dts = df["DateTime"]
                if dts.dt.tz is None:
                    dts = dts.dt.tz_localize("UTC")
                diff = (dts - target).abs()
                idx  = diff.idxmin()
                gap  = diff[idx]
                if gap <= tol:
                    prev_ccu  = int(df.loc[idx, "Players"])
                    ref_dt    = df.loc[idx, "DateTime"]
                    source    = "csv_exact"
                    ref_label = (ref_dt.strftime(f"{days}d ago (%d %b %Y)")
                                 if hasattr(ref_dt, "strftime") else f"{days}d ago")

        # Priority 3: last CSV row (when no data is old enough)
        if prev_ccu is None and app_id in raw_data:
            last = get_csv_last_ccu(raw_data, app_id)
            if last is not None:
                prev_ccu  = last
                source    = "csv_last"
                ref_label = "latest CSV row"

        if prev_ccu is None or prev_ccu == 0:
            continue

        delta     = curr_ccu - prev_ccu
        delta_pct = (delta / prev_ccu * 100)
        result[app_id] = {
            "prev_ccu":     prev_ccu,
            "curr_ccu":     curr_ccu,
            "delta":        delta,
            "delta_pct":    delta_pct,
            "period_label": ref_label,
            "source":       source,
        }
    return result

# ─────────────────────────────────────────────────────────────
# REPORT ARCHIVE  (weekly auto-archive + monthly comparison)
# Reports saved to /data/report_archive/ as JSON files.
# ─────────────────────────────────────────────────────────────

def _archive_dir() -> Path:
    base = DATA_DIR if DATA_DIR and DATA_DIR.exists() else Path(__file__).parent
    d = base / "report_archive"
    try:
        d.mkdir(exist_ok=True)
    except Exception:
        pass
    return d

def list_archived_reports() -> list[dict]:
    """Return archive metadata dicts, newest first."""
    reports = []
    try:
        for p in sorted(_archive_dir().glob("report_*.json"), reverse=True):
            try:
                with open(p) as f:
                    meta = json.load(f)
                meta["_filename"] = p.name
                reports.append(meta)
            except Exception:
                pass
    except Exception:
        pass
    return reports

def save_report_to_archive(report_md: str, label: str, genre: str,
                            ccu_data: list[dict]) -> str | None:
    """Save report to archive. Returns filename or None on failure."""
    date_str = datetime.utcnow().strftime("%Y-%m-%d")
    filename = f"report_{date_str}_{genre.lower()}.json"
    # Recompute WoW for archival if available
    _wow_snap = {}
    try:
        _raw_snap = load_all_raw()
        _lmap     = {r["app_id"]: r["ccu"] for r in ccu_data}
        _wow_snap = compute_period_diff(_raw_snap, _lmap, days=7)
    except Exception:
        pass

    payload  = {
        "date":        date_str,
        "label":       label,
        "genre":       genre,
        "report_md":   report_md,
        "ccu_snapshot": [{
            "app_id": r["app_id"], "name": r["name"], "ccu": r["ccu"],
            "yoy":    r.get("yoy", "N/A"),
            "mom":    (r.get("hist_summary") or {}).get("mom_trend", "—"),
            "wow":    (f"+{round(_wow_snap[r['app_id']]['delta_pct'])}%"
                       if r["app_id"] in _wow_snap and _wow_snap[r["app_id"]]["delta_pct"] >= 0
                       else f"{round(_wow_snap[r['app_id']]['delta_pct'])}%"
                       if r["app_id"] in _wow_snap else "N/A"),
        } for r in ccu_data],
        "archived_at": datetime.utcnow().isoformat(),
    }
    try:
        with open(_archive_dir() / filename, "w") as f:
            json.dump(payload, f, indent=2)
        return filename
    except Exception:
        return None

def load_archived_report(filename: str) -> dict | None:
    """Load a full archived report by filename."""
    try:
        with open(_archive_dir() / filename) as f:
            return json.load(f)
    except Exception:
        return None

def should_auto_archive(genre: str) -> bool:
    """True if today is Monday and this week+genre hasn't been archived yet."""
    today = datetime.utcnow()
    if today.weekday() != 0:  # Monday = 0
        return False
    week_str = today.strftime("%Y-%m-%d")
    return not (_archive_dir() / f"report_{week_str}_{genre.lower()}.json").exists()

def data_hash(ccu_data: list[dict]) -> str:
    """Stable hash of CCU values for cache-key purposes."""
    import hashlib, json
    payload = json.dumps([{"id": r["app_id"], "ccu": r["ccu"]} for r in ccu_data], sort_keys=True)
    return hashlib.md5(payload.encode()).hexdigest()[:12]

# ─────────────────────────────────────────────────────────────
# DAILY CACHE  (persists ccu_data + ai_report across refreshes)
# Stored at /data/daily_cache.json. Re-fetched once per 24h or
# when roster/genre changes. Force-refresh via sidebar button.
# ─────────────────────────────────────────────────────────────

CACHE_TTL_HOURS = 24

def _cache_path() -> Path:
    base = DATA_DIR if DATA_DIR and DATA_DIR.exists() else Path(__file__).parent
    return base / "daily_cache.json"

def load_daily_cache(genre: str, roster_ids: list[int]) -> dict | None:
    """Load cache if < CACHE_TTL_HOURS old and roster/genre match."""
    try:
        p = _cache_path()
        if not p.exists():
            return None
        with open(p) as f:
            cache = json.load(f)
        cached_at = datetime.fromisoformat(cache.get("cached_at", "2000-01-01"))
        age_hours = (datetime.utcnow() - cached_at).total_seconds() / 3600
        if age_hours > CACHE_TTL_HOURS:
            return None
        if cache.get("genre") != genre:
            return None
        if sorted(cache.get("roster_ids", [])) != sorted(roster_ids):
            return None
        return cache
    except Exception:
        return None

def save_daily_cache(genre: str, roster_ids: list[int],
                     ccu_data: list[dict], ai_report: str,
                     report_label: str = "") -> None:
    """Persist fetch results to disk so next page load skips re-fetching."""
    try:
        payload = {
            "cached_at":    datetime.utcnow().isoformat(),
            "genre":        genre,
            "roster_ids":   roster_ids,
            "ccu_data":     ccu_data,
            "ai_report":    ai_report,
            "report_label": report_label,
        }
        with open(_cache_path(), "w") as f:
            json.dump(payload, f)
    except Exception:
        pass

def cache_age_str() -> str:
    """Human-readable string for how old the current cache is."""
    try:
        p = _cache_path()
        if not p.exists():
            return ""
        with open(p) as f:
            cached_at = datetime.fromisoformat(json.load(f).get("cached_at", ""))
        delta = datetime.utcnow() - cached_at
        mins  = int(delta.total_seconds() // 60)
        if mins < 60:
            return f"{mins}m ago"
        hours = mins // 60
        if hours < 24:
            return f"{hours}h {mins % 60}m ago"
        return f"{delta.days}d ago"
    except Exception:
        return ""


# ─────────────────────────────────────────────────────────────
# BACKGROUND SCHEDULER  (Monday 09:00 UTC auto-archive)
# Uses APScheduler BackgroundScheduler so it fires even when no
# user is actively on the page — as long as the Streamlit process
# is alive (e.g. Streamlit Cloud always-on).
# ─────────────────────────────────────────────────────────────

def _run_monday_archive() -> None:
    """Called by the scheduler every Monday at 09:00 UTC.
    Fetches live CCU for all genres, generates a report via Claude,
    and saves it to the archive.  Requires ANTHROPIC_API_KEY in
    st.secrets (accessed via env var fallback)."""
    import os
    try:
        api_key = os.environ.get("ANTHROPIC_API_KEY", "")
        if not api_key:
            return  # no key configured — skip silently

        # Build full BOTH roster
        both_ids    = list(dict.fromkeys(FPS_ROSTER_IDS + TPS_ROSTER_IDS))
        roster      = [{"app_id": a, **GAME_CATALOG[a]} for a in both_ids if a in GAME_CATALOG]
        historical  = load_all_historical()
        raw_data    = load_all_raw()
        results     = []

        for game in roster:
            ccu  = fetch_ccu(game["app_id"])
            hist_df  = historical.get(game["app_id"])
            has_hist = hist_df is not None and not hist_df.empty
            if has_hist:
                yoy_str, yoy_pct = compute_yoy(hist_df)
                hist_summary = get_historical_summary(hist_df)
            else:
                hist_summary = {}
                yoy_str, yoy_pct = "N/A", 0
            ss = fetch_steamspy(game["app_id"])
            pos = ss.get("positive", 0) or 0
            neg = ss.get("negative", 0) or 0
            tot = pos + neg
            review_pct = round(pos / tot * 100) if tot else None
            if review_pct is None:
                review_pct = fetch_steam_reviews(game["app_id"])
            if not ccu and game["app_id"] in raw_data and not raw_data[game["app_id"]].empty:
                ccu = int(raw_data[game["app_id"]].dropna(subset=["Players"])["Players"].iloc[-1])
            results.append({
                **game,
                "ccu": ccu or 0, "ccu_from_csv": False, "ccu_live": ccu is not None,
                "yoy": yoy_str, "yoy_val": yoy_pct,
                "has_hist": has_hist, "hist_summary": hist_summary,
                "avg_2w_hrs": 0, "review_pct": review_pct,
                "pos_reviews": pos, "neg_reviews": neg,
            })
            time.sleep(0.4)

        results.sort(key=lambda x: x["ccu"], reverse=True)
        save_ccu_snapshot(results)

        # Compute WoW and inject into session-like dict for prompt builder
        live_ccu_map = {r["app_id"]: r["ccu"] for r in results}
        wow_diff     = compute_period_diff(raw_data, live_ccu_map, days=7)

        # Build prompt manually (can't use st.session_state in background thread)
        from datetime import datetime as _dt, timezone as _tz
        today_str = _dt.now(_tz.utc).strftime("%B %d, %Y")
        rows = []
        for rank, r in enumerate(results[:25], 1):
            hs   = r.get("hist_summary", {})
            wow_d = wow_diff.get(r["app_id"])
            wow_s = (f"+{round(wow_d['delta_pct'])}%" if wow_d and wow_d["delta_pct"] > 0
                     else f"{round(wow_d['delta_pct'])}%" if wow_d else "N/A")
            rows.append(
                f"{rank}. {r['name']} ({r['sub']}): {r['ccu']:,} CCU | "
                f"WoW {wow_s} | MoM {hs.get('mom_trend','—')} | "
                f"YoY {r.get('yoy','N/A')} | Review {r.get('review_pct','?')}%"
            )
        prompt = (
            f"You are producing SEGA's internal weekly shooter market intelligence report "
            f"for the week of {today_str}.\n\n"
            f"LIVE CCU SNAPSHOT — {today_str}:\n" + "\n".join(rows) +
            "\n\n---\n\n## SECTION 1: EXECUTIVE SUMMARY\n\n"
            "Write 150–200 words covering:\n"
            "- Overall market mood this week (Rising / Flat / Declining)\n"
            "- 3 headline bullet findings\n"
            "- Story of the Week\n\n---\n\n"
            "## SECTION 2: SHOOTERS RANKED BY CCU\n\n"
            "| Rank | Title | Sub-genre | Live CCU | WoW | MoM | YoY | Review Score | Notes |\n\n"
            "Rules:\n- Use ONLY the CCU figures provided.\n"
            "- Notes: one short observation per title.\n"
            "- Flag YoY > +50% or YoY < -30% in Notes.\n\n"
            "Do not write Section 3 or beyond."
        )

        import anthropic as _anth_sched
        client = _anth_sched.AnthropicBedrock(
            aws_access_key=st.secrets.get("AWS_ACCESS_KEY_ID_API", ""),
            aws_secret_key=st.secrets.get("AWS_SECRET_ACCESS_KEY_API", ""),
            aws_region=st.secrets.get("AWS_BEDROCK_REGION", "us-east-1"),
        )
        resp   = client.messages.create(
            model="us.anthropic.claude-sonnet-4-6",
            max_tokens=6000,
            messages=[{"role": "user", "content": prompt}],
        )
        report_md = resp.content[0].text

        # Save for each genre (BOTH + FPS + TPS)
        for _genre in ["BOTH", "FPS", "TPS"]:
            if should_auto_archive(_genre):
                label = f"Weekly Report — {_genre} — {today_str}"
                save_report_to_archive(report_md, label, _genre, results)

    except Exception as _exc:
        # Log to a file so failures are discoverable
        try:
            log_path = (_archive_dir() / "scheduler.log")
            with open(log_path, "a") as _lf:
                _lf.write(f"{datetime.utcnow().isoformat()} ERROR: {_exc}\n")
        except Exception:
            pass


_scheduler_started = False   # module-level flag — common.py is imported once per process

def _start_scheduler() -> None:
    """Start the APScheduler background thread (idempotent — safe to call from
    every page's import of common.py; the module-level flag above ensures the
    actual scheduler only starts once per process, regardless of how many
    pages import this module)."""
    global _scheduler_started
    try:
        from apscheduler.schedulers.background import BackgroundScheduler
        from apscheduler.triggers.cron import CronTrigger

        if _scheduler_started:
            return
        _scheduler_started = True

        scheduler = BackgroundScheduler(timezone="UTC")
        scheduler.add_job(
            _run_monday_archive,
            CronTrigger(day_of_week="mon", hour=9, minute=0, timezone="UTC"),
            id="monday_archive",
            replace_existing=True,
            misfire_grace_time=3600,  # fire up to 1h late if process was down
        )
        scheduler.start()

        import atexit
        atexit.register(scheduler.shutdown)

    except ImportError:
        pass  # APScheduler not installed — scheduler silently disabled
    except Exception:
        pass  # Already running or other issue — fail silently


# Start scheduler on module load
_start_scheduler()

# ─────────────────────────────────────────────────────────────
# CCU SNAPSHOT PERSISTENCE
# Saves live CCU readings to /data/ccu_snapshots.json on every
# fetch. Used as the WoW reference when the CSV has no matching
# row for exactly 7 days ago.
# ─────────────────────────────────────────────────────────────

def _snapshot_path() -> Path:
    d = DATA_DIR if DATA_DIR and DATA_DIR.exists() else Path(__file__).parent
    return d / "ccu_snapshots.json"

def load_ccu_snapshots() -> list[dict]:
    """Load all saved CCU snapshots [{ts, data:{app_id:ccu}}]."""
    p = _snapshot_path()
    try:
        if p.exists():
            with open(p, "r") as f:
                return json.load(f)
    except Exception:
        pass
    return []

def save_ccu_snapshot(ccu_data: list[dict]) -> None:
    """Append a new snapshot entry with the current UTC timestamp.

    Each entry stores:
      data    — {app_id: live_ccu}              for WoW diff and snapshot YoY
      reviews — {app_id: total_review_count}     for review velocity (pos + neg)

    Retention is 400 days so compute_yoy_from_snapshots() has a full 12-month
    window plus a 30-day tolerance band on either side.
    """
    snapshots = load_ccu_snapshots()
    entry = {
        "ts":      datetime.utcnow().isoformat(),
        "data":    {str(r["app_id"]): r["ccu"] for r in ccu_data},
        "reviews": {
            str(r["app_id"]): (r.get("pos_reviews", 0) or 0) + (r.get("neg_reviews", 0) or 0)
            for r in ccu_data
        },
    }
    snapshots.append(entry)
    # 400-day window supports real YoY (12 months + ±30 day tolerance)
    cutoff    = (datetime.utcnow() - timedelta(days=400)).isoformat()
    snapshots = [s for s in snapshots if s["ts"] >= cutoff]
    try:
        with open(_snapshot_path(), "w") as f:
            json.dump(snapshots, f)
    except Exception:
        pass  # read-only filesystem — fail silently

def find_snapshot_near(snapshots: list[dict], app_id: int, days: int = 7,
                       tolerance_hours: int = 96) -> int | None:
    """Return the CCU for app_id from the snapshot closest to `days` ago.
    Returns None if no snapshot is within tolerance."""
    if not snapshots:
        return None
    target = datetime.utcnow() - timedelta(days=days)
    tol    = timedelta(hours=tolerance_hours)
    best   = None
    best_gap = None
    for s in snapshots:
        try:
            ts  = datetime.fromisoformat(s["ts"])
            gap = abs(ts - target)
            ccu = s["data"].get(str(app_id))
            if ccu is None:
                continue
            if best_gap is None or gap < best_gap:
                best_gap = gap
                best     = int(ccu)
        except Exception:
            continue
    if best_gap is not None and best_gap <= tol:
        return best
    return None

def get_csv_last_ccu(raw_data: dict[int, pd.DataFrame], app_id: int) -> int | None:
    """Return the most recent Players value from the raw CSV for app_id."""
    df = raw_data.get(app_id)
    if df is None or df.empty:
        return None
    try:
        return int(df.dropna(subset=["Players"])["Players"].iloc[-1])
    except Exception:
        return None

PRESET_QUERIES = [
    {
        "id": "ccu_mecha",
        "label": "CCU Trends & Mecha-Shooter Demand",
        "tag": "Market",
        "tag_class": "tag",
        "desc": "Analyze the top 10 shooters on Steam and compare CCU trends to last year. What does this say about current demand for mecha-shooters?",
        "prompt_key": "ccu_mecha",
    },
    {
        "id": "competitive_gap",
        "label": "Competitive Gap Analysis — Top 5",
        "tag": "Strategy",
        "tag_class": "tag amber",
        "desc": "Identify the exact mechanics, monetisation patterns, and content cadences separating the top 5 by CCU. Surfaces addressable product gaps SEGA could exploit.",
        "prompt_key": "competitive_gap",
    },
    {
        "id": "table_stakes",
        "label": "2026 Netcode & Server Table Stakes",
        "tag": "Tech",
        "tag_class": "tag amber",
        "desc": "What are the non-negotiable 'table stakes' for a competitive shooter in 2026 regarding netcode and server architecture to satisfy Western competitive integrity standards?",
        "prompt_key": "table_stakes",
    },
    {
        "id": "social_metrics",
        "label": "Social Media Metrics for Day-1 Success",
        "tag": "Social",
        "tag_class": "tag pos",
        "desc": "Based on recent investor reports and market data, what are the primary social media metrics to track to predict a new shooter's Day 1 success?",
        "prompt_key": "social_metrics",
    },
    {
        "id": "weekly_report",
        "label": "Weekly Retention & Engagement Report Template",
        "tag": "Report",
        "tag_class": "tag purple",
        "desc": "Create a template for a weekly market report that tracks retention and engagement KPIs across the top 100 shooters, highlighting any 'breakout' indie titles.",
        "prompt_key": "weekly_report",
    },
]

# ─────────────────────────────────────────────────────────────
# STEAM LIVE CCU FETCH
# ─────────────────────────────────────────────────────────────

STEAMSPY_URL = "https://steamspy.com/api.php"

@st.cache_data(ttl=300, show_spinner=False)
def fetch_ccu(app_id: int) -> int | None:
    """Fetch live concurrent player count from the Steam public API.
    Retries transient failures (timeout, connection error, non-2xx) up to
    3 times with backoff before giving up — see _http_get_with_retry()."""
    r = _http_get_with_retry(CCU_URL, {"appid": app_id}, timeout=8)
    if r is None:
        return None
    try:
        return r.json().get("response", {}).get("player_count")
    except Exception:
        return None

@st.cache_data(ttl=3600, show_spinner=False)
def fetch_steam_reviews(app_id: int) -> int | None:
    """Fallback: fetch all-time review score from Steam store API."""
    r = _http_get_with_retry(
        f"https://store.steampowered.com/appreviews/{app_id}",
        {"json": 1, "language": "all", "review_type": "all", "purchase_type": "all"},
        timeout=8,
    )
    if r is None:
        return None
    try:
        qs = r.json().get("query_summary", {})
        total = qs.get("total_reviews", 0) or 0
        pos   = qs.get("total_positive", 0) or 0
        if total > 0:
            return round(pos / total * 100)
    except Exception:
        pass
    return None

@st.cache_data(ttl=3600, show_spinner=False)
def fetch_steamspy(app_id: int) -> dict:
    """
    Fetch game data from SteamSpy (no API key required, updates daily).
    Returns fields including:
      - average_forever: avg playtime (mins) all-time
      - average_2weeks:  avg playtime (mins) past 2 weeks
      - owners:          estimated owner band e.g. '1,000,000 .. 2,000,000'
      - positive / negative: review counts
    """
    r = _http_get_with_retry(
        STEAMSPY_URL, {"request": "appdetails", "appid": app_id}, timeout=12,
    )
    if r is None:
        return {}
    try:
        return r.json()
    except Exception:
        return {}

def parse_yoy_from_steamspy(ss: dict) -> tuple[str, int]:
    """
    Derive a YoY proxy from SteamSpy's playtime data.
    Compares average_2weeks (recent engagement) to average_forever (all-time baseline).
    A ratio > 1 implies the game is played MORE than its historical average → growing.
    A ratio < 1 implies declining engagement.
    Returns (display_str, numeric_pct).
    """
    avg_all  = ss.get("average_forever", 0) or 0
    avg_2w   = ss.get("average_2weeks",  0) or 0

    if avg_all == 0:
        return "N/A", 0

    # Normalise: if 2-week avg is X% of all-time avg, that's the engagement index
    ratio = avg_2w / avg_all
    # Map to a YoY-style percentage: ratio 1.0 = flat (0%), 1.5 = +50%, 0.5 = -50%
    pct = round((ratio - 1.0) * 100)
    # Cap display at ±99% to avoid misleading outliers
    pct = max(-99, min(99, pct))
    sign = "+" if pct >= 0 else ""
    return f"{sign}{pct}%", pct


def _fetch_one_game(
    game: dict,
    historical: dict[int, pd.DataFrame],
    raw_data: dict[int, pd.DataFrame],
    snapshots: list[dict],
) -> dict:
    """Fetch all external data for a single game.

    Designed to run inside a ThreadPoolExecutor worker.  All network calls are
    wrapped in @st.cache_data (thread-safe in Streamlit ≥ 1.18) so concurrent
    cache-hits are free and only genuine misses hit the network.

    Both Steam CCU and SteamSpy requests are throttled through their own
    semaphores (_STEAM_CCU_SEM / _STEAMSPY_SEM) so a burst of worker threads
    never sends more than a handful of simultaneous requests to either API —
    Steam's public endpoints are generally tolerant, but a single-IP burst of
    a dozen-plus simultaneous requests (what unthrottled concurrency would
    produce) is a much more plausible trigger for rate-limiting or transient
    blocks than the original sequential fetch (one request every ~0.4s) ever
    was. fetch_ccu/fetch_steamspy themselves retry transient failures with
    backoff — see _http_get_with_retry().

    YoY priority: snapshot-based (real) → SteamDB CSV → SteamSpy proxy
    """
    app_id = game["app_id"]

    # ── Steam live CCU ──
    with _STEAM_CCU_SEM:
        ccu = fetch_ccu(app_id)
        time.sleep(0.08)          # polite pause while holding the semaphore

    # ── SteamSpy ── (one call; reused for both YoY proxy and reviews)
    with _STEAMSPY_SEM:
        ss = fetch_steamspy(app_id)
        time.sleep(0.12)          # polite pause while holding the semaphore

    # ── Twitch live viewers ── (graceful no-op when credentials absent)
    twitch_viewers = fetch_twitch_viewers(app_id, game["name"])

    # ── Historical CSV data ──
    hist_df  = historical.get(app_id)
    has_hist = hist_df is not None and not hist_df.empty

    # ── YoY: snapshot → CSV → SteamSpy (priority order) ──────────────────
    # After ~13 months of daily snapshots the real YoY takes precedence.
    # CSV YoY is preferred over SteamSpy while snapshot history builds up.
    snap_yoy = compute_yoy_from_snapshots(snapshots, app_id)
    if snap_yoy is not None:
        yoy_str, yoy_pct = snap_yoy
        yoy_source   = "snapshot"
        hist_summary = get_historical_summary(hist_df) if has_hist else {}
    elif has_hist:
        yoy_str, yoy_pct = compute_yoy(hist_df)
        yoy_source   = "csv"
        hist_summary = get_historical_summary(hist_df)
    else:
        yoy_str, yoy_pct = parse_yoy_from_steamspy(ss)
        yoy_source   = "steamspy"
        hist_summary = {}

    # ── Reviews ──
    pos_reviews = ss.get("positive", 0) or 0
    neg_reviews = ss.get("negative", 0) or 0
    total_rev   = pos_reviews + neg_reviews
    review_pct  = round(pos_reviews / total_rev * 100) if total_rev else None
    if review_pct is None:
        review_pct = fetch_steam_reviews(app_id)   # fallback to Steam store API

    # ── Review velocity (reviews/day over last 7 days from snapshots) ──
    review_velocity = compute_review_velocity(snapshots, app_id, days=7)

    avg_2w_hrs = round((ss.get("average_2weeks", 0) or 0) / 60, 1)

    # ── CSV CCU fallback ── (when live API returns 0, e.g. Deadlock beta period)
    ccu_from_csv = False
    if not ccu and app_id in raw_data and not raw_data[app_id].empty:
        ccu = int(raw_data[app_id].dropna(subset=["Players"])["Players"].iloc[-1])
        ccu_from_csv = True

    return {
        **game,
        "ccu":             ccu if ccu else 0,
        "ccu_from_csv":    ccu_from_csv,
        "ccu_live":        ccu is not None,
        "twitch_viewers":  twitch_viewers,           # int | None
        "yoy":             yoy_str,
        "yoy_val":         yoy_pct,
        "yoy_source":      yoy_source,               # "snapshot" | "csv" | "steamspy"
        "has_hist":        has_hist,
        "hist_summary":    hist_summary,
        "avg_2w_hrs":      avg_2w_hrs,
        "review_pct":      review_pct,
        "review_velocity": review_velocity,          # float | None  (reviews/day)
        "pos_reviews":     pos_reviews,
        "neg_reviews":     neg_reviews,
    }

# ─────────────────────────────────────────────────────────────
# REPORT HELPERS (HTML + PDF)
# ─────────────────────────────────────────────────────────────

def report_to_html(md_text: str) -> str:
    body = _md_lib.markdown(md_text, extensions=["tables", "fenced_code"]) if MARKDOWN_AVAILABLE else md_text.replace("\n", "<br>")
    return f"""<!DOCTYPE html>
<html lang="en">
<head>
<meta charset="utf-8">
<title>SEGA Shooter Intel — Report</title>
<style>
body{{background:#0a0c1a;color:#eef0fa;font-family:'Segoe UI',sans-serif;max-width:860px;margin:40px auto;padding:0 24px;line-height:1.7}}
h1,h2,h3{{font-weight:900;letter-spacing:-.02em;color:#fff}}
h1{{font-size:2rem;border-bottom:2px solid #4080ff;padding-bottom:.4rem}}
h2{{font-size:1.4rem;color:#7ab0ff;margin-top:2rem}}
h3{{font-size:1.1rem;color:#b8bcd4}}
code,pre{{background:#141728;border:1px solid #232640;border-radius:4px;padding:.15em .4em;font-size:.88em;color:#4080ff}}
pre code{{padding:0}}
table{{border-collapse:collapse;width:100%}}
th,td{{border:1px solid #232640;padding:.5rem .8rem;text-align:left}}
th{{background:#141728;color:#7ab0ff;font-size:.78rem;letter-spacing:.1em;text-transform:uppercase}}
hr{{border:none;border-top:1px solid #232640;margin:1.5rem 0}}
.footer{{margin-top:3rem;padding-top:.8rem;border-top:1px solid #232640;font-size:.72rem;color:#5a5f82}}
</style>
</head>
<body>
{body}
<div class="footer">SEGA Shooter Intelligence — Internal analytics use only</div>
</body>
</html>"""

def report_to_pdf(md_text: str) -> bytes | None:
    if not _REPORTLAB_AVAILABLE:
        return None
    buf = _rl_io.BytesIO()
    doc = SimpleDocTemplate(buf, pagesize=A4,
                            leftMargin=2*cm, rightMargin=2*cm,
                            topMargin=2*cm, bottomMargin=2*cm)
    styles = getSampleStyleSheet()
    _h1 = ParagraphStyle("h1", parent=styles["Heading1"],
                         fontSize=18, textColor=_rl_colors.HexColor("#0f3460"),
                         spaceAfter=8, spaceBefore=14)
    _h2 = ParagraphStyle("h2", parent=styles["Heading2"],
                         fontSize=14, textColor=_rl_colors.HexColor("#0f3460"),
                         spaceAfter=6, spaceBefore=10)
    _h3 = ParagraphStyle("h3", parent=styles["Heading3"],
                         fontSize=12, textColor=_rl_colors.HexColor("#1a1a2e"),
                         spaceAfter=4, spaceBefore=8)
    _body = ParagraphStyle("body", parent=styles["Normal"],
                           fontSize=10, leading=15, spaceAfter=6)
    _bullet = ParagraphStyle("bullet", parent=_body,
                             leftIndent=16, bulletIndent=6, spaceAfter=3)
    _code = ParagraphStyle("code", parent=styles["Code"],
                           fontSize=8, leading=12,
                           backColor=_rl_colors.HexColor("#f0f0f8"),
                           leftIndent=12, rightIndent=12, spaceAfter=6)
    story = []
    in_code, code_lines = False, []
    for line in md_text.split("\n"):
        if line.startswith("```"):
            if in_code:
                story.append(Preformatted("\n".join(code_lines), _code))
                story.append(Spacer(1, 4))
                code_lines = []; in_code = False
            else:
                in_code = True
            continue
        if in_code:
            code_lines.append(line); continue
        if line.startswith("### "):
            story.append(Paragraph(line[4:], _h3))
        elif line.startswith("## "):
            story.append(HRFlowable(width="100%", thickness=0.5,
                         color=_rl_colors.HexColor("#c0c0d8"), spaceAfter=2))
            story.append(Paragraph(line[3:], _h2))
        elif line.startswith("# "):
            story.append(Paragraph(line[2:], _h1))
        elif line.startswith("- ") or line.startswith("* "):
            t = line[2:].replace("**", "<b>", 1).replace("**", "</b>", 1)
            story.append(Paragraph(f"• {t}", _bullet))
        elif line.strip() in ("---", "***", "___"):
            story.append(HRFlowable(width="100%", thickness=0.5,
                         color=_rl_colors.HexColor("#c0c0d8")))
            story.append(Spacer(1, 4))
        elif line.strip() == "":
            story.append(Spacer(1, 6))
        else:
            import re as _re
            t = _re.sub(r"\*\*(.+?)\*\*", r"<b>\1</b>", line)
            t = _re.sub(r"\*(.+?)\*",   r"<i>\1</i>", t)
            t = _re.sub(r"`(.+?)`", r"<font name='Courier'>\1</font>", t)
            story.append(Paragraph(t, _body))
    doc.build(story)
    return buf.getvalue()

# ─────────────────────────────────────────────────────────────
# CLAUDE PROMPTS
# ─────────────────────────────────────────────────────────────

def build_system_prompt(language: str = "English") -> str:
    lang_instruction = (
        " IMPORTANT: Write your entire response in Japanese (日本語). "
        "Use professional business Japanese suitable for senior management. "
        "All section headers, bullet points, tables (including table column headers), "
        "and analysis must be in Japanese. "
        "Game titles may be kept in their original English/romanised form where commonly known. "
        "Never switch back to English partway through — the full report must be Japanese from "
        "the first line to the last."
    ) if language == "Japanese" else ""
    return (
        "You are a senior games market analyst at SEGA's internal strategy team. "
        "You specialise in the competitive shooter genre across Steam, console, and mobile. "
        "Your analysis is data-driven, commercially sharp, and directly actionable for a publishing team. "
        "Use markdown for all output. Use headers, bullet points, tables, and bold highlights where appropriate. "
        "Be specific — cite titles, numbers, dates, and named competitors whenever possible. "
        "Avoid vague generalisations. Outputs will be read by product leads and senior management."
        + lang_instruction
    )

def build_ccu_mecha_prompt(ccu_data: list[dict], genre: str = "FPS") -> str:
    rows = []
    for r in ccu_data:
        hs = r.get("hist_summary", {})
        data_src = "SteamDB CSV" if r.get("has_hist") else "SteamSpy proxy"
        row = (
            f"- **{r['name']}** ({r['sub']}): {r['ccu']:,} live CCU | "
            f"YoY {r.get('yoy','N/A')} [{data_src}] | "
            f"Peak ever {hs.get('peak_ever', '?'):,} | " if hs.get('peak_ever') else
            f"- **{r['name']}** ({r['sub']}): {r['ccu']:,} live CCU | "
            f"YoY {r.get('yoy','N/A')} [{data_src}] | "
        )
        if hs:
            row += (
                f"Peak 12m {hs.get('peak_12m','?'):,} | "
                f"Avg CCU 12m {hs.get('avg_12m','?'):,} | "
                f"MoM {hs.get('mom_trend','—')} | "
                f"{hs.get('months_data',0)} months of data"
            )
        row += f" | Review score {r.get('review_pct','?')}%"
        rows.append(row)
    rows_str = "\n".join(rows)

    if genre == "TPS":
        genre_label    = "Third-Person Shooter (TPS)"
        genre_context  = "third-person perspective, cover mechanics, and over-the-shoulder combat"
        mecha_section  = (
            "3. **Mecha-Shooter Demand (TPS)** — Assess demand for third-person mecha titles "
            "(e.g. Armored Core VI in TPS mode, MechWarrior, Daemon X Machina). "
            "What does the TPS CCU data imply about player appetite for mech-based third-person action?"
        )
        sega_question  = (
            "4. **Strategic Implications** — If SEGA were to greenlight a new third-person mecha-shooter today, "
            "what CCU targets are realistic for Year 1 based on comparable TPS titles? "
            "What monetisation and player-acquisition strategies are implied by this data?"
        )
    else:
        genre_label    = "First-Person Shooter (FPS)"
        genre_context  = "first-person perspective and direct gunplay"
        mecha_section  = (
            "3. **Mecha-Shooter Demand (FPS)** — Focus on titles with mech/robot/exosuit themes "
            "(e.g. Titanfall 2, Armored Core VI). What does the FPS CCU data say about current demand "
            "for mecha-shooters in first-person? Is the sub-genre growing or niche?"
        )
        sega_question  = (
            "4. **Strategic Implications** — If SEGA were to greenlight a new first-person mecha-shooter today, "
            "what CCU targets are realistic for Year 1? What player acquisition strategies are implied by the data?"
        )

    return f"""## Task: CCU Trend Analysis — Top {genre_label} Titles on Steam

You have been provided with **live Steam CCU data** and **real historical SteamDB data** (where available) for the top {genre_label} titles as of today. This roster focuses on games with {genre_context}. YoY comparisons marked [SteamDB CSV] are genuine same-month comparisons from raw data; those marked [SteamSpy proxy] are engagement-based estimates.

### Live CCU Snapshot + Historical Context ({genre_label})
{rows_str}

### Questions to answer (structure your report with these sections):

1. **CCU Overview** — Rank the titles. Which are growing YoY? Which are declining? Identify clear winners and losers within the {genre_label} space.
2. **{genre_label} Genre Health** — What do these numbers collectively say about {genre_label} health on Steam in 2026? Are players migrating between sub-genres?
{mecha_section}
{sega_question}
5. **Recommended Watch List** — List 3 {genre_label} titles to monitor closely over the next 90 days and why.

Be quantitative. Use tables where useful."""

def build_table_stakes_prompt() -> str:
    return """## Task: 2026 Technical Table Stakes — Competitive Shooters

Western PC/console players in 2026 have extremely high expectations for competitive integrity. This analysis is for SEGA's internal game design and engineering teams evaluating a potential new competitive shooter.

### Provide a comprehensive analysis covering:

1. **Netcode Fundamentals**
   - What server tick rates are now considered minimum vs. best-in-class? (Reference CS2, Valorant, Apex, R6 Siege)
   - Client-side prediction vs. server reconciliation — what is the acceptable latency ceiling for Western esports audiences?
   - Rollback netcode vs. delay-based: which is relevant for shooters and why?

2. **Server Infrastructure**
   - Regional server coverage expectations (NA, EU, APAC minimums)
   - Server ownership vs. peer-to-peer: what is the current industry standard?
   - DDoS protection requirements for ranked/competitive modes

3. **Anti-Cheat**
   - Which anti-cheat solutions are now considered table stakes vs. competitive differentiators?
   - Kernel-level vs. user-level: player sentiment trade-offs
   - Replay validation and server-side hit detection requirements

4. **Competitive Integrity Features**
   - Ranked system design standards (MMR transparency, placement matches, ranked resets)
   - In-game reporting and review workflows
   - Tournament/LAN mode requirements for esports potential

5. **Post-Launch Patch Standards**
   - Acceptable cadence for balance patches and hotfixes
   - Communication standards players expect (patch notes format, dev blogs)

6. **Cost Estimates** — Rough infrastructure budget ranges (low/mid/high tier) for a title targeting 10K, 100K, and 1M peak CCU.

Format as a structured technical brief with a clear PASS/FAIL checklist at the end."""

def build_social_metrics_prompt() -> str:
    return """## Task: Social Media Predictive Metrics — Shooter Day 1 Success

This analysis synthesises findings from recent investor reports (Embracer, EA, Take-Two, Krafton, Nexon), GDC talks, and published post-mortems to identify the most reliable social media metrics for predicting a new shooter's Day 1 commercial performance.

### Analyse and structure your report as follows:

1. **Pre-Launch Predictive Signals (T-90 to T-0)**
   - Which metrics have the highest correlation with Day 1 peak CCU and revenue?
   - Provide specific benchmarks (e.g. "10K wishlist adds per week in final month = X CCU at launch")
   - Differentiate between vanity metrics and actionable leading indicators

2. **Platform-Specific Metrics**
   - **Steam**: Wishlists, review velocity, concurrent viewers on launch day
   - **Twitch/YouTube**: Hours watched, unique streamers, clip virality
   - **X/Twitter**: Hashtag impressions, sentiment ratio, influencer amplification coefficient
   - **TikTok/Shorts**: View-to-wishlist conversion rate, trend longevity
   - **Reddit/Discord**: Community growth rate, DAU/MAU ratio, organic posts vs. seeded content

3. **Red Flags — What Signals a Soft Launch**
   - Which patterns in social data predicted underperformance for recent titles?
   - Name specific failed or underperforming shooters and their pre-launch signals

4. **Green Flags — The Viral Flywheel**
   - What social conditions created outsized launches? (e.g. HELLDIVERS 2, BattleBit Remastered)
   - What was different about their pre-launch social fingerprint?

5. **Recommended KPI Dashboard**
   - Design a 10-metric social dashboard SEGA should build for any new shooter pre-launch
   - Include: metric name, platform, measurement method, target threshold, and refresh cadence

6. **Budget Implication** — Given these metrics, what social/influencer spend is required to hit minimum viable social velocity for a Western competitive shooter launch?"""

def build_weekly_report_prompt(ccu_data: list[dict], language: str = "English") -> str:
    from datetime import datetime, timezone
    today    = datetime.now(timezone.utc)
    date_str = today.strftime("%B %d, %Y")

    # Pull wow_diff from session state if available
    _wow = st.session_state.get("_wow_diff_cache", {})

    rows = []
    for rank, r in enumerate(ccu_data, 1):
        hs     = r.get("hist_summary", {})
        src_map  = {"snapshot": "Snapshot", "csv": "SteamDB", "steamspy": "est."}
        data_src = src_map.get(r.get("yoy_source", "steamspy"), "est.")
        wow_d  = _wow.get(r["app_id"])
        wow_str = f"+{round(wow_d['delta_pct'])}%" if wow_d and wow_d["delta_pct"] > 0 else (f"{round(wow_d['delta_pct'])}%" if wow_d else "N/A")
        twitch  = f"{r['twitch_viewers']:,}" if r.get("twitch_viewers") is not None else "—"
        vel     = f"{r['review_velocity']}/day" if r.get("review_velocity") is not None else "N/A"
        line = (
            f"{rank}. {r['name']} ({r['sub']}): {r['ccu']:,} live CCU | "
            f"Twitch {twitch}v | "
            f"WoW {wow_str} | "
            f"MoM {hs.get('mom_trend','—')} | "
            f"YoY {r.get('yoy','N/A')} [{data_src}] | "
            f"Review {r.get('review_pct','?')}% | Rev/day {vel}"
        )
        if hs.get("peak_12m"):
            line += f" | Peak 12m {hs['peak_12m']:,}"
        rows.append(line)
    rows_str = "\n".join(rows)

    # Table column headers — keep these consistent with the UI language
    if language == "Japanese":
        table_header = "| 順位 | タイトル | サブジャンル | ライブCCU | WoW | MoM | YoY | レビュースコア | コメント |"
        section3_prompt = """## セクション3: SEGAへの戦略的示唆

上記のCCUデータとランキングに基づき、SEGAのパブリッシング・製品リーダーシップ向けに戦略ブリーフを作成してください。

1. **市場の機会窓** — CCUが成長中であるが、まだ市場が固定されていないサブジャンルや特定ニッチを1〜2つ特定してください。おおまかな緊急度も示してください。
2. **競合の脆弱性** — 今週のデータの中で、CCUの低下や指標の弱体化によりプレイヤーが代替を求めているシグナルを示すタイトルを2〜3つ挙げ、その層を引き付けられるゲームのタイプを述べてください。
3. **SEGAのポジショニング** — SEGAのアクション・チームベース・ライブサービスゲームにおけるパブリッシング実績を踏まえ、今週のデータのどのセグメントがSEGAの強みと最も合致するか、またSEGAブランドのタイトルが現在のリーダーと何を差別化すべきかを述べてください。
4. **90日ウォッチリスト** — 次のレポートの早期警告シグナルとして監視すべきタイトルを2つ、追跡する具体的な指標とともに挙げてください。
5. **50字の戦略的提言** — SEGAリーダーシップ向けの簡潔な提言を1文で。

直接的かつ商業的に鋭くまとめ、セクション1と2の具体的なCCU数値を必ず参照してください。"""
        lang_note = (
            "\nLANGUAGE: Write the entire report in Japanese (日本語), including the executive "
            "summary, all table column headers, and the Notes/コメント column. Game titles may stay "
            "in English. Do not switch to English at any point."
        )
    else:
        table_header = "| Rank | Title | Sub-genre | Live CCU | WoW | MoM | YoY | Review Score | Notes |"
        section3_prompt = """## SECTION 3: SEGA STRATEGIC IMPLICATIONS

Based on the CCU data and rankings above, provide a strategic brief for SEGA's publishing and product leadership.

1. **Market Windows** — Identify 1–2 sub-genres or niches where CCU is growing but the market is not yet locked (no single title holds >60% share). Frame each as an opportunity window with a rough urgency level (act within 12 months / 24 months / watch and wait).
2. **Competitive Vulnerabilities** — Name 2–3 titles in this week's data whose declining CCU or weakening metrics signal their player base is ready to migrate. What type of game could credibly capture that audience?
3. **SEGA Positioning** — Given SEGA's publishing history in action, team-based, and live-service games, which segment in this week's data aligns best with SEGA's demonstrated strengths? What would a SEGA-branded entry need to do differently from the current leaders to earn a top-10 CCU position?
4. **90-Day Watchlist** — 2 titles to monitor as early warning signals, with the specific metric to track for each.
5. **One-sentence recommendation** — A single sharp strategic recommendation for SEGA leadership based solely on this week's data.

Be direct and commercially sharp. Reference specific CCU numbers from Sections 1 and 2. Avoid generic observations."""
        lang_note = ""

    return f"""You are producing SEGA's internal weekly shooter market intelligence report for the week of {date_str}.

IMPORTANT: All data below is LIVE as of {date_str}. Use ONLY these figures. Do not reference any other dates. Do not invent or estimate CCU numbers.{lang_note}

LIVE CCU SNAPSHOT — {date_str}:
{rows_str}

Produce the three sections below. Do not add any other sections or commentary.

---

## SECTION 1: EXECUTIVE SUMMARY

Write 150–200 words covering:
- Overall market mood this week (Rising / Flat / Declining) — justify with the data above
- 3 headline bullet findings drawn directly from the numbers above
- Story of the Week: the single most notable move or trend visible in this data

---

## SECTION 2: SHOOTERS RANKED BY CCU

Produce a markdown table with these exact columns:
{table_header}

Rules:
- Use ONLY the CCU figures provided above — do not invent or estimate
- WoW, MoM, YoY: use the values from the data above; N/A if not available
- Notes: one short observation per title based on the data (e.g. "Declining 3 months", "New season spike", "Near all-time peak")
- Flag in Notes any title with YoY > +50% or YoY < -30%
- The table MUST include every title from the snapshot and MUST be complete — never stop partway through a row

---

{section3_prompt}"""

def build_competitive_gap_prompt(ccu_data: list[dict], language: str = "English") -> str:
    """Detailed breakdown of the top-5 titles: mechanics, monetisation, content cadence,
    and the product gaps SEGA could exploit. Written for a product/publishing audience."""
    from datetime import datetime, timezone
    today    = datetime.now(timezone.utc)
    date_str = today.strftime("%B %d, %Y")

    _wow = st.session_state.get("_wow_diff_cache", {})

    rows = []
    for rank, r in enumerate(ccu_data[:5], 1):
        hs   = r.get("hist_summary", {})
        wow_d = _wow.get(r["app_id"])
        wow_str = (f"+{round(wow_d['delta_pct'])}%" if wow_d and wow_d["delta_pct"] > 0
                   else f"{round(wow_d['delta_pct'])}%" if wow_d else "N/A")
        peak = f" | 12m peak {hs['peak_12m']:,}" if isinstance(hs.get("peak_12m"), int) else ""
        rows.append(
            f"{rank}. {r['name']} ({r['sub']})\n"
            f"   CCU: {r['ccu']:,} | F2P: {'Yes' if r.get('f2p') else 'No'} | "
            f"Publisher: {r['publisher']} | Year: {r.get('year','?')}\n"
            f"   WoW {wow_str} | MoM {hs.get('mom_trend','—')} | "
            f"YoY {r.get('yoy','N/A')} | Review {r.get('review_pct','?')}%{peak}"
        )
    rows_str = "\n\n".join(rows)

    lang_note = (
        "\n\nLANGUAGE: Write the entire analysis in Japanese (日本語) using professional "
        "business Japanese. Game titles may stay in English. Do not switch to English at any point."
    ) if language == "Japanese" else ""

    return f"""You are a senior competitive intelligence analyst at SEGA, evaluating the top 5 shooters on Steam as of {date_str}.{lang_note}

TOP 5 SHOOTERS BY LIVE CCU — {date_str}:

{rows_str}

Produce a COMPETITIVE GAP ANALYSIS with the following five sections:

## 1. Mechanics Driving CCU

For each of the top 5 titles, identify the 2–3 core gameplay mechanics or design decisions that most directly explain its current CCU position. Be specific — not "good gunplay" but "variable TTK and destructible cover creating emergent ranked meta". Note where mechanics overlap between titles, as this signals shared audience and potential migration risk.

## 2. Monetisation Patterns

Compare the monetisation models across all 5. Which model (F2P battle pass, premium + DLC, F2P cosmetics-only, etc.) is delivering the best retention-to-revenue outcome, based on the CCU and review data? Flag any outlier where premium pricing is sustaining or growing CCU against F2P competition, and explain why.

## 3. Content Cadence & Live-Service Signals

Infer the content update cadence for each title from its WoW / MoM / YoY trajectory. Which show "season spike" patterns (CCU surges then decays)? Which show a "stable floor" suggesting genuine long-term retention? What cadence frequency does the data suggest is the current player expectation for a top-tier shooter?

## 4. Addressable Gaps

Based on mechanics, monetisation, and cadence analysis above, identify 2–3 specific combination gaps — things the top 5 are NOT providing. Frame each as a one-sentence product brief: "A [sub-genre] title with [mechanic], [monetisation model], and [update cadence] targeting [audience segment] has no direct incumbent in the current top 5."

## 5. SEGA Entry Recommendation

If SEGA were to greenlight a new competitive shooter targeting the Western PC market for a 2027 launch, what single market position does this analysis most clearly support? Specify: sub-genre, 3 non-negotiable mechanics, pricing model, first-year content cadence, and the single incumbent title whose audience represents the primary acquisition target.

Be specific and data-driven throughout. Reference actual CCU numbers, YoY/MoM figures, and review scores from the data above. Avoid generic games industry observations."""



def generate_pptx_bytes(report_md: str, ccu_data: list[dict], label: str) -> bytes | None:
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
        from io import BytesIO
    except ImportError:
        return None

    BG_DARK = RGBColor(0x05, 0x08, 0x18)
    BG_NAVY = RGBColor(0x0D, 0x11, 0x26)
    BLUE    = RGBColor(0x00, 0x57, 0xFF)
    MUTED   = RGBColor(0xC3, 0xC5, 0xD5)
    TEXT    = RGBColor(0xE1, 0xEA, 0xFF)
    POS     = RGBColor(0x20, 0xC6, 0x5A)
    NEG     = RGBColor(0xFF, 0x4D, 0x6D)

    W = Inches(13.33)
    H = Inches(7.5)
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    blank = prs.slide_layouts[6]

    def new_slide():
        return prs.slides.add_slide(blank)

    def bg(slide, color):
        f = slide.background.fill
        f.solid()
        f.fore_color.rgb = color

    def rect(slide, x, y, w, h, color):
        sh = slide.shapes.add_shape(1, x, y, w, h)
        sh.fill.solid()
        sh.fill.fore_color.rgb = color
        sh.line.fill.background()
        return sh

    def tb(slide, x, y, w, h, text, size=14, bold=False,
           color=TEXT, align=PP_ALIGN.LEFT, italic=False):
        box = slide.shapes.add_textbox(x, y, w, h)
        box.word_wrap = True
        tf = box.text_frame
        tf.word_wrap = True
        p  = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size   = Pt(size)
        run.font.bold   = bold
        run.font.italic = italic
        run.font.color.rgb = color
        run.font.name = "Calibri"
        return box

    def stripe(slide):
        rect(slide, 0, 0, Inches(0.18), H, BLUE)

    def fmt_ccu(n):
        if n >= 1_000_000: return f"{n/1_000_000:.2f}M"
        if n >= 1_000:     return f"{n/1_000:.1f}K"
        return str(n)

    # Parse report_md into ## sections
    sections = []
    cur_title, cur_body = "", []
    for line in report_md.splitlines():
        if line.startswith("## "):
            if cur_title:
                sections.append((cur_title, "\n".join(cur_body).strip()))
            cur_title = line.lstrip("# ").strip()
            cur_body  = []
        elif line.startswith("# "):
            pass
        else:
            cur_body.append(line)
    if cur_title:
        sections.append((cur_title, "\n".join(cur_body).strip()))

    # Slide 1: Title
    s = new_slide()
    bg(s, BG_DARK)
    stripe(s)
    rect(s, 0, 0, W, Inches(0.06), BLUE)
    tb(s, Inches(0.42), Inches(1.6),  Inches(9), Inches(1.1), "SHOOTER MARKET", 52, bold=True, color=TEXT)
    tb(s, Inches(0.42), Inches(2.55), Inches(9), Inches(1.1), "INTELLIGENCE",   52, bold=True, color=BLUE)
    tb(s, Inches(0.42), Inches(3.85), Inches(9), Inches(0.5), label, 15, color=MUTED)
    tb(s, Inches(0.42), Inches(4.35), Inches(9), Inches(0.4),
       "SEGA Publishing & Strategy  \u00b7  Steam CCU Analysis", 12, color=MUTED, italic=True)

    # Slide 2: CCU snapshot cards
    top10 = sorted(ccu_data, key=lambda r: r["ccu"], reverse=True)[:10]
    s = new_slide()
    bg(s, BG_NAVY)
    stripe(s)
    rect(s, Inches(0.18), 0, W - Inches(0.18), Inches(0.8), BLUE)
    tb(s, Inches(0.42), Inches(0.12), Inches(10), Inches(0.56),
       "LIVE CCU SNAPSHOT", 18, bold=True, color=TEXT)
    cols = 5
    cw, ch = Inches(2.4), Inches(1.4)
    xoff, yoff, gap = Inches(0.35), Inches(1.0), Inches(0.12)
    for i, r in enumerate(top10):
        col, row = i % cols, i // cols
        cx = xoff + col * (cw + gap)
        cy = yoff + row * (ch + gap)
        card_col = BG_DARK if row == 0 else RGBColor(0x08, 0x0C, 0x1E)
        rect(s, cx, cy, cw, ch, card_col)
        top_col = POS if r.get("yoy_val", 0) >= 0 else NEG
        rect(s, cx, cy, cw, Inches(0.04), top_col)
        tb(s, cx + Inches(0.12), cy + Inches(0.06), cw - Inches(0.14), Inches(0.35),
           r["name"][:22], 9, color=MUTED)
        tb(s, cx + Inches(0.12), cy + Inches(0.38), cw - Inches(0.14), Inches(0.55),
           fmt_ccu(r["ccu"]), 22, bold=True, color=TEXT)
        yoy = r.get("yoy", "N/A")
        yoy_col = POS if str(yoy).startswith("+") else (NEG if str(yoy).startswith("-") else MUTED)
        tb(s, cx + Inches(0.12), cy + Inches(0.92), cw - Inches(0.14), Inches(0.35),
           f"YoY {yoy}", 10, color=yoy_col)

    # Slides 3+: Report sections
    for i, (title, body) in enumerate(sections):
        s = new_slide()
        bg(s, BG_DARK if i % 2 == 0 else BG_NAVY)
        stripe(s)
        rect(s, Inches(0.18), 0, W - Inches(0.18), Inches(0.75), RGBColor(0x00, 0x22, 0x66))
        tb(s, Inches(0.42), Inches(0.1), Inches(11), Inches(0.55),
           title.upper(), 16, bold=True, color=TEXT)
        clean = []
        for line in body.splitlines():
            line = line.strip()
            if not line:
                clean.append("")
                continue
            line = line.lstrip("#").strip()
            if line.startswith("- ") or line.startswith("* "):
                line = "\u2022  " + line[2:]
            clean.append(line)
        body_text = "\n".join(clean).strip()
        if len(body_text) > 1400:
            body_text = body_text[:1400] + "\n\n[\u2026continued in full report]"
        tb(s, Inches(0.42), Inches(0.9), Inches(12.5), Inches(6.3),
           body_text, 13, color=MUTED)

    # Final slide
    s = new_slide()
    bg(s, BG_DARK)
    stripe(s)
    tb(s, Inches(0.42), Inches(2.8), Inches(12), Inches(1.0),
       "SEGA", 80, bold=True, color=RGBColor(0x0D, 0x11, 0x26))
    tb(s, Inches(0.42), Inches(3.85), Inches(12), Inches(0.5),
       "Shooter Market Intelligence  \u00b7  Confidential", 13, color=MUTED, italic=True)

    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()

# ─────────────────────────────────────────────────────────────
# PPTX SNAPSHOT EXPORT  (Playwright headless Chromium)
# Renders a dark-themed HTML replica of the dashboard — the AI
# report with coloured tables, KPI tiles, the WoW table, and the
# top-25 CCU chart — then screenshots each section and assembles
# the images into full-bleed 16:9 slides. Tall sections are
# sliced into multiple slides so nothing is cut off.
# ─────────────────────────────────────────────────────────────

_PW_INSTALL_TRIED = False

def _ensure_playwright_chromium() -> bool:
    """True if Playwright + Chromium are usable; installs Chromium once if missing."""
    global _PW_INSTALL_TRIED
    try:
        from playwright.sync_api import sync_playwright
    except ImportError:
        return False
    try:
        with sync_playwright() as pw:
            b = pw.chromium.launch(headless=True, args=["--no-sandbox"])
            b.close()
        return True
    except Exception as e:
        msg = str(e)
        if ("Executable doesn't exist" in msg or "playwright install" in msg) and not _PW_INSTALL_TRIED:
            _PW_INSTALL_TRIED = True
            import subprocess
            try:
                r = subprocess.run(
                    [sys.executable, "-m", "playwright", "install", "chromium"],
                    capture_output=True, text=True, timeout=300,
                )
                return r.returncode == 0
            except Exception:
                return False
        return False


def _fig_to_b64_png(fig, width=1760, height=520) -> str | None:
    """Plotly figure → base64 PNG (requires kaleido). Returns None on failure."""
    try:
        png = fig.to_image(format="png", width=width, height=height, scale=2)
        return base64.b64encode(png).decode()
    except Exception:
        return None


def _snapshot_md_section_to_html(md: str) -> str:
    """Markdown section body → HTML, with markdown tables rendered through
    html_table() so they keep the dashboard's conditional colouring."""
    lines = md.splitlines(keepends=True)
    out, i = [], 0
    while i < len(lines):
        line = lines[i]
        if "|" in line and line.strip().startswith("|"):
            table_lines = []
            while i < len(lines) and "|" in lines[i]:
                table_lines.append(lines[i]); i += 1
            parsed = parse_md_table("".join(table_lines))
            if parsed:
                headers, rows = parsed
                out.append(html_table(rows, headers))
            else:
                out.append("<pre>" + "".join(table_lines) + "</pre>")
        else:
            text = ""
            while i < len(lines) and not ("|" in lines[i] and lines[i].strip().startswith("|")):
                text += lines[i]; i += 1
            if text.strip():
                if MARKDOWN_AVAILABLE:
                    out.append(_md_lib.markdown(text, extensions=["tables", "fenced_code"]))
                else:
                    out.append("<p>" + text.replace("\n", "<br>") + "</p>")
    return "".join(out)


def build_snapshot_html(report_md: str, ccu_data: list[dict], label: str,
                        wow_diff: dict, language: str = "English") -> str:
    """Build a self-contained dark HTML replica of the dashboard.
    Each <section class='snap'> becomes one (or more) slides."""
    jp = language == "Japanese"
    today_str = datetime.utcnow().strftime("%Y-%m-%d")

    # ── KPI numbers ──
    total_ccu = sum(r["ccu"] for r in ccu_data)
    wow_up    = sum(1 for v in wow_diff.values() if v["delta"] > 0)
    wow_down  = sum(1 for v in wow_diff.values() if v["delta"] < 0)
    growing   = sum(1 for r in ccu_data if r.get("yoy_val", 0) > 0)
    declining = sum(1 for r in ccu_data if r.get("yoy_val", 0) < 0)
    health_ratios = [r["ccu"] / r["hist_summary"]["peak_ever"] * 100
                     for r in ccu_data if r.get("hist_summary", {}).get("peak_ever") and r["ccu"] > 0]
    avg_health = sum(health_ratios) / len(health_ratios) if health_ratios else 0

    def kpi(label_, value, sub):
        return (f"<div class='kpi'><div class='kpi-l'>{label_}</div>"
                f"<div class='kpi-v'>{value}</div><div class='kpi-s'>{sub}</div></div>")

    L = {
        "total":  ("合計CCU（追跡対象）", "Total CCU (Tracked)")[0 if jp else 1],
        "tot_s":  (f"シューター {len(ccu_data)} タイトル合計", f"Across {len(ccu_data)} shooter titles")[0 if jp else 1],
        "health": ("平均ピーク健全度", "Avg Peak Health")[0 if jp else 1],
        "hea_s":  ("現在CCU ÷ 過去最高ピーク", "current CCU vs. all-time peak")[0 if jp else 1],
        "wow":    ("WoW成長バランス", "WoW Growth Balance")[0 if jp else 1],
        "wow_s":  (f"成長 {wow_up} / 下落 {wow_down}", f"{wow_up} growing / {wow_down} declining")[0 if jp else 1],
        "yoy":    ("YoY成長バランス", "YoY Growth Balance")[0 if jp else 1],
        "yoy_s":  (f"成長 {growing} / 下落 {declining}", f"{growing} growing / {declining} declining")[0 if jp else 1],
        "wow_h":  (f"WoW CCU変化 — 成長 {wow_up} / 下落 {wow_down}",
                   f"WoW CCU Change — {wow_up} growing, {wow_down} declining")[0 if jp else 1],
        "rank_h": ("シューターCCUランキング", "Top Shooter CCU Stack-Ranked")[0 if jp else 1],
        "chart_h":("ライブCCU上位タイトル", "Top Titles by Live CCU")[0 if jp else 1],
        "sub":    ("シューター市場インテリジェンス", "Shooter Market Intelligence")[0 if jp else 1],
        "conf":   ("社内分析専用", "Internal analytics use only")[0 if jp else 1],
    }
    COLS = {
        "title":  ("タイトル", "Title")[0 if jp else 1],
        "year":   ("発売年", "Year")[0 if jp else 1],
        "live":   ("ライブCCU", "Live CCU")[0 if jp else 1],
        "ago":    ("7日前", "7 Days Ago")[0 if jp else 1],
        "dccu":   ("変化 (CCU)", "Change (CCU)")[0 if jp else 1],
        "wpct":   ("週次変化", "Weekly Change")[0 if jp else 1],
        "rank":   "#",
        "sub":    ("サブジャンル", "Sub-Genre")[0 if jp else 1],
        "pub":    ("パブリッシャー", "Publisher")[0 if jp else 1],
        "yoy":    "YoY",
        "mom":    "MoM",
        "rev":    ("レビュー", "Review")[0 if jp else 1],
    }

    sections_html: list[str] = []

    # ── Section 1: cover + KPI tiles ──
    sections_html.append(f"""
<section class="snap">
  <div class="topbar"><span class="logo"><span class="seg">SEGA</span> SHOOTER INTELLIGENCE</span>
    <span class="tlabel">{L['sub']}</span><span class="pill">{today_str}</span></div>
  <div class="hero-t">{label}</div>
  <div class="kpis">
    {kpi(L['total'], f"{total_ccu:,}", L['tot_s'])}
    {kpi(L['health'], f"{avg_health:.0f}%", L['hea_s'])}
    {kpi(L['wow'], f"{wow_up} / {wow_down}", L['wow_s'])}
    {kpi(L['yoy'], f"{growing} / {declining}", L['yoy_s'])}
  </div>
</section>""")

    # ── Sections 2..n: AI report split by ## headings ──
    cur_title, cur_body, report_sections = "", [], []
    for line in report_md.splitlines():
        if line.startswith("## "):
            if cur_title or cur_body:
                report_sections.append((cur_title, "\n".join(cur_body).strip()))
            cur_title = line.lstrip("# ").strip()
            cur_body  = []
        elif line.startswith("# "):
            continue
        else:
            cur_body.append(line)
    if cur_title or any(b.strip() for b in cur_body):
        report_sections.append((cur_title, "\n".join(cur_body).strip()))

    for title, body in report_sections:
        if not body.strip() and not title:
            continue
        body_html = _snapshot_md_section_to_html(body)
        head_html = f"<div class='sec-h'><span class='dot'></span>{title.upper()}</div>" if title else ""
        sections_html.append(f"<section class='snap'>{head_html}{body_html}</section>")

    # ── WoW table section ──
    wow_rows = []
    for r in ccu_data:
        d = wow_diff.get(r["app_id"])
        if d:
            wow_rows.append({
                COLS["title"]: r["name"],
                COLS["year"]:  str(r.get("year", "—")),
                COLS["live"]:  f"{d['curr_ccu']:,}",
                COLS["ago"]:   f"{d['prev_ccu']:,}",
                COLS["dccu"]:  (f"+{d['delta']:,}" if d["delta"] > 0 else f"{d['delta']:,}"),
                COLS["wpct"]:  (f"+{round(d['delta_pct'])}%" if d["delta_pct"] > 0 else f"{round(d['delta_pct'])}%"),
                "_sort": d["delta_pct"],
            })
    if wow_rows:
        wow_rows.sort(key=lambda x: x["_sort"], reverse=True)
        for r2 in wow_rows: r2.pop("_sort", None)
        sections_html.append(
            f"<section class='snap'><div class='sec-h'><span class='dot'></span>{L['wow_h']}</div>"
            + html_table(wow_rows, [COLS['title'], COLS['year'], COLS['live'],
                                    COLS['ago'], COLS['dccu'], COLS['wpct']])
            + "</section>")

    # ── Top-25 chart section (kaleido) ──
    top_n = sorted(ccu_data, key=lambda r: r["ccu"], reverse=True)[:25]
    def _bar_col(r):
        d = wow_diff.get(r["app_id"])
        if d and d["delta"] > 0: return "#20c65a"
        if d and d["delta"] < 0: return "#ff4d4d"
        return "#888aaa"
    fig = go.Figure(go.Bar(
        x=[f"#{i+1} {r['name'][:18]}" for i, r in enumerate(top_n)],
        y=[r["ccu"] for r in top_n],
        marker_color=[_bar_col(r) for r in top_n],
        text=[f"{r['ccu']:,}" for r in top_n],
        textposition="outside",
        textfont=dict(size=10, color="#b8bcd4"),
    ))
    fig.update_layout(
        paper_bgcolor="#0f1120", plot_bgcolor="#0f1120",
        font=dict(family="sans-serif", color="#eef0fa"),
        margin=dict(l=20, r=20, t=20, b=110),
        xaxis=dict(showgrid=False, tickfont=dict(size=10), tickangle=-35, linecolor="#232640"),
        yaxis=dict(showgrid=True, gridcolor="#1a1e30", tickformat=","),
        showlegend=False,
    )
    chart_b64 = _fig_to_b64_png(fig)
    if chart_b64:
        sections_html.append(
            f"<section class='snap'><div class='sec-h'><span class='dot'></span>{L['chart_h']}</div>"
            f"<img src='data:image/png;base64,{chart_b64}' style='width:100%;border:1px solid #232640;border-radius:8px;'/>"
            "</section>")

    # ── Stack-ranked table section ──
    rank_rows = [{
        COLS["rank"]:  i + 1,
        COLS["title"]: r["name"],
        COLS["sub"]:   r["sub"],
        COLS["pub"]:   r["publisher"],
        COLS["live"]:  f"{r['ccu']:,}",
        COLS["yoy"]:   r.get("yoy", "N/A"),
        COLS["mom"]:   r.get("hist_summary", {}).get("mom_trend", "—"),
        COLS["rev"]:   f"{r['review_pct']}%" if r.get("review_pct") else "—",
    } for i, r in enumerate(ccu_data)]
    sections_html.append(
        f"<section class='snap'><div class='sec-h'><span class='dot'></span>{L['rank_h']}</div>"
        + html_table(rank_rows, [COLS['rank'], COLS['title'], COLS['sub'], COLS['pub'],
                                 COLS['live'], COLS['yoy'], COLS['mom'], COLS['rev']])
        + f"<div class='foot'>SEGA Shooter Intelligence · {L['conf']} · {today_str}</div></section>")

    return f"""<!DOCTYPE html><html><head><meta charset="utf-8">
<style>
@import url('https://fonts.googleapis.com/css2?family=Inter+Tight:wght@700;800;900&family=Poppins:wght@300;400;500;600&family=Noto+Sans+JP:wght@400;500;700&display=swap');
* {{ margin:0; padding:0; box-sizing:border-box;
    font-family:'Poppins','Noto Sans JP',sans-serif; }}
body {{ background:#0a0c1a; color:#eef0fa; width:1920px; }}
section.snap {{ background:#0a0c1a; padding:48px 64px; border-bottom:1px solid #0a0c1a; }}
.topbar {{ display:flex; align-items:center; gap:24px; border-bottom:1px solid #232640;
          padding-bottom:18px; margin-bottom:32px; }}
.logo {{ font-family:'Inter Tight',sans-serif; font-weight:900; font-size:22px;
        letter-spacing:.12em; }}
.logo .seg {{ color:#4080ff; }}
.tlabel {{ font-size:13px; color:#5a5f82; letter-spacing:.2em; text-transform:uppercase; }}
.pill {{ margin-left:auto; background:rgba(64,128,255,.16); border:1px solid rgba(64,128,255,.28);
        border-radius:20px; padding:4px 16px; font-size:12px; font-weight:700;
        letter-spacing:.14em; color:#4080ff; }}
.hero-t {{ font-family:'Inter Tight',sans-serif; font-size:44px; font-weight:900;
          letter-spacing:-.02em; margin-bottom:36px; }}
.kpis {{ display:flex; gap:20px; }}
.kpi {{ flex:1; background:#0f1120; border:1px solid #232640; border-top:2px solid #4080ff;
       border-radius:8px; padding:26px 30px; }}
.kpi-l {{ font-size:12px; font-weight:700; letter-spacing:.22em; text-transform:uppercase;
         color:#5a5f82; margin-bottom:12px; }}
.kpi-v {{ font-family:'Inter Tight',sans-serif; font-size:44px; font-weight:900; line-height:1; }}
.kpi-s {{ font-size:13px; color:#5a5f82; margin-top:8px; }}
.sec-h {{ font-family:'Inter Tight','Noto Sans JP',sans-serif; font-size:16px; font-weight:800;
         letter-spacing:.2em; text-transform:uppercase; color:#b8bcd4;
         border-bottom:1px solid #232640; padding-bottom:12px; margin-bottom:24px;
         display:flex; align-items:center; gap:12px; }}
.sec-h .dot {{ width:8px; height:8px; background:#4080ff; border-radius:2px;
              box-shadow:0 0 8px #4080ff; display:inline-block; }}
h1,h2,h3 {{ font-family:'Inter Tight','Noto Sans JP',sans-serif; color:#fff; margin:18px 0 10px; }}
h2 {{ color:#7ab0ff; }}
p, li {{ font-size:17px; line-height:1.75; color:#d8dcf0; margin-bottom:10px; }}
ul, ol {{ margin:0 0 14px 28px; }}
strong {{ color:#fff; }}
table {{ width:100%; border-collapse:collapse; font-size:15px; }}
.foot {{ margin-top:28px; padding-top:14px; border-top:1px solid #232640;
        font-size:12px; color:#5a5f82; letter-spacing:.1em; }}
</style></head><body>
{''.join(sections_html)}
</body></html>"""


def generate_pptx_snapshot_bytes(report_md: str, ccu_data: list[dict], label: str,
                                 wow_diff: dict, language: str = "English") -> bytes | None:
    """Snapshot PPTX: render the dashboard replica in headless Chromium,
    screenshot each section, slice tall ones, place as full-bleed slides.
    Returns None if playwright/pillow/python-pptx are unavailable."""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Emu
        from pptx.dml.color import RGBColor
        from PIL import Image
        from io import BytesIO
        from playwright.sync_api import sync_playwright
    except ImportError:
        return None

    if not _ensure_playwright_chromium():
        return None

    html = build_snapshot_html(report_md, ccu_data, label, wow_diff, language)

    # ── Render & capture each section ──
    shots: list[bytes] = []
    try:
        with sync_playwright() as pw:
            browser = pw.chromium.launch(
                headless=True,
                args=["--no-sandbox", "--disable-dev-shm-usage", "--disable-gpu",
                      "--force-device-scale-factor=1"],
            )
            page = browser.new_page(viewport={"width": 1920, "height": 1080})
            try:
                page.set_content(html, wait_until="networkidle", timeout=25000)
            except Exception:
                # Fonts CDN blocked / slow — fall back to basic load
                page.set_content(html, wait_until="load", timeout=25000)
            page.wait_for_timeout(900)  # let webfonts settle
            for el in page.query_selector_all("section.snap"):
                try:
                    shots.append(el.screenshot(type="png"))
                except Exception:
                    continue
            browser.close()
    except Exception:
        return None

    if not shots:
        return None

    # ── Assemble slides ──
    SLIDE_W_IN, SLIDE_H_IN = 13.333, 7.5
    prs = Presentation()
    prs.slide_width  = Inches(SLIDE_W_IN)
    prs.slide_height = Inches(SLIDE_H_IN)
    blank = prs.slide_layouts[6]
    BG = RGBColor(0x0A, 0x0C, 0x1A)

    def _add_image_slide(img: "Image.Image"):
        slide = prs.slides.add_slide(blank)
        f = slide.background.fill
        f.solid(); f.fore_color.rgb = BG
        buf = BytesIO(); img.save(buf, format="PNG"); buf.seek(0)
        disp_h_in = SLIDE_W_IN * img.height / img.width
        slide.shapes.add_picture(buf, 0, 0,
                                 width=Inches(SLIDE_W_IN), height=Inches(disp_h_in))

    for shot in shots:
        img = Image.open(BytesIO(shot)).convert("RGB")
        # px height of one 16:9 slide at this image width
        seg_h = int(img.width * SLIDE_H_IN / SLIDE_W_IN)
        if img.height <= seg_h + 40:        # fits on one slide (small tolerance)
            _add_image_slide(img)
        else:                                # slice tall sections
            top = 0
            while top < img.height:
                bottom = min(top + seg_h, img.height)
                if img.height - top < seg_h * 0.18 and prs.slides:
                    break                    # skip a sliver-thin trailing strip
                _add_image_slide(img.crop((0, top, img.width, bottom)))
                top = bottom

    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()

# ─────────────────────────────────────────────────────────────
# SESSION STATE
# ─────────────────────────────────────────────────────────────

def init_session_defaults() -> None:
    """Populate st.session_state with defaults for any key not already set.
    Idempotent — safe to call at the top of every page."""
    defaults = {
        "claude_key": st.secrets.get("CLAUDE_KEY", os.environ.get("CLAUDE_KEY", "")),
        "ccu_data": [],
        "active_query": None,
        "ai_report": "",
        "ai_chat_history": [],
        "ai_chat_pending": False,
        "report_label": "",
        "custom_query": "",
        "report_language": "English",
        "report_cache": {},
        "uploaded_csvs": {},   # app_id -> bytes, from sidebar uploader
        "roster_genre":  "BOTH",            # "FPS", "TPS", or "BOTH"
        "roster_filter": [],                # list of app_ids to include (empty = all)
        "drilldown_game":   None,           # app_id of selected game
        "drilldown_report": "",             # cached drilldown report
        "drilldown_cache":  {},             # {app_id: report_text}
    }
    for k, v in defaults.items():
        if k not in st.session_state:
            st.session_state[k] = v

# ─────────────────────────────────────────────────────────────
# TRANSLATIONS  (EN / JP)
# Every user-facing string in the app routes through T() so the
# whole page — not just the AI report — follows the EN/JP toggle.
# ─────────────────────────────────────────────────────────────

TRANSLATIONS = {
    "English": {
        # Topbar / nav
        "topbar_subtitle":        "Market &amp; Tech Analysis",
        # Sidebar
        "sidebar_config":         "Configuration",
        "admin_header":           "ADMIN & SETTINGS",
        "admin_scheduler_header": "BACKGROUND SCHEDULER",
        "admin_scheduler_desc":   "Auto-archives a weekly report every Monday at 09:00 UTC, for all three genres (FPS, TPS, BOTH). Requires ANTHROPIC_API_KEY as an environment variable (Bedrock secrets are not visible to the background thread).",
        "admin_log_header":       "Scheduler Log (last 20 entries)",
        "admin_log_none":         "No scheduler log yet — entries appear here after the first Monday run, or after any error.",
        "admin_log_clear":       "Clear log",
        "fetch_health_warning":  "⚠️ Live CCU fetch failed for all {n} titles this run — every value below is 0 because the Steam API call didn't succeed, not because these games genuinely have no players. Check the connectivity diagnostics on the Admin page.",
        "fetch_health_partial":  "ℹ️ Live CCU succeeded for only {pct}% of titles this run ({live}/{total}); the rest fell back to cached/CSV data or show as 0. This can happen during a transient Steam/SteamSpy rate limit — try Refresh CCU Data again in a minute.",
        "conn_check_header":     "CONNECTIVITY CHECK",
        "conn_check_desc":       "Tests each upstream API directly with a known title (Counter-Strike 2, app 730) and shows the raw result — use this to tell whether a fetch failure is a network/firewall issue, a rate limit, or something else.",
        "conn_check_btn":        "Run Connectivity Check",
        "conn_check_running":    "Testing Steam, SteamSpy, and Twitch connectivity…",
        "pipeline_check_header": "Pipeline Check (advanced)",
        "pipeline_check_desc":   "Calls the exact function the real fetch uses — once directly, once inside a worker thread like the Dashboard does — so you can tell whether a fetch failure is a bug in the fetch logic itself (fails both ways) or something specific to running inside a thread (succeeds direct, fails threaded). Run this if the Connectivity Check above passes but the Dashboard still shows all-zero CCU.",
        "pipeline_check_btn":    "Run Pipeline Check",
        "pipeline_check_running": "Testing the fetch pipeline directly and inside a worker thread…",
        "pipeline_check_thread_specific": "⚠️ The main-thread call succeeded but the worker-thread call failed — this points to something specific to running inside a background thread (e.g. a caching/session quirk), not the fetch logic itself. Worth trying a lower max_workers or removing st.cache_data from the affected function as a workaround.",
        "pipeline_check_both_fail": "⚠️ Both calls failed — this is a bug in the fetch logic itself, not a threading issue. Check the error details above; if you can share them, that pinpoints the exact line to fix.",
        "pipeline_check_both_ok": "✅ Both succeeded — the fetch pipeline itself is working correctly for this title right now. If the Dashboard still shows all-zero CCU, try Refresh CCU Data again; it may have been a transient issue.",
        "pipeline_check_both_no_data": "🟡 Both calls completed without error, but neither got a live value back from Steam — the fetch logic itself is fine, but Steam's CCU endpoint isn't returning data for this title right now (could be a transient API issue, or this title genuinely has an API quirk). Try a different title or wait a few minutes.",
        "pipeline_check_main_only_exception": "The main-thread call raised an exception but the worker-thread call didn't — this is unusual and may just mean the underlying failure was transient between the two calls. Try running the check again.",
        "admin_csv_note":        "Uploaded here, used everywhere: CSVs uploaded on this page are available immediately to the Dashboard's next fetch — no need to re-upload per page.",
        "api_loaded":             "Anthropic API key loaded",
        "api_missing":            "Anthropic API key missing",
        "bedrock_ok":             "✓ AWS Bedrock credentials loaded",
        "bedrock_missing":        "AWS Bedrock credentials missing",
        "model_caption":          "Model: claude-sonnet-4-20250514",
        "ccu_caption":            "CCU: Steam public API (5 min cache)",
        "engagement_caption":     "Engagement: SteamSpy API (1 hr cache)",
        "cache_age":              "📦 Data cached: {age}",
        "cache_refetch":          "Re-fetches automatically every 24 hours.",
        "cache_none":             "No cache yet — fetching on load.",
        "refresh_now":            "🔄 Refresh Now",
        "refresh_now_help":       "Force a fresh CCU fetch, ignoring the 24-hour cache",
        "csvs_loaded":            "SteamDB CSVs: {n}/{total} loaded",
        "csv_missing":            "Missing: {names}",
        "csv_drop_hint":          "Drop steamdb_chart_{appid}.csv into /data to update",
        "upload_header":          "**Upload SteamDB CSVs**",
        "upload_caption":         "Upload steamdb_chart_{appid}.csv files directly — no repo access needed.",
        "upload_loaded":          "{n} CSV(s) loaded: {names}",
        "signed_in_as":           "Signed in as",
        "sign_out":               "Sign out",
        "watchlist_header":       "My Watchlist",
        "watchlist_max":          "Max 5 pinned titles",
        "lang_header":            "Report Language",
        "last_fetched":           "CCU last fetched: {time}",
        # Hero
        "hero_line1":             "SHOOTER MARKET",
        "hero_line2":             "INTELLIGENCE",
        "hero_sub":               "Live Steam CCU data · AI-powered analysis · Competitive benchmarks · Weekly reporting templates — all in one tool for SEGA's publishing and strategy teams.",
        # Section headers
        "select_analysis":        "SELECT ANALYSIS TYPE",
        "live_ccu_header":        "LIVE STEAM CCU SNAPSHOT",
        "ai_analysis_header":     "AI ANALYSIS — {label}",
        "dataset_header":         "DATASET",
        # Dataset picker
        "btn_fps":                "First-Person",
        "btn_tps":                "Third-Person",
        "btn_both":               "FPS + TPS",
        "games_included":         "Games included: {n} / {total}",
        "select_all":             "Select all",
        "clear_all":              "Clear all",
        "overlap_note":           "ℹ {n} titles also appear in the other list: {names}",
        # Fetch button
        "fetch_ccu_btn":          "Fetch Live CCU Data",
        "fetch_spinner":          "Pulling live CCU from Steam + SteamSpy…",
        "fetching_game":          "Fetching: {name}…",
        "fetch_done":             "Fetched {n} titles",
        "refresh_ccu_btn":        "Refresh CCU Data",
        "cache_loaded_toast":     "Loaded from cache ({age})",
        # KPI cards
        "kpi_total_ccu":          "Total CCU (Tracked)",
        "kpi_total_sub":          "Across {n} shooter titles",
        "kpi_wow":                "WoW Growth Balance",
        "kpi_wow_sub":            "of {n} titles with CSV data",
        "kpi_wow_none":           "No CSV data loaded",
        "kpi_yoy":                "YoY Growth Balance",
        "kpi_yoy_sub":            "of {n} titles with YoY data",
        "kpi_mom":                "MoM Growth Balance",
        "kpi_mom_sub":            "of {n} titles with CSV data",
        "kpi_mom_none":           "No CSV data loaded",
        "kpi_csvs":               "SteamDB CSVs Loaded",
        "kpi_csvs_sub":           "titles with full historical data",
        "kpi_health":             "Avg Peak Health",
        "kpi_health_sub":         "current CCU vs. all-time peak",
        "kpi_twitch":             "Total Twitch Viewers",
        "kpi_twitch_sub":         "across {n} tracked titles (top-100 streams)",
        "kpi_twitch_none":        "Add TWITCH_CLIENT_ID to secrets to enable",
        "kpi_best_grower":        "Biggest YoY Grower",
        "kpi_best_sub":           "{pct} YoY (SteamDB)",
        "kpi_worst_decline":      "Biggest YoY Decline",
        "kpi_worst_sub":          "{pct} YoY (SteamDB)",
        # Expanders
        "yoy_expander":           "YoY CCU Change — {up} growing, {down} declining",
        "wow_expander":           "WoW CCU Change — {up} growing, {down} declining",
        "mom_expander":           "MoM CCU Change — {up} growing, {down} declining",
        "growth_label":           "Growth",
        "decline_label":          "Decline",
        "wow_caption":            "Comparing latest CSV value vs. the row closest to exactly 7 days prior. Source: SteamDB 10-minute interval data.",
        "wow_none":               "No CSV data loaded yet. Add steamdb_chart_{appid}.csv files to the /data folder.",
        "heatmap_expander":       "Sub-Genre CCU Heat Map",
        "heatmap_caption":        "Source: Aggregated from Steam API live CCU, grouped by sub-genre tag in roster.",
        "table_expander":         "Top Shooter CCU Stack-Ranked — {genre}",
        "table_footnote":         "Review = all-time positive ÷ total reviews.  Rev/day = reviews/day over 7 days (N/A until 7+ days of snapshots).  Twitch = live viewers across top-100 streams (requires TWITCH_CLIENT_ID).  — = no data.  * = latest CSV value used instead of live API.",
        "history_expander":       "Monthly Peak CCU History — SteamDB Data",
        "history_caption":        "Source: SteamDB 10-min interval CSVs, aggregated to monthly peak{note}",
        "history_note_few":       " | Dashed lines = key events (hover for details)",
        "history_note_many":      " | Hover each line for key events",
        "chart_history_title":    "Monthly Peak CCU — Last 24 Months (SteamDB)",
        "formulas_caption":       "**WoW%** = (Live CCU − CCU 7 days ago) ÷ CCU 7 days ago × 100  |  **MoM%** = (Last month avg − Prior month avg) ÷ Prior month avg × 100  |  **YoY%** = (Current month avg − Same month last year avg) ÷ Same month last year avg × 100  |  Source: SteamDB CSV data · WoW falls back to latest CSV row if no 7-day snapshot exists",
        # Table columns
        "col_rank":               "#",
        "col_title":              "Title",
        "col_subgenre":           "Sub-Genre",
        "col_publisher":          "Publisher",
        "col_f2p":                "F2P",
        "col_live_ccu":           "Live CCU",
        "col_yoy":                "YoY",
        "col_year":               "Year",
        "col_data_source":        "Data Source",
        "col_peak_ever":          "All-Time Peak",
        "col_peak_12m":           "12m Peak",
        "col_avg_ccu_12m":        "12m Avg",
        "col_mom":                "MoM",
        "col_review":             "Review",
        "col_rev_velocity":       "Rev/day",
        "col_twitch":             "Twitch",
        "col_yoy_source":         "YoY Source",
        "col_owners":             "Est. Owners",
        "col_7d_ago":             "7 Days Ago",
        "col_change_ccu":         "Change (CCU)",
        "col_weekly_change":      "Weekly Change",
        "col_steam_page":         "Steam Page",
        "col_1yr_ago":            "1 Year Ago",
        "col_annual_change":      "Annual Change",
        "col_month_change":       "Month Change",
        "col_delta_ccu":          "Δ CCU",
        "col_delta_pct":          "Δ %",
        "col_direction":          "Direction",
        "col_reference":          "Reference",
        "store_link":             "↗ Store",
        # Bar charts
        "chart_top25_title":      "Top {n} {genre} Titles by Live CCU",
        "chart_ranks_title":      "Ranks {start}-{end} {genre} Titles by Live CCU",
        "chart_caption":          "Green = WoW up  |  Red = WoW down  |  Grey = no CSV data  |  {genre} roster",
        # Analysis presets
        "run_analysis":           "Run Analysis",
        "custom_label":           "Or ask a custom question",
        "custom_placeholder":     "e.g. Compare monetisation models across the top 5 F2P shooters on Steam…",
        "preset_labels": {
            "ccu_mecha":        "CCU Trends & Mecha-Shooter Demand",
            "competitive_gap":  "Competitive Gap Analysis — Top 5",
            "table_stakes":     "2026 Netcode & Server Table Stakes",
            "social_metrics":   "Social Media Metrics for Day-1 Success",
            "weekly_report":    "Weekly Retention & Engagement Report Template",
        },
        "preset_descs": {
            "ccu_mecha":        "Analyze the top 10 shooters on Steam and compare CCU trends to last year. What does this say about current demand for mecha-shooters?",
            "competitive_gap":  "Identify the exact mechanics, monetisation patterns, and content cadences separating the top 5 by CCU. Surfaces addressable product gaps SEGA could exploit.",
            "table_stakes":     "What are the non-negotiable 'table stakes' for a competitive shooter in 2026 regarding netcode and server architecture to satisfy Western competitive integrity standards?",
            "social_metrics":   "Based on recent investor reports and market data, what are the primary social media metrics to track to predict a new shooter's Day 1 success?",
            "weekly_report":    "Create a template for a weekly market report that tracks retention and engagement KPIs across the top 100 shooters, highlighting any 'breakout' indie titles.",
        },
        "preset_tags": {
            "ccu_mecha":        "Market",
            "competitive_gap":  "Strategy",
            "table_stakes":     "Tech",
            "social_metrics":   "Social",
            "weekly_report":    "Report",
        },
        "run_btn":                "Run",
        "custom_query_label":     "Custom Query",
        # AI report
        "cache_notice":           "Loaded from cache — data unchanged since last run. Re-fetch CCU to force refresh.",
        "no_ccu_warning":         "Please fetch live CCU data first.",
        "spinner_generating":     "Claude is generating your analysis…",
        "no_key_warning":         "AWS Bedrock credentials not found. Add AWS_ACCESS_KEY_ID_API to .streamlit/secrets.toml to run AI analysis.",
        "no_anthropic_error":     "Install the `anthropic` package: `pip install anthropic`",
        "auth_error":             "Invalid API key. Check the credentials in .streamlit/secrets.toml.",
        "rate_limit_error":       "Rate limit hit. Wait a moment and try again.",
        "analysis_failed":        "Analysis failed: {e}",
        "auto_archived":          "📋 Auto-archived: {f}",
        # Downloads
        "download_report_header": "DOWNLOAD REPORT",
        "dl_md":                  "Download Markdown",
        "dl_html":                "Download HTML",
        "dl_pdf":                 "Download PDF",
        "dl_pptx_btn":            "Download PowerPoint",
        "dl_pptx_file":           "Download .pptx",
        "dl_pptx_error":          "PPTX generation failed. Ensure python-pptx is installed: pip install python-pptx",
        "dl_pdf_missing":         "PDF: install `reportlab`",
        "spinner_pptx":           "Building slides… (taking a snapshot of the page)",
        "pptx_fallback_note":     "Snapshot export unavailable (pip install playwright pillow kaleido, then `playwright install chromium`) — generated the text-based deck instead.",
        "regen_btn":              "Regenerate",
        "archive_btn":            "🗄 Archive",
        "archive_saved":          "Saved: {f}",
        "archive_failed":         "Archive failed",
        # Follow-up chat
        "chat_header":            "FOLLOW-UP CHAT",
        "chat_subtext":           "— ask Claude follow-up questions about this report",
        "chat_placeholder":       "Ask a follow-up question about this report…",
        "chat_clear":             "Clear chat history",
        "chat_error":             "Chat error: {e}",
        # Empty state
        "empty_title":            "NO ANALYSIS SELECTED",
        "empty_sub":              "Fetch live CCU data above, then choose an analysis type or enter a custom question to generate your report.",
        # Drilldown
        "drilldown_header":       "GAME INTELLIGENCE DEEP DIVE",
        "back_btn":               "← Back to Dashboard",
        "drilldown_select":       "Select a title to deep dive...",
        "drilldown_btn":          "Deep Dive",
        "drilldown_no_data":      "Game data not found — please fetch CCU data first.",
        "drilldown_no_key":       "AWS Bedrock credentials not found. Add them to .streamlit/secrets.toml to run AI analysis.",
        "drilldown_spinner":      "Generating deep-dive analysis…",
        "drilldown_failed":       "Deep dive failed: {e}",
        "drilldown_dl":           "Download Deep Dive (.md)",
        "drilldown_no_csv":       "No SteamDB CSV for {name}. Drop steamdb_chart_{appid}.csv into /data to enable the history chart.",
        "chart_drilldown_title":  "{name} — Monthly Peak CCU",
        "pill_peak_ever":         "All-time peak",
        "pill_peak_12m":          "12m peak",
        "pill_avg_12m":           "12m avg CCU",
        "pill_mom":               "MoM trend",
        "trace_peak":             "Peak CCU",
        "trace_avg":              "Avg CCU",
        "trace_live":             "Live CCU",
        "no_hist_info":           "No historical CSV data for this title. Drop steamdb_chart_{appid}.csv into /data.",
        "yoy_caption":            "Snapshot = real 12-month rolling average (active after ~13 months) · SteamDB CSV = genuine same-month YoY · SteamSpy proxy = engagement estimate",
        "yoy_none":               "No YoY data available — fetch CCU data first.",
        # Monthly analysis
        "monthly_header":         "MONTHLY ANALYSIS",
        "monthly_none":           "No archived reports yet. Reports are auto-archived every Monday, or use the Archive button above to save the current report.",
        "monthly_count":          "{n} archived reports across {m} month(s) · stored in /data/report_archive/",
        "monthly_expander":       "📅 Monthly Comparison — compare weekly reports & run accuracy check",
        "monthly_select":         "Select month to analyse",
        "monthly_reports_for":    "**{n} report(s)** archived for {month}:",
        "monthly_titles":         "{n} titles",
        "monthly_drift":          "**Month-long CCU drift** ({a} → {b}):",
        "monthly_one":            "Only 1 report this month — need at least 2 to show drift. More will accumulate each Monday.",
        "monthly_ai_header":      "**AI Accuracy & Trend Check**",
        "monthly_ai_caption":     "Claude compares the weekly archived reports for this month and flags any divergence from the month-long drift.",
        "monthly_no_key":         "Add your AWS Bedrock credentials to run the monthly analysis.",
        "monthly_run_btn":        "Run Monthly Analysis",
        "monthly_spinner":        "Running monthly analysis…",
        "monthly_failed":         "Monthly analysis failed: {e}",
        "monthly_dl":             "⬇ Download Monthly Report (.md)",
        "col_monthly_drift":      "Monthly Drift",
        "col_latest_wow":         "Latest WoW",
        # Footer
        "footer_brand":           "SEGA SHOOTER INTELLIGENCE",
        "footer_note":            "Data sourced from Steam public API · Powered by Claude · Internal analytics use only",
        # Watchlist section
        "watchlist_section":      "WATCHLISTED TITLES",
        # Values
        "yes": "Yes",
        "no":  "No",
        "up":  "Up",
        "down":"Down",
        "flat":"Flat",
    },
    "Japanese": {
        # Topbar / nav
        "topbar_subtitle":        "市場・技術分析",
        # Sidebar
        "sidebar_config":         "設定",
        "admin_header":           "管理・設定",
        "admin_scheduler_header": "バックグラウンドスケジューラー",
        "admin_scheduler_desc":   "毎週月曜09:00 UTCに、FPS・TPS・BOTHの3ジャンルすべてで週次レポートを自動アーカイブします。バックグラウンドスレッドはBedrockのsecretsを参照できないため、環境変数 ANTHROPIC_API_KEY が必要です。",
        "admin_log_header":       "スケジューラーログ（直近20件）",
        "admin_log_none":        "スケジューラーログはまだありません — 最初の月曜実行後、またはエラー発生時にここに表示されます。",
        "admin_log_clear":       "ログをクリア",
        "fetch_health_warning":  "⚠️ 今回の取得で {n} タイトル全てのライブCCU取得が失敗しました — 以下の値が0なのはSteam APIの呼び出しが成功しなかったためであり、実際にプレイヤーが0人というわけではありません。Adminページの接続診断を確認してください。",
        "fetch_health_partial":  "ℹ️ 今回はライブCCUが {live}/{total} タイトル（{pct}%）でのみ成功しました。残りはキャッシュ/CSVデータにフォールバックするか0として表示されています。Steam/SteamSpyの一時的なレート制限が原因の場合があります — 1分ほど待ってから「CCUデータを更新」を再試行してください。",
        "conn_check_header":     "接続診断",
        "conn_check_desc":       "既知のタイトル（Counter-Strike 2, app 730）を使って各上流APIを直接テストし、生の結果を表示します — 取得失敗がネットワーク/ファイアウォールの問題か、レート制限か、その他の原因かを判断する際に使用してください。",
        "conn_check_btn":        "接続診断を実行",
        "conn_check_running":    "Steam、SteamSpy、Twitchへの接続をテスト中…",
        "pipeline_check_header": "パイプラインチェック（詳細）",
        "pipeline_check_desc":   "実際の取得処理が使うのと全く同じ関数を、1回は直接、1回はDashboardと同様にワーカースレッド内で呼び出します。取得失敗が処理ロジック自体のバグ（両方失敗）なのか、スレッド内実行に特有の問題（直接は成功、スレッドでは失敗）なのかを判別できます。上の接続診断が成功してもDashboardが全て0のままの場合に実行してください。",
        "pipeline_check_btn":    "パイプラインチェックを実行",
        "pipeline_check_running": "取得処理を直接およびワーカースレッド内でテスト中…",
        "pipeline_check_thread_specific": "⚠️ メインスレッドの呼び出しは成功しましたが、ワーカースレッドの呼び出しは失敗しました — これは取得ロジック自体ではなく、バックグラウンドスレッド内での実行に特有の問題（キャッシュ/セッションの不具合など）を示しています。回避策として max_workers を下げるか、該当関数から st.cache_data を外すことを検討してください。",
        "pipeline_check_both_fail": "⚠️ 両方の呼び出しが失敗しました — これは取得ロジック自体のバグであり、スレッドの問題ではありません。上記のエラー詳細を確認してください。",
        "pipeline_check_both_ok": "✅ 両方成功しました — 現時点でこのタイトルの取得処理は正しく動作しています。Dashboardが依然として全て0を表示する場合は「CCUデータを更新」を再試行してください。一時的な問題だった可能性があります。",
        "pipeline_check_both_no_data": "🟡 両方の呼び出しがエラーなく完了しましたが、どちらもSteamからライブ値を取得できませんでした — 取得ロジック自体は問題ありませんが、現在このタイトルに対してSteamのCCUエンドポイントがデータを返していません（一時的なAPIの問題、またはこのタイトル特有の問題の可能性があります）。別のタイトルで試すか、数分待ってから再試行してください。",
        "pipeline_check_main_only_exception": "メインスレッドの呼び出しはエラーになりましたが、ワーカースレッドの呼び出しは成功しました — これは珍しいケースで、2回の呼び出しの間で一時的な問題が発生しただけかもしれません。もう一度実行してみてください。",
        "admin_csv_note":        "ここでアップロードすれば全ページで利用可能: このページでアップロードしたCSVは、Dashboardの次回取得で即座に反映されます — ページごとに再アップロードする必要はありません。",
        "api_loaded":             "Anthropic APIキー読み込み済み",
        "api_missing":            "Anthropic APIキーがありません",
        "bedrock_ok":             "✓ AWS Bedrock認証情報 読み込み済み",
        "bedrock_missing":        "AWS Bedrock認証情報がありません",
        "model_caption":          "モデル: claude-sonnet-4-20250514",
        "ccu_caption":            "CCU: Steam公開API（5分キャッシュ）",
        "engagement_caption":     "エンゲージメント: SteamSpy API（1時間キャッシュ）",
        "cache_age":              "📦 データキャッシュ: {age}",
        "cache_refetch":          "24時間ごとに自動で再取得します。",
        "cache_none":             "キャッシュなし — 読み込み時に取得します。",
        "refresh_now":            "🔄 今すぐ更新",
        "refresh_now_help":       "24時間キャッシュを無視してCCUを再取得します",
        "csvs_loaded":            "SteamDB CSV: {n}/{total} 読み込み済み",
        "csv_missing":            "未取得: {names}",
        "csv_drop_hint":          "steamdb_chart_{appid}.csv を /data に配置すると更新されます",
        "upload_header":          "**SteamDB CSVをアップロード**",
        "upload_caption":         "steamdb_chart_{appid}.csv を直接アップロードできます — リポジトリへのアクセスは不要です。",
        "upload_loaded":          "{n} 件のCSVを読み込み: {names}",
        "signed_in_as":           "ログイン中",
        "sign_out":               "サインアウト",
        "watchlist_header":       "ウォッチリスト",
        "watchlist_max":          "最大5タイトルまで",
        "lang_header":            "レポート言語",
        "last_fetched":           "CCU最終取得: {time}",
        # Hero
        "hero_line1":             "シューター市場",
        "hero_line2":             "インテリジェンス",
        "hero_sub":               "Steamライブ CCUデータ · AI分析 · 競合ベンチマーク · 週次レポートテンプレート — SEGAのパブリッシング・戦略チーム向けオールインワンツール。",
        # Section headers
        "select_analysis":        "分析タイプを選択",
        "live_ccu_header":        "STEAMライブ CCUスナップショット",
        "ai_analysis_header":     "AI分析 — {label}",
        "dataset_header":         "データセット",
        # Dataset picker
        "btn_fps":                "FPS（一人称）",
        "btn_tps":                "TPS（三人称）",
        "btn_both":               "FPS + TPS",
        "games_included":         "対象タイトル: {n} / {total}",
        "select_all":             "すべて選択",
        "clear_all":              "すべて解除",
        "overlap_note":           "ℹ {n} タイトルはもう一方のリストにも含まれています: {names}",
        # Fetch button
        "fetch_ccu_btn":          "ライブCCUデータを取得",
        "fetch_spinner":          "Steam / SteamSpyからライブCCUを取得中…",
        "fetching_game":          "取得中: {name}…",
        "fetch_done":             "{n} タイトルを取得しました",
        "refresh_ccu_btn":        "CCUデータを更新",
        "cache_loaded_toast":     "キャッシュから読み込みました（{age}）",
        # KPI cards
        "kpi_total_ccu":          "合計CCU（追跡対象）",
        "kpi_total_sub":          "シューター {n} タイトルの合計",
        "kpi_wow":                "WoW成長バランス",
        "kpi_wow_sub":            "CSVデータあり {n} タイトル中",
        "kpi_wow_none":           "CSVデータ未読み込み",
        "kpi_yoy":                "YoY成長バランス",
        "kpi_yoy_sub":            "YoYデータあり {n} タイトル中",
        "kpi_mom":                "MoM成長バランス",
        "kpi_mom_sub":            "CSVデータあり {n} タイトル中",
        "kpi_mom_none":           "CSVデータ未読み込み",
        "kpi_csvs":               "SteamDB CSV読み込み数",
        "kpi_csvs_sub":           "完全な履歴データを持つタイトル数",
        "kpi_health":             "平均ピーク健全度",
        "kpi_health_sub":         "現在CCU ÷ 過去最高ピーク",
        "kpi_twitch":             "Twitch総視聴者数",
        "kpi_twitch_sub":         "{n} タイトル（上位100ストリーム合計）",
        "kpi_twitch_none":        "TWITCH_CLIENT_ID を設定すると有効になります",
        "kpi_best_grower":        "YoY最大成長",
        "kpi_best_sub":           "{pct} YoY（SteamDB）",
        "kpi_worst_decline":      "YoY最大下落",
        "kpi_worst_sub":          "{pct} YoY（SteamDB）",
        # Expanders
        "yoy_expander":           "YoY CCU変化 — 成長 {up} / 下落 {down}",
        "wow_expander":           "WoW CCU変化 — 成長 {up} / 下落 {down}",
        "mom_expander":           "MoM CCU変化 — 成長 {up} / 下落 {down}",
        "growth_label":           "成長",
        "decline_label":          "下落",
        "wow_caption":            "最新CSV値と、ちょうど7日前に最も近い行を比較。出典: SteamDB 10分間隔データ。",
        "wow_none":               "CSVデータがまだ読み込まれていません。steamdb_chart_{appid}.csv を /data フォルダに追加してください。",
        "heatmap_expander":       "サブジャンル別CCUヒートマップ",
        "heatmap_caption":        "出典: Steam APIライブCCUをロスター内のサブジャンルタグ別に集計。",
        "table_expander":         "シューターCCUランキング — {genre}",
        "table_footnote":         "レビュー = 好評数 ÷ 総レビュー数。レビュー/日 = 7日間の取得ペース（スナップショット7日以上必要）。Twitch = 上位100ストリームの視聴者数合計（TWITCH_CLIENT_ID要）。— = データなし。* = CSV値使用。",
        "history_expander":       "月間ピークCCU推移 — SteamDBデータ",
        "history_caption":        "出典: SteamDB 10分間隔CSVを月間ピークに集計{note}",
        "history_note_few":       " | 破線 = 主要イベント（ホバーで詳細）",
        "history_note_many":      " | 各ラインにホバーすると主要イベントを表示",
        "chart_history_title":    "月間ピークCCU — 直近24ヶ月（SteamDB）",
        "formulas_caption":       "**WoW%** =（ライブCCU − 7日前CCU）÷ 7日前CCU × 100  |  **MoM%** =（先月平均 − 前々月平均）÷ 前々月平均 × 100  |  **YoY%** =（当月平均 − 前年同月平均）÷ 前年同月平均 × 100  |  出典: SteamDB CSVデータ · 7日前のスナップショットがない場合、WoWは最新CSV行を使用",
        # Table columns
        "col_rank":               "#",
        "col_title":              "タイトル",
        "col_subgenre":           "サブジャンル",
        "col_publisher":          "パブリッシャー",
        "col_f2p":                "基本無料",
        "col_live_ccu":           "ライブCCU",
        "col_yoy":                "YoY",
        "col_year":               "発売年",
        "col_data_source":        "データソース",
        "col_peak_ever":          "過去最高",
        "col_peak_12m":           "12ヶ月ピーク",
        "col_avg_ccu_12m":        "12ヶ月平均",
        "col_mom":                "MoM",
        "col_review":             "レビュー",
        "col_rev_velocity":       "レビュー/日",
        "col_twitch":             "Twitch",
        "col_yoy_source":         "YoYソース",
        "col_owners":             "推定所有者数",
        "col_7d_ago":             "7日前",
        "col_change_ccu":         "変化 (CCU)",
        "col_weekly_change":      "週次変化",
        "col_steam_page":         "Steamページ",
        "col_1yr_ago":            "1年前",
        "col_annual_change":      "年間変化",
        "col_month_change":       "月間変化",
        "col_delta_ccu":          "Δ CCU",
        "col_delta_pct":          "Δ %",
        "col_direction":          "方向",
        "col_reference":          "参照",
        "store_link":             "↗ ストア",
        # Bar charts
        "chart_top25_title":      "ライブCCU上位{n} {genre}タイトル",
        "chart_ranks_title":      "ライブCCU {start}〜{end}位 {genre}タイトル",
        "chart_caption":          "緑 = WoW増加  |  赤 = WoW減少  |  グレー = CSVデータなし  |  {genre}ロスター",
        # Analysis presets
        "run_analysis":           "分析を実行",
        "custom_label":           "またはカスタム質問を入力",
        "custom_placeholder":     "例: Steam上位5つの基本無料シューターのマネタイズモデルを比較…",
        "preset_labels": {
            "ccu_mecha":        "CCUトレンドとメカシューター需要",
            "competitive_gap":  "競合ギャップ分析 — 上位5タイトル",
            "table_stakes":     "2026年 ネットコード・サーバー必須要件",
            "social_metrics":   "Day-1成功を予測するSNS指標",
            "weekly_report":    "週次リテンション＆エンゲージメントレポート",
        },
        "preset_descs": {
            "ccu_mecha":        "Steam上位10シューターのCCUトレンドを前年と比較し、メカシューターの現在の需要を分析します。",
            "competitive_gap":  "上位5タイトルのCCUを分ける具体的なメカニクス、マネタイズパターン、コンテンツケイデンスを特定し、SEGAが活用できる製品ギャップを明らかにします。",
            "table_stakes":     "2026年の競技シューターにおいて、欧米の競技的公平性基準を満たすために必須となるネットコードとサーバーアーキテクチャの要件を分析します。",
            "social_metrics":   "最新の投資家向け資料と市場データに基づき、新作シューターのDay-1成功を予測する主要SNS指標を特定します。",
            "weekly_report":    "上位100シューターのリテンションとエンゲージメントKPIを追跡し、ブレイクアウト中のインディータイトルを抽出する週次レポートテンプレートを作成します。",
        },
        "preset_tags": {
            "ccu_mecha":        "市場",
            "competitive_gap":  "戦略",
            "table_stakes":     "技術",
            "social_metrics":   "SNS",
            "weekly_report":    "レポート",
        },
        "run_btn":                "実行",
        "custom_query_label":     "カスタムクエリ",
        # AI report
        "cache_notice":           "キャッシュから読み込みました — 前回実行からデータに変更はありません。CCUを再取得すると更新されます。",
        "no_ccu_warning":         "先にライブCCUデータを取得してください。",
        "spinner_generating":     "Claudeが分析を生成しています…",
        "no_key_warning":         "AWS Bedrock認証情報が見つかりません。.streamlit/secrets.toml に AWS_ACCESS_KEY_ID_API を追加してください。",
        "no_anthropic_error":     "`anthropic` パッケージをインストールしてください: `pip install anthropic`",
        "auth_error":             "APIキーが無効です。.streamlit/secrets.toml の認証情報を確認してください。",
        "rate_limit_error":       "レート制限に達しました。しばらく待ってから再試行してください。",
        "analysis_failed":        "分析に失敗しました: {e}",
        "auto_archived":          "📋 自動アーカイブ: {f}",
        # Downloads
        "download_report_header": "レポートをダウンロード",
        "dl_md":                  "Markdownをダウンロード",
        "dl_html":                "HTMLをダウンロード",
        "dl_pdf":                 "PDFをダウンロード",
        "dl_pptx_btn":            "PowerPointをダウンロード",
        "dl_pptx_file":           ".pptxをダウンロード",
        "dl_pptx_error":          "PPTXの生成に失敗しました。python-pptx がインストールされているか確認してください: pip install python-pptx",
        "dl_pdf_missing":         "PDF: `reportlab` をインストールしてください",
        "spinner_pptx":           "スライドを作成中…（ページのスナップショットを撮影しています）",
        "pptx_fallback_note":     "スナップショット出力が利用できません（pip install playwright pillow kaleido の後 `playwright install chromium` を実行してください）— テキストベースの資料を生成しました。",
        "regen_btn":              "再生成",
        "archive_btn":            "🗄 アーカイブ",
        "archive_saved":          "保存しました: {f}",
        "archive_failed":         "アーカイブに失敗しました",
        # Follow-up chat
        "chat_header":            "フォローアップチャット",
        "chat_subtext":           "— このレポートについてClaudeに質問できます",
        "chat_placeholder":       "このレポートについて質問する…",
        "chat_clear":             "チャット履歴をクリア",
        "chat_error":             "チャットエラー: {e}",
        # Empty state
        "empty_title":            "分析が選択されていません",
        "empty_sub":              "上でライブCCUデータを取得し、分析タイプを選択するかカスタム質問を入力してレポートを生成してください。",
        # Drilldown
        "drilldown_header":       "ゲーム別ディープダイブ分析",
        "back_btn":               "← ダッシュボードに戻る",
        "drilldown_select":       "ディープダイブするタイトルを選択...",
        "drilldown_btn":          "ディープダイブ",
        "drilldown_no_data":      "ゲームデータが見つかりません — 先にCCUデータを取得してください。",
        "drilldown_no_key":       "AWS Bedrock認証情報が見つかりません。.streamlit/secrets.toml に追加してください。",
        "drilldown_spinner":      "ディープダイブ分析を生成しています…",
        "drilldown_failed":       "ディープダイブに失敗しました: {e}",
        "drilldown_dl":           "ディープダイブをダウンロード (.md)",
        "drilldown_no_csv":       "{name} のSteamDB CSVがありません。steamdb_chart_{appid}.csv を /data に配置すると推移チャートが有効になります。",
        "chart_drilldown_title":  "{name} — 月間ピークCCU",
        "pill_peak_ever":         "過去最高ピーク",
        "pill_peak_12m":          "12ヶ月ピーク",
        "pill_avg_12m":           "12ヶ月平均CCU",
        "pill_mom":               "MoMトレンド",
        "trace_peak":             "ピークCCU",
        "trace_avg":              "平均CCU",
        "trace_live":             "ライブCCU",
        "no_hist_info":           "このタイトルの履歴CSVデータがありません。steamdb_chart_{appid}.csv を /data に配置してください。",
        "yoy_caption":            "スナップショット = アプリ蓄積の実測12ヶ月データ（約13ヶ月後に有効）· SteamDB CSV = 前年同月実測YoY · SteamSpyプロキシ = エンゲージメント推定",
        "yoy_none":               "YoYデータがありません — 先にCCUデータを取得してください。",
        # Monthly analysis
        "monthly_header":         "月次分析",
        "monthly_none":           "アーカイブ済みレポートはまだありません。毎週月曜に自動アーカイブされるほか、上のアーカイブボタンで現在のレポートを保存できます。",
        "monthly_count":          "{n} 件のアーカイブ済みレポート / {m} ヶ月分 · /data/report_archive/ に保存",
        "monthly_expander":       "📅 月次比較 — 週次レポートの比較と精度チェック",
        "monthly_select":         "分析する月を選択",
        "monthly_reports_for":    "{month} のアーカイブ済みレポート: **{n} 件**",
        "monthly_titles":         "{n} タイトル",
        "monthly_drift":          "**月間CCUドリフト**（{a} → {b}）:",
        "monthly_one":            "今月のレポートは1件のみです — ドリフト表示には2件以上必要です。毎週月曜に蓄積されます。",
        "monthly_ai_header":      "**AI精度・トレンドチェック**",
        "monthly_ai_caption":     "Claudeが当月の週次レポートを比較し、月間ドリフトとの乖離を指摘します。",
        "monthly_no_key":         "月次分析を実行するにはAWS Bedrock認証情報を設定してください。",
        "monthly_run_btn":        "月次分析を実行",
        "monthly_spinner":        "月次分析を実行しています…",
        "monthly_failed":         "月次分析に失敗しました: {e}",
        "monthly_dl":             "⬇ 月次レポートをダウンロード (.md)",
        "col_monthly_drift":      "月間ドリフト",
        "col_latest_wow":         "最新WoW",
        # Footer
        "footer_brand":           "SEGA シューターインテリジェンス",
        "footer_note":            "データ出典: Steam公開API · Powered by Claude · 社内分析専用",
        # Watchlist section
        "watchlist_section":      "ウォッチリスト登録タイトル",
        # Values
        "yes": "はい",
        "no":  "いいえ",
        "up":  "上昇",
        "down":"下降",
        "flat":"横ばい",
    },
}


def T(key: str, **kwargs) -> str:
    lang = st.session_state.get("report_language", "English")
    if lang not in TRANSLATIONS:
        lang = "English"
    text = TRANSLATIONS[lang].get(key, TRANSLATIONS["English"].get(key, key))
    return text.format(**kwargs) if kwargs else text

# ─────────────────────────────────────────────────────────────
# GAME DRILL-DOWN PROMPT
# ─────────────────────────────────────────────────────────────

def build_drilldown_prompt(game: dict, historical: dict, language: str = "English") -> str:
    hs   = game.get("hist_summary", {})
    mdf  = historical.get(game["app_id"])
    src  = "SteamDB CSV" if game.get("has_hist") else "SteamSpy proxy"

    # Last 12 months table
    history_table = ""
    if mdf is not None and not mdf.empty:
        last12 = mdf.sort_values("month").tail(12)
        rows = []
        for _, row in last12.iterrows():
            rows.append(f" {row['month'].strftime('%Y-%m')}  peak={int(row['peak_ccu']):,}  avg={int(row['avg_ccu']):,}")
        history_table = "\n".join(rows)
    else:
        history_table = " (no CSV data available)"

    lang_instruction = (
        " IMPORTANT: Write the entire analysis in Japanese (日本語). "
        "Use professional business Japanese suitable for senior management. "
        "All section headers, bullet points, and analysis must be in Japanese. "
        "Game titles may be kept in their original English/romanised form. "
        "Never switch back to English partway through the report."
    ) if language == "Japanese" else ""

    return f"""You are a senior games market analyst at SEGA's internal strategy team.{lang_instruction}

Produce a focused deep-dive intelligence report on ONE game for SEGA's competitive analysis team.

GAME: {game['name']}
Sub-genre: {game.get('sub', 'N/A')}
Publisher: {game.get('publisher', 'N/A')}
F2P: {'Yes' if game.get('f2p') else 'No'}
Current CCU: {game.get('ccu', 0):,}
YoY trend: {game.get('yoy', 'N/A')} (source: {src})
Avg playtime (2 weeks): {game.get('avg_2w_hrs', 0)} hrs
Review score: {f"{game.get('review_pct')}%" if game.get('review_pct') else 'N/A'}
All-time peak CCU: {hs.get('peak_ever', 'N/A'):,} (if int, else {hs.get('peak_ever', 'N/A')})
12-month peak CCU: {hs.get('peak_12m', 'N/A')}
12-month avg CCU:  {hs.get('avg_12m', 'N/A')}
MoM trend:         {hs.get('mom_trend', 'N/A')}

MONTHLY CCU HISTORY (last 12 months, SteamDB):
{history_table}

Write a structured report with these sections:
## 1. Executive Summary
One paragraph — current health, trajectory, key risk or opportunity for SEGA.

## 2. CCU Trend Analysis
Interpret the monthly history. Identify growth phases, declines, seasonal patterns, notable spikes/drops. Quantify where possible.

## 3. Player Engagement & Retention Signals
Use playtime and review score to assess engagement depth and churn risk.

## 4. Competitive Position
How does this title's CCU and trajectory compare to its sub-genre peers? Is it gaining or losing share?

## 5. Monetisation Model Assessment
Evaluate the F2P vs premium model in context of this title's performance. What does it signal about player willingness-to-pay in this sub-genre?

## 6. Strategic Implications for SEGA
3–5 concrete, actionable bullet points for SEGA's product / publishing strategy based on this data.

Be specific, data-driven, and concise. Avoid generic observations. Attribute data sources inline."""



# ─────────────────────────────────────────────────────────────
# OTP AUTHENTICATION  (hardened)
#
# Security changes vs. the original single-file version:
#   1. COOKIE_SIGNING_KEY has NO fallback default. If it isn't set in
#      secrets.toml, the app refuses to start rather than signing
#      cookies with a publicly-known string.
#   2. OTP send + verify attempts are persisted server-side (JSON file
#      keyed by email), not just in st.session_state. A page refresh
#      or new browser tab no longer resets the attempt counter.
#   3. Cookie expiry slides forward on every authenticated page load,
#      so active users stay signed in; inactive users still expire
#      24h after their last visit.
# ─────────────────────────────────────────────────────────────

# Filename of the entry-point script (the Dashboard page), used by every
# "go back to Dashboard" page_link across the other pages. If you rename the
# main script when deploying, update this constant to match — it's the only
# place that needs to change.
HOME_PAGE = "shooter_sentiment.py"

ALLOWED_DOMAIN     = "@segaamerica.com"
OTP_EXPIRY_SECS    = 600       # 10 minutes
COOKIE_EXPIRY_SECS = 86400     # 1 day, sliding
COOKIE_NAME        = "sega_shooter_auth"

MAX_VERIFY_ATTEMPTS  = 5        # per OTP code, server-side enforced
MAX_SEND_ATTEMPTS     = 8       # code re-sends per rolling window
SEND_WINDOW_SECS      = 900     # 15 minutes


def _otp_attempts_path() -> Path:
    return DATA_DIR / "otp_attempts.json"


def _load_otp_attempts() -> dict:
    p = _otp_attempts_path()
    try:
        if p.exists():
            with open(p) as f:
                return json.load(f)
    except Exception:
        pass
    return {}


def _save_otp_attempts(data: dict) -> None:
    try:
        with open(_otp_attempts_path(), "w") as f:
            json.dump(data, f)
    except Exception:
        pass  # read-only filesystem — fail silently, same posture as other state files


def _check_send_allowed(email: str) -> tuple[bool, int]:
    """Server-side rate limit on OTP *sends*. Returns (allowed, sends_remaining)."""
    data  = _load_otp_attempts()
    now   = time.time()
    entry = data.get(email, {})
    window_start = entry.get("send_window_start", 0)
    if now - window_start > SEND_WINDOW_SECS:
        entry["send_window_start"] = now
        entry["send_count"] = 0
    entry["send_count"] = entry.get("send_count", 0) + 1
    data[email] = entry
    _save_otp_attempts(data)
    remaining = MAX_SEND_ATTEMPTS - entry["send_count"]
    return remaining >= 0, max(0, remaining)


def _check_verify_allowed(email: str) -> tuple[bool, int]:
    """Server-side rate limit on OTP *verify* attempts for the current code.
    Returns (allowed, attempts_remaining_before_this_one)."""
    data  = _load_otp_attempts()
    entry = data.get(email, {})
    used  = entry.get("verify_count", 0)
    return used < MAX_VERIFY_ATTEMPTS, MAX_VERIFY_ATTEMPTS - used


def _record_verify_attempt(email: str) -> None:
    data  = _load_otp_attempts()
    entry = data.get(email, {})
    entry["verify_count"] = entry.get("verify_count", 0) + 1
    data[email] = entry
    _save_otp_attempts(data)


def _reset_verify_attempts(email: str) -> None:
    data = _load_otp_attempts()
    if email in data:
        data[email]["verify_count"] = 0
        _save_otp_attempts(data)


def _send_otp(email: str, code: str) -> bool:
    try:
        import boto3
        ses = boto3.client(
            "ses",
            region_name=st.secrets.get("AWS_SES_REGION", "us-east-1"),
            aws_access_key_id=st.secrets.get("AWS_ACCESS_KEY_ID", ""),
            aws_secret_access_key=st.secrets.get("AWS_SECRET_ACCESS_KEY", ""),
        )
        ses.send_email(
            Source=st.secrets.get("EMAIL_FROM", "noreply@segaamerica.com"),
            Destination={"ToAddresses": [email]},
            Message={
                "Subject": {"Data": "SEGA Shooter Intelligence — Your verification code", "Charset": "UTF-8"},
                "Body": {
                    "Text": {"Data": f"Your verification code is: {code}\n\nExpires in 10 minutes.", "Charset": "UTF-8"},
                    "Html": {"Data": f"""
                    <div style="font-family:Arial,sans-serif;max-width:480px;margin:0 auto;padding:32px 24px;">
                      <div style="font-size:22px;font-weight:900;letter-spacing:.1em;color:#0057FF;margin-bottom:4px;">SEGA</div>
                      <div style="font-size:14px;color:#444;margin-bottom:28px;">Shooter Market Intelligence</div>
                      <div style="font-size:14px;color:#222;margin-bottom:16px;">Your verification code is:</div>
                      <div style="font-size:42px;font-weight:900;letter-spacing:.18em;color:#050818;
                                  background:#EEF3FF;border-radius:8px;padding:18px 24px;
                                  display:inline-block;margin-bottom:24px;">{code}</div>
                      <div style="font-size:12px;color:#888;">Expires in 10 minutes. If you didn't request this, ignore this email.</div>
                    </div>""", "Charset": "UTF-8"},
                },
            },
        )
        return True
    except Exception as e:
        st.error(f"Failed to send email: {e}")
        return False


def _cookie_secret() -> str:
    """Return the cookie-signing secret. Hard-fails (no fallback) if unset —
    see require_auth() for where this is enforced before any cookie logic runs."""
    return st.secrets.get("COOKIE_SIGNING_KEY", "")


def _sign_cookie(email: str) -> str:
    secret  = _cookie_secret()
    expiry  = int(time.time()) + COOKIE_EXPIRY_SECS
    payload = f"{email}|{expiry}"
    sig     = hmac.new(secret.encode(), payload.encode(), hashlib.sha256).hexdigest()
    return base64.urlsafe_b64encode(f"{payload}|{sig}".encode()).decode()


def _verify_cookie(token: str) -> tuple[str, int] | None:
    """Validate a signed cookie token. Returns (email, expiry_unix_ts) on
    success, or None if missing, tampered, or expired. Exposing the expiry
    (rather than just the email) lets the caller decide whether a refresh is
    actually needed, instead of re-signing on every single page load."""
    try:
        secret  = _cookie_secret()
        decoded = base64.urlsafe_b64decode(token.encode()).decode()
        email, expiry_str, sig = decoded.rsplit("|", 2)
        payload  = f"{email}|{expiry_str}"
        expected = hmac.new(secret.encode(), payload.encode(), hashlib.sha256).hexdigest()
        if not hmac.compare_digest(sig, expected):
            return None
        expiry = int(expiry_str)
        if int(time.time()) > expiry:
            return None
        return email, expiry
    except Exception:
        return None


def require_auth() -> None:
    """Gate a page behind OTP authentication.

    Call this immediately after st.set_page_config() + inject_css() at the
    top of EVERY page (Streamlit's pages/ directory model runs each page as
    an independent script — there is no shared "main file ran first"
    guarantee, so every page must perform this check itself).

    On success: renders the "Signed in as / Sign out" sidebar widget and
    returns normally. On failure: renders the login form and calls
    st.stop() — execution of the page never continues past that point.
    """
    # ── Hard-fail if the cookie-signing secret isn't configured ──
    # No fallback default: signing cookies with a guessable string would let
    # anyone forge a valid session, so we refuse to run instead.
    if not _cookie_secret():
        st.error(
            "Server misconfiguration: `COOKIE_SIGNING_KEY` is not set in "
            "`.streamlit/secrets.toml`. Authentication cannot proceed securely "
            "until this is configured. Contact an administrator."
        )
        st.stop()

    try:
        import extra_streamlit_components as stx
        _cookie_manager  = stx.CookieManager(key="auth_cookies")
        _existing_cookie = _cookie_manager.get(COOKIE_NAME)
    except Exception:
        _cookie_manager  = None
        _existing_cookie = None

    _verified = _verify_cookie(_existing_cookie) if _existing_cookie else None
    _cookie_email, _cookie_expiry = _verified if _verified else (None, None)

    for _k, _v in [
        ("auth_verified", False), ("auth_email", ""),
        ("otp_code", ""),        ("otp_email", ""),
        ("otp_expiry", 0),       ("otp_sent", False),
    ]:
        if _k not in st.session_state:
            st.session_state[_k] = _v

    if _cookie_email and not st.session_state.auth_verified:
        st.session_state.auth_verified = True
        st.session_state.auth_email    = _cookie_email

    # ── Sliding expiry: re-issue the cookie once its remaining lifetime
    # drops below half of COOKIE_EXPIRY_SECS, so active users never get
    # logged out mid-session while inactive users still expire ~24h after
    # their last visit. Only refreshing when actually needed (rather than
    # on every single page load) avoids hammering the underlying cookie
    # component with writes. The key is a STABLE string, not time-based —
    # using a key that changes every second (as an earlier version of this
    # function did) makes the custom component remount on every single
    # rerun instead of updating in place, which is why the cookie wasn't
    # reliably persisting across page navigations.
    if _cookie_email and _cookie_manager and _cookie_expiry is not None:
        _remaining = _cookie_expiry - int(time.time())
        if _remaining < COOKIE_EXPIRY_SECS / 2:
            _fresh_token = _sign_cookie(_cookie_email)
            _cookie_manager.set(COOKIE_NAME, _fresh_token,
                expires_at=None, key="slide_auth_cookie")

    if not st.session_state.auth_verified:
        st.markdown("""
        <style>
        .auth-wrap { max-width:420px;margin:6rem auto;padding:2.5rem 2.5rem 2rem;
                     background:var(--surface);border:1px solid var(--border);
                     border-top:3px solid var(--blue);border-radius:0 0 10px 10px; }
        .auth-logo  { font-family:'Inter Tight',sans-serif;font-size:1.6rem;font-weight:900;
                      letter-spacing:.12em;color:var(--blue) !important;margin-bottom:.2rem; }
        .auth-title { font-family:'Inter Tight',sans-serif;font-size:1rem;font-weight:700;
                      color:var(--text) !important;margin-bottom:.25rem; }
        .auth-sub   { font-size:.8rem;color:var(--muted) !important;margin-bottom:1.5rem; }
        .auth-note  { font-size:.72rem;color:var(--muted) !important;margin-top:1rem;
                      text-align:center;line-height:1.5; }
        </style>
        """, unsafe_allow_html=True)

        _lc, _mc, _rc = st.columns([1, 2, 1])
        with _mc:
            st.markdown("""
            <div class="auth-wrap">
              <div class="auth-logo">SEGA</div>
              <div class="auth-title">Shooter Market Intelligence</div>
              <div class="auth-sub">Sign in with your SEGA America email to continue</div>
            </div>
            """, unsafe_allow_html=True)

            if not st.session_state.otp_sent:
                _email_in = st.text_input("Email address", placeholder="you@segaamerica.com",
                    label_visibility="collapsed", key="auth_email_input")
                if st.button("Send verification code", use_container_width=True):
                    _email_norm = _email_in.strip().lower()
                    if not _email_norm.endswith(ALLOWED_DOMAIN):
                        st.error(f"Access restricted to {ALLOWED_DOMAIN} addresses.")
                    else:
                        _allowed, _remaining = _check_send_allowed(_email_norm)
                        if not _allowed:
                            st.error(
                                "Too many verification codes requested. "
                                "Please wait 15 minutes and try again."
                            )
                        else:
                            _code = str(random.randint(100000, 999999))
                            if _send_otp(_email_norm, _code):
                                _reset_verify_attempts(_email_norm)
                                st.session_state.otp_code   = _code
                                st.session_state.otp_email  = _email_norm
                                st.session_state.otp_expiry = time.time() + OTP_EXPIRY_SECS
                                st.session_state.otp_sent   = True
                                st.rerun()
            else:
                st.info(f"Code sent to **{st.session_state.otp_email}** — check your inbox.")
                _code_in = st.text_input("6-digit code", placeholder="123456",
                    label_visibility="collapsed", max_chars=6, key="auth_code_input")
                if st.button("Verify code", use_container_width=True):
                    _email = st.session_state.otp_email
                    _allowed, _remaining = _check_verify_allowed(_email)
                    if not _allowed:
                        st.error("Too many attempts. Please request a new code.")
                        st.session_state.otp_sent = False
                    elif time.time() > st.session_state.otp_expiry:
                        st.error("Code has expired. Please request a new one.")
                        st.session_state.otp_sent = False
                    elif _code_in.strip() != st.session_state.otp_code:
                        _record_verify_attempt(_email)
                        _, _rem_after = _check_verify_allowed(_email)
                        st.error(f"Incorrect code. {_rem_after} attempt{'s' if _rem_after != 1 else ''} remaining.")
                    else:
                        _reset_verify_attempts(_email)
                        st.session_state.auth_verified = True
                        st.session_state.auth_email    = _email
                        st.session_state.otp_code      = ""
                        if _cookie_manager:
                            _token = _sign_cookie(st.session_state.auth_email)
                            _cookie_manager.set(COOKIE_NAME, _token,
                                expires_at=None, key="set_auth_cookie")
                        st.rerun()

                if st.button("← Use a different email", key="auth_back"):
                    st.session_state.otp_sent = False
                    st.session_state.otp_code = ""
                    st.rerun()

            st.markdown(
                f'<div class="auth-note">Restricted to {ALLOWED_DOMAIN} addresses only.<br>'
                f'Codes expire after 10 minutes · Session lasts 24 hours.</div>',
                unsafe_allow_html=True)

        st.stop()

    with st.sidebar:
        st.markdown(
            f'<div style="font-size:.7rem;font-weight:600;color:var(--muted);margin-bottom:.5rem;">'
            f'{T("signed_in_as")}<br>'
            f'<span style="color:var(--text);font-weight:700;">{st.session_state.auth_email}</span>'
            f'</div>',
            unsafe_allow_html=True)
        if st.button(T("sign_out"), key="auth_sign_out_btn"):
            if _cookie_manager:
                _cookie_manager.delete(COOKIE_NAME, key="delete_auth_cookie")
            for _k in ["auth_verified", "auth_email", "otp_sent", "otp_code",
                       "otp_email", "otp_expiry"]:
                st.session_state[_k] = False if _k == "auth_verified" else ""
            st.rerun()

# ─────────────────────────────────────────────────────────────
# SHARED CHROME  (topbar + language toggle, footer)
# Every page calls render_topbar() near the top and render_footer()
# at the bottom, so branding and the EN/JP toggle stay consistent
# across the multipage app.
# ─────────────────────────────────────────────────────────────

def render_topbar() -> None:
    _tc = st.columns([7, 1])
    with _tc[0]:
        st.markdown("""
<div class="topbar">
  <div class="topbar-logo"><span class="seg">SEGA</span> &nbsp;SHOOTER INTELLIGENCE</div>
  <div class="topbar-divider"></div>
  <div class="topbar-label">{subtitle}</div>
</div>""".format(subtitle=T("topbar_subtitle")), unsafe_allow_html=True)
    with _tc[1]:
        _lang = st.selectbox(
            "Language", options=["EN", "JP"],
            index=1 if st.session_state.report_language == "Japanese" else 0,
            label_visibility="collapsed", key="lang_toggle",
        )
        _new_lang = "Japanese" if _lang == "JP" else "English"
        if _new_lang != st.session_state.report_language:
            st.session_state.report_language = _new_lang
            st.session_state.ai_report = ""
            st.session_state.report_cache = {}
            st.rerun()


# (key, page path, label, icon) for every page in the app, in display order.
# HOME_PAGE is used here rather than a literal filename so this stays correct
# if the entry script is ever renamed (see the HOME_PAGE constant above).
_NAV_PAGES = [
    ("dashboard",      HOME_PAGE,                       "Dashboard",       "📊"),
    ("weekly_report",  "pages/1_Weekly_Report.py",       "Weekly Report",   "📝"),
    ("deep_dive",      "pages/2_Deep_Dive.py",           "Deep Dive",       "🔍"),
    ("monthly",        "pages/3_Monthly_Analysis.py",    "Monthly Analysis","📅"),
    ("admin",          "pages/4_Admin.py",                "Admin",          "⚙️"),
]


def render_nav_tabs(current: str) -> None:
    """Persistent tab-style navigation bar, identical on every page.

    st.tabs() itself can't navigate between separate multipage-app pages —
    it only switches content within a single page — so this is a styled row
    of safe_page_link calls instead, with the current page shown as a plain
    highlighted label (not a link to itself). `current` must be one of the
    keys in _NAV_PAGES ("dashboard", "weekly_report", "deep_dive", "monthly",
    "admin").
    """
    st.markdown('<div class="nav-tabs">', unsafe_allow_html=True)
    _cols = st.columns(len(_NAV_PAGES))
    for _col, (_key, _path, _label, _icon) in zip(_cols, _NAV_PAGES):
        with _col:
            if _key == current:
                st.markdown(
                    f'<div class="nav-tab-active">{_icon} {_label}</div>',
                    unsafe_allow_html=True,
                )
            else:
                safe_page_link(_path, label=_label, icon=_icon)
    st.markdown('</div>', unsafe_allow_html=True)


def summarize_fetch_health(ccu_data: list[dict]) -> dict:
    """Summarize how the live fetch actually went, so a systemic failure
    (e.g. Steam/SteamSpy unreachable, or rate-limited by burst concurrency)
    is visible as a clear signal instead of just looking like every title
    genuinely has 0 concurrent players — which is indistinguishable from a
    real (if surprising) result unless something explicitly flags it.

    live_count   — titles where the Steam CCU API call itself succeeded
    csv_fallback — titles where the live call failed but a CSV/snapshot
                   number filled in instead (ccu may still be meaningful)
    zero_count   — titles where ccu ended up at exactly 0 (live failed AND
                   no fallback data existed)
    """
    total        = len(ccu_data)
    live_count   = sum(1 for r in ccu_data if r.get("ccu_live"))
    csv_fallback = sum(1 for r in ccu_data if r.get("ccu_from_csv"))
    zero_count   = sum(1 for r in ccu_data if (r.get("ccu") or 0) == 0)
    return {
        "total": total,
        "live_count": live_count,
        "csv_fallback": csv_fallback,
        "zero_count": zero_count,
        "live_pct": round(live_count / total * 100) if total else 0,
        "looks_systemic": total > 0 and live_count == 0 and csv_fallback == 0,
    }


def run_connectivity_probe(test_app_id: int = 730, test_name: str = "Counter-Strike 2") -> list[dict]:
    """Connectivity test against each upstream API — used by the Admin page's
    connectivity check.

    Calls the REAL application functions (fetch_ccu, fetch_steamspy), not raw
    requests.get() — an earlier version of this probe bypassed fetch_ccu()
    entirely, which meant a passing probe only proved the network and Steam's
    API were reachable, NOT that the actual fetch pipeline (retry wrapper,
    cache decorator, semaphore throttling) worked. That gap is exactly why a
    user could see this probe succeed while the real Dashboard fetch still
    returned all zeros — the probe wasn't testing the same code.

    Clears the relevant @st.cache_data caches first so a stale cached failure
    can't be mistaken for a fresh one.
    """
    results = []

    # Steam CCU — via the real fetch_ccu(), main thread
    try:
        fetch_ccu.clear()
    except Exception:
        pass
    t0 = time.time()
    try:
        count = fetch_ccu(test_app_id)
        elapsed = round((time.time() - t0) * 1000)
        results.append({"api": "Steam CCU (fetch_ccu)", "ok": count is not None,
                         "detail": f"player_count={count} · {elapsed}ms"})
    except Exception as e:
        elapsed = round((time.time() - t0) * 1000)
        results.append({"api": "Steam CCU (fetch_ccu)", "ok": False,
                         "detail": f"{type(e).__name__}: {e} · {elapsed}ms"})

    # SteamSpy — via the real fetch_steamspy(), main thread
    try:
        fetch_steamspy.clear()
    except Exception:
        pass
    t0 = time.time()
    try:
        data = fetch_steamspy(test_app_id)
        elapsed = round((time.time() - t0) * 1000)
        results.append({"api": "SteamSpy (fetch_steamspy)", "ok": bool(data),
                         "detail": f"keys={list(data.keys())[:5]} · {elapsed}ms"})
    except Exception as e:
        elapsed = round((time.time() - t0) * 1000)
        results.append({"api": "SteamSpy (fetch_steamspy)", "ok": False,
                         "detail": f"{type(e).__name__}: {e} · {elapsed}ms"})

    # Twitch (only if configured)
    if st.secrets.get("TWITCH_CLIENT_ID"):
        t0 = time.time()
        try:
            v = fetch_twitch_viewers(test_app_id, test_name)
            elapsed = round((time.time() - t0) * 1000)
            results.append({"api": "Twitch", "ok": v is not None,
                             "detail": f"viewer_count={v} · {elapsed}ms"})
        except Exception as e:
            elapsed = round((time.time() - t0) * 1000)
            results.append({"api": "Twitch", "ok": False,
                             "detail": f"{type(e).__name__}: {e} · {elapsed}ms"})
    else:
        results.append({"api": "Twitch", "ok": None,
                         "detail": "Not configured — TWITCH_CLIENT_ID not in secrets (this is fine if you don't use Twitch data)"})

    return results


def run_pipeline_probe(test_app_id: int = 730, test_name: str = "Counter-Strike 2") -> dict:
    """Call _fetch_one_game() — the exact function the real Dashboard fetch
    uses — directly, once in the main thread and once inside an actual
    ThreadPoolExecutor worker, and report the full result or exception from
    each.

    This exists specifically to distinguish two different failure modes that
    look identical from the Dashboard (every title at 0):
      1. A bug in _fetch_one_game() itself (e.g. a call to something that
         doesn't actually exist) — would fail in BOTH the main-thread and
         worker-thread call below.
      2. Something specific to running inside a worker thread (e.g. a
         Streamlit caching/session quirk without an active ScriptRunContext)
         — would succeed in the main-thread call but fail in the worker-
         thread call.

    run_connectivity_probe() only proves the network and upstream APIs are
    reachable; it doesn't exercise _fetch_one_game() or the thread pool at
    all, so it can pass while the real fetch still fails. This probe closes
    that gap.
    """
    game = {"app_id": test_app_id, "name": test_name, "sub": "Tactical / Competitive",
             "publisher": "Valve", "f2p": True, "year": 2023}

    # Use real historical/raw data if a CSV happens to exist for this app_id,
    # otherwise empty dicts — _fetch_one_game() handles both cases.
    try:
        historical = load_all_historical(frozenset({test_app_id}))
    except Exception:
        historical = {}
    try:
        raw_data = load_all_raw(frozenset({test_app_id}))
    except Exception:
        raw_data = {}
    try:
        snapshots = load_ccu_snapshots()
    except Exception:
        snapshots = []

    out = {"main_thread": None, "worker_thread": None}

    # ── Main thread — direct call ──
    try:
        result = _fetch_one_game(game, historical, raw_data, snapshots)
        if result.get("ccu_live"):
            out["main_thread"] = {"status": "ok",
                "detail": f"ccu={result['ccu']:,} (live) · yoy_source={result['yoy_source']}"}
        else:
            out["main_thread"] = {"status": "ran_but_no_live_data",
                "detail": f"completed without error, but the live Steam call itself returned nothing "
                          f"(ccu={result['ccu']}, ccu_live=False, yoy_source={result['yoy_source']})"}
    except Exception as e:
        out["main_thread"] = {"status": "exception", "detail": f"{type(e).__name__}: {e}"}

    # ── Worker thread — via a real one-off ThreadPoolExecutor, same as the
    # Dashboard's actual fetch loop uses ──
    try:
        with concurrent.futures.ThreadPoolExecutor(max_workers=1) as _pool:
            _fut = _pool.submit(_fetch_one_game, game, historical, raw_data, snapshots)
            result = _fut.result(timeout=20)
        if result.get("ccu_live"):
            out["worker_thread"] = {"status": "ok",
                "detail": f"ccu={result['ccu']:,} (live) · yoy_source={result['yoy_source']}"}
        else:
            out["worker_thread"] = {"status": "ran_but_no_live_data",
                "detail": f"completed without error, but the live Steam call itself returned nothing "
                          f"(ccu={result['ccu']}, ccu_live=False, yoy_source={result['yoy_source']})"}
    except Exception as e:
        out["worker_thread"] = {"status": "exception", "detail": f"{type(e).__name__}: {e}"}

    return out



def safe_page_link(page: str, label: str, icon: str | None = None, **kwargs) -> None:
    """st.page_link wrapper that degrades to a disabled-looking caption instead
    of raising StreamlitPageNotFoundError and crashing the whole page.

    This matters specifically for this app: files have repeatedly gone out of
    sync between what's generated and what's actually present in the deployed
    repo (a missing or misnamed file under pages/, or the entry script renamed
    during deployment). One missing target should never take down an entire
    page — it should just show a flat note so the rest of the page still works.
    """
    try:
        st.page_link(page, label=label, icon=icon, **kwargs)
    except Exception:
        st.caption(
            f"{icon or ''} {label} — page link unavailable "
            f"(`{page}` not found in this deployment; check it exists in the repo with this exact name)."
        )



def render_footer() -> None:
    st.markdown(f"""
<div class="footer">
  <div class="footer-brand">{T("footer_brand")}</div>
  <div class="footer-note">{T("footer_note")}</div>
</div>
""", unsafe_allow_html=True)
