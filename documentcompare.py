"""
Document Comparator — SEGA-branded Streamlit App
=================================================
Run with:  streamlit run doc_compare.py

Required:  pip install streamlit anthropic python-docx openpyxl pymupdf boto3 streamlit-cookies-controller python-pptx
"""

import io
import base64
import tempfile
import os
import random
import time
import hashlib
import hmac

import streamlit as st

# ─────────────────────────────────────────────────────────────
# PAGE CONFIG
# ─────────────────────────────────────────────────────────────

st.set_page_config(
    page_title="SEGA Doc Comparator",
    page_icon=":material/compare:",
    layout="wide",
    initial_sidebar_state="collapsed",
)

# ─────────────────────────────────────────────────────────────
# SEGA BRAND STYLES
# ─────────────────────────────────────────────────────────────

st.markdown("""
<style>
@import url('https://fonts.googleapis.com/css2?family=Inter+Tight:wght@400;700;800;900&family=Poppins:wght@300;400;500;600&display=swap');

:root,
html[data-theme="light"],
html[data-theme="dark"],
[data-theme="light"],
[data-theme="dark"] {
    color-scheme: dark !important;
    --bg:           #0a0c1a;
    --surface:      #0f1120;
    --surface2:     #141728;
    --surface3:     #1a1e30;
    --border:       #232640;
    --border-hi:    #323760;
    --blue:         #4080ff;
    --blue-lo:      #1a3acc;
    --blue-glow:    rgba(64,128,255,0.16);
    --blue-glow-hi: rgba(64,128,255,0.32);
    --text:         #eef0fa;
    --text-dim:     #b8bcd4;
    --muted:        #5a5f82;
    --pos:          #20c65a;
    --pos-dim:      rgba(32,198,90,0.14);
    --neg:          #ff3d52;
    --neg-dim:      rgba(255,61,82,0.14);
    --amber:        #f0a500;
    --amber-dim:    rgba(240,165,0,0.14);
}

html, body {
    background: var(--bg) !important;
    color: var(--text) !important;
    color-scheme: dark !important;
}
.stApp,
.stApp > div,
section[data-testid="stAppViewContainer"],
section[data-testid="stAppViewContainer"] > div,
div[data-testid="stMain"],
div[data-testid="stVerticalBlock"],
div[data-testid="stHorizontalBlock"],
.main .block-container,
.block-container {
    background-color: var(--bg) !important;
    color: var(--text) !important;
}

*, *::before, *::after { font-family: 'Poppins', sans-serif; box-sizing: border-box; }

p, span, div, li, td, th, label,
h1, h2, h3, h4, h5, h6,
.stMarkdown, .stMarkdown p, .stMarkdown span,
[data-testid="stText"],
[data-testid="stMarkdownContainer"],
[data-testid="stMarkdownContainer"] p,
[data-testid="stMarkdownContainer"] span,
[data-testid="stMarkdownContainer"] li,
[data-testid="stMarkdownContainer"] strong,
[data-testid="stMarkdownContainer"] em,
[class*="css"] { color: var(--text) !important; }

.stCaption, [data-testid="stCaptionContainer"],
[data-testid="stCaptionContainer"] p { color: var(--muted) !important; }

code { background: var(--surface3) !important; color: var(--blue) !important; padding: 0.1em 0.4em; border-radius: 3px; }

/* Ensure markdown bold/italic render correctly and are never overridden */
[data-testid="stMarkdownContainer"] strong,
[data-testid="stMarkdownContainer"] b { font-weight: 700 !important; font-style: normal !important; }
[data-testid="stMarkdownContainer"] em,
[data-testid="stMarkdownContainer"] i { font-style: italic !important; font-weight: normal !important; }
[data-testid="stMarkdownContainer"] strong em,
[data-testid="stMarkdownContainer"] em strong { font-weight: 700 !important; font-style: italic !important; }

#MainMenu, footer, header { visibility: hidden; }
.block-container { padding: 0 2.5rem 4rem !important; max-width: 1440px !important; }

::-webkit-scrollbar { width: 5px; height: 5px; }
::-webkit-scrollbar-track { background: transparent; }
::-webkit-scrollbar-thumb { background: var(--border-hi); border-radius: 4px; }
::-webkit-scrollbar-thumb:hover { background: var(--muted); }

/* TOP NAV */
.topbar {
    background: var(--surface);
    border-bottom: 1px solid var(--border);
    padding: 0.8rem 2.5rem;
    margin: 0 -2.5rem 1.75rem;
    display: flex;
    align-items: center;
    gap: 1.25rem;
    position: relative;
}
.topbar::after {
    content: '';
    position: absolute;
    bottom: -1px; left: 0; right: 0; height: 1px;
    background: linear-gradient(90deg, var(--blue) 0%, rgba(64,128,255,0) 55%);
}
.topbar-logo { font-family: 'Inter Tight', sans-serif; font-size: 0.95rem; font-weight: 900; color: var(--text) !important; letter-spacing: 0.12em; text-transform: uppercase; }
.topbar-logo .seg { color: var(--blue); }
.topbar-divider { width: 1px; height: 18px; background: var(--border-hi); flex-shrink: 0; }
.topbar-label { font-size: 0.6rem; font-weight: 600; color: var(--muted) !important; letter-spacing: 0.2em; text-transform: uppercase; }
.topbar-pill { margin-left: auto; background: var(--blue-glow); border: 1px solid rgba(64,128,255,0.28); border-radius: 20px; padding: 0.18rem 0.7rem; font-size: 0.58rem; font-weight: 700; letter-spacing: 0.14em; text-transform: uppercase; color: var(--blue) !important; }

/* HERO */
.hero { padding: 1.5rem 0 0.75rem; }
.hero-title { font-family: 'Inter Tight', sans-serif; font-size: 2.4rem; font-weight: 900; line-height: 1.05; color: var(--text) !important; letter-spacing: -0.03em; margin-bottom: 0.5rem; }
.hero-title .accent { color: var(--blue); }
.hero-sub { font-size: 0.87rem; font-weight: 300; color: var(--muted) !important; max-width: 520px; line-height: 1.65; }

/* UPLOAD ZONE */
.upload-block {
    background: var(--surface);
    border: 1px solid var(--border);
    border-top: 2px solid var(--blue);
    border-radius: 0 0 10px 10px;
    padding: 1.4rem 1.75rem 1.25rem;
    margin: 1.25rem 0 0;
}
.field-label { font-size: 0.58rem; font-weight: 700; letter-spacing: 0.22em; text-transform: uppercase; color: var(--muted) !important; margin-bottom: 0.3rem; }

/* FILE BADGE */
.file-badge {
    display: inline-flex; align-items: center; gap: 0.5rem;
    background: var(--surface2); border: 1px solid var(--border-hi);
    border-radius: 6px; padding: 0.4rem 0.85rem;
    font-size: 0.78rem; color: var(--text-dim) !important; margin-top: 0.3rem;
}
.file-badge .ext { font-family: 'Inter Tight', sans-serif; font-weight: 800; font-size: 0.7rem; color: var(--blue) !important; letter-spacing: 0.08em; }

/* FORM CONTROLS */
.stTextArea textarea {
    background: var(--bg) !important;
    border: 1px solid var(--border) !important;
    border-radius: 6px !important;
    color: var(--text) !important;
    font-family: 'Poppins', sans-serif !important;
    font-size: 0.88rem !important;
    caret-color: var(--blue) !important;
}
.stTextArea textarea:focus {
    border-color: var(--blue) !important;
    box-shadow: 0 0 0 3px var(--blue-glow) !important;
}
textarea::placeholder { color: var(--muted) !important; opacity: 0.6 !important; }

div[data-baseweb="select"] > div,
div[data-baseweb="select"] > div > div {
    background: var(--bg) !important;
    border-color: var(--border) !important;
    color: var(--text) !important;
}
div[data-baseweb="select"] svg { fill: var(--muted) !important; }
div[data-baseweb="select"] span,
div[data-baseweb="select"] input { color: var(--text) !important; }
div[data-baseweb="menu"],
div[data-baseweb="popover"] { background: var(--surface2) !important; border: 1px solid var(--border-hi) !important; box-shadow: 0 8px 32px rgba(0,0,0,0.5) !important; }
div[data-baseweb="menu"] li,
div[data-baseweb="menu"] [role="option"] { color: var(--text) !important; background: transparent !important; }
div[data-baseweb="menu"] li:hover,
div[data-baseweb="menu"] [aria-selected="true"] { background: var(--surface3) !important; color: var(--text) !important; }

/* BUTTONS */
.stButton > button {
    background: var(--blue) !important; color: #fff !important; border: none !important;
    border-radius: 6px !important; font-family: 'Inter Tight', sans-serif !important;
    font-size: 0.78rem !important; font-weight: 800 !important; letter-spacing: 0.12em !important;
    text-transform: uppercase !important; padding: 0.5rem 1.5rem !important;
    transition: background 0.15s, box-shadow 0.15s, transform 0.1s !important;
    box-shadow: 0 2px 10px rgba(64,128,255,0.3) !important;
}
.stButton > button:hover { background: #2d6aee !important; box-shadow: 0 4px 18px rgba(64,128,255,0.45) !important; transform: translateY(-1px) !important; }
.stButton > button:active { transform: translateY(0px) !important; }
.stButton > button:disabled { background: var(--surface3) !important; color: var(--muted) !important; box-shadow: none !important; transform: none !important; }

.stDownloadButton > button {
    background: transparent !important; color: var(--blue) !important;
    border: 1px solid rgba(64,128,255,0.35) !important; border-radius: 6px !important;
    font-family: 'Inter Tight', sans-serif !important; font-size: 0.72rem !important;
    font-weight: 700 !important; letter-spacing: 0.1em !important; text-transform: uppercase !important;
    transition: all 0.15s !important; box-shadow: none !important;
}
.stDownloadButton > button:hover { background: var(--blue-glow) !important; border-color: var(--blue) !important; transform: none !important; }

/* SECTION HEADER */
.section-header {
    display: flex; align-items: center; gap: 0.55rem;
    font-family: 'Inter Tight', sans-serif; font-size: 0.72rem; font-weight: 800;
    letter-spacing: 0.18em; text-transform: uppercase; color: var(--text-dim) !important;
    margin: 1.6rem 0 0.9rem;
}
.section-header .dot { width: 6px; height: 6px; border-radius: 50%; background: var(--blue); flex-shrink: 0; }

/* RESULT CARD */
.result-card {
    background: var(--surface);
    border: 1px solid var(--border);
    border-radius: 10px;
    padding: 1.5rem 1.75rem;
    margin-top: 1.25rem;
    line-height: 1.75;
}
.result-card h2, .result-card h3 {
    font-family: 'Inter Tight', sans-serif !important;
    font-weight: 800 !important;
    letter-spacing: -0.01em !important;
    color: var(--text) !important;
}

/* PROGRESS */
.stProgress > div > div { background: var(--blue) !important; }

/* TABS */
button[data-baseweb="tab"] {
    font-family: 'Inter Tight', sans-serif !important;
    font-size: 0.72rem !important; font-weight: 800 !important;
    letter-spacing: 0.14em !important; text-transform: uppercase !important;
    color: var(--muted) !important; background: transparent !important;
}
button[data-baseweb="tab"][aria-selected="true"] { color: var(--blue) !important; }
div[data-baseweb="tab-highlight"] { background: var(--blue) !important; }
div[data-baseweb="tab-border"] { background: var(--border) !important; }

/* FILE UPLOADER */
[data-testid="stFileUploader"] {
    background: var(--surface2) !important;
    border: 1px dashed var(--border-hi) !important;
    border-radius: 8px !important;
    padding: 0.5rem !important;
}
[data-testid="stFileUploader"]:hover {
    border-color: var(--blue) !important;
}
[data-testid="stFileUploaderDropzone"] {
    background: transparent !important;
}
[data-testid="stFileUploaderDropzoneInstructions"] span,
[data-testid="stFileUploaderDropzoneInstructions"] p {
    color: var(--muted) !important;
}

/* EMPTY STATE */
.empty-state { text-align: center; padding: 5rem 2rem; }
.empty-title { font-family: 'Inter Tight', sans-serif; font-size: 1.6rem; font-weight: 900; letter-spacing: -0.02em; color: var(--border-hi) !important; margin-bottom: 0.75rem; }
.empty-sub { font-size: 0.88rem; color: var(--muted) !important; max-width: 400px; margin: 0 auto; line-height: 1.65; }

/* FOOTER */
.footer { margin-top: 4rem; padding: 1.5rem 0; border-top: 1px solid var(--border); display: flex; align-items: center; justify-content: space-between; }
.footer-brand { font-family: 'Inter Tight', sans-serif; font-size: 0.72rem; font-weight: 900; letter-spacing: 0.18em; color: var(--muted) !important; }
.footer-note { font-size: 0.65rem; color: var(--muted) !important; }
</style>
""", unsafe_allow_html=True)

# ─────────────────────────────────────────────────────────────
# DEPENDENCY CHECKS
# ─────────────────────────────────────────────────────────────

try:
    import anthropic as _anthropic
    ANTHROPIC_AVAILABLE = True
except ImportError:
    ANTHROPIC_AVAILABLE = False

try:
    import docx as _docx
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    import openpyxl as _openpyxl
    XLSX_AVAILABLE = True
except ImportError:
    XLSX_AVAILABLE = False

try:
    import fitz as _fitz  # PyMuPDF
    PDF_AVAILABLE = True
except ImportError:
    PDF_AVAILABLE = False

# ─────────────────────────────────────────────────────────────
# TEXT EXTRACTION
# ─────────────────────────────────────────────────────────────

def extract_text(uploaded_file) -> tuple[str, str]:
    """
    Returns (text_content, file_type) from an uploaded Streamlit file.
    file_type is one of: 'pdf', 'docx', 'xlsx', 'csv', 'pptx', 'unknown'
    """
    name = uploaded_file.name.lower()
    raw  = uploaded_file.read()

    if name.endswith(".pdf"):
        if not PDF_AVAILABLE:
            return "[ERROR: PyMuPDF not installed — run: pip install pymupdf]", "pdf"
        try:
            doc  = _fitz.open(stream=raw, filetype="pdf")
            text = "\n\n".join(page.get_text() for page in doc)
            doc.close()
            return text.strip() or "[PDF contained no extractable text]", "pdf"
        except Exception as e:
            return f"[PDF extraction error: {e}]", "pdf"

    elif name.endswith(".docx"):
        if not DOCX_AVAILABLE:
            return "[ERROR: python-docx not installed — run: pip install python-docx]", "docx"
        try:
            doc  = _docx.Document(io.BytesIO(raw))
            text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
            return text.strip() or "[DOCX contained no extractable text]", "docx"
        except Exception as e:
            return f"[DOCX extraction error: {e}]", "docx"

    elif name.endswith((".xlsx", ".xls", ".xlsm")):
        if not XLSX_AVAILABLE:
            return "[ERROR: openpyxl not installed — run: pip install openpyxl]", "xlsx"
        try:
            wb    = _openpyxl.load_workbook(io.BytesIO(raw), read_only=True, data_only=True, keep_vba=True)
            lines = []
            for sheet in wb.worksheets:
                lines.append(f"\n[Sheet: {sheet.title}]")
                for row in sheet.iter_rows():
                    for cell in row:
                        if cell.value is not None and str(cell.value).strip():
                            coord = f"{sheet.title}!{cell.coordinate}"
                            lines.append(f"{coord}: {cell.value}")
            wb.close()
            return "\n".join(lines).strip() or "[XLSX contained no extractable text]", "xlsx"
        except Exception as e:
            return f"[XLSX extraction error: {e}]", "xlsx"

    elif name.endswith(".csv"):
        try:
            import csv as _csv
            text_raw = raw.decode("utf-8-sig", errors="replace")
            reader   = _csv.reader(io.StringIO(text_raw))
            rows     = list(reader)
            if not rows:
                return "[CSV contained no data]", "csv"
            lines = []
            headers = rows[0]
            for ri, row in enumerate(rows):
                # Label each cell as Row,Col so Claude can cite exact locations
                for ci, val in enumerate(row):
                    if val.strip():
                        col_label = headers[ci] if ri > 0 and ci < len(headers) else f"Col{ci+1}"
                        lines.append(f"Row{ri+1},{col_label}: {val}")
            return "\n".join(lines).strip() or "[CSV contained no extractable text]", "csv"
        except Exception as e:
            return f"[CSV extraction error: {e}]", "csv"

    elif name.endswith((".pptx", ".ppt")):
        try:
            from pptx import Presentation as _Presentation
            prs   = _Presentation(io.BytesIO(raw))
            lines = []
            for si, slide in enumerate(prs.slides, start=1):
                slide_texts = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            text = para.text.strip()
                            if text:
                                slide_texts.append(text)
                if slide_texts:
                    lines.append(f"[Slide {si}]")
                    lines.extend(slide_texts)
            return "\n".join(lines).strip() or "[PPTX contained no extractable text]", "pptx"
        except ImportError:
            return "[ERROR: python-pptx not installed — run: pip install python-pptx]", "pptx"
        except Exception as e:
            return f"[PPTX extraction error: {e}]", "pptx"

    else:
        return "[Unsupported file type]", "unknown"


def truncate(text: str, max_chars: int = 40_000) -> str:
    if len(text) <= max_chars:
        return text
    half = max_chars // 2
    return text[:half] + f"\n\n[… {len(text) - max_chars:,} characters omitted for brevity …]\n\n" + text[-half:]


# ─────────────────────────────────────────────────────────────
# SESSION STATE
# ─────────────────────────────────────────────────────────────

for key, default in [
    ("comparison_result", ""),
    ("chat_history",      []),
    ("chat_pending",      False),
    ("doc_a_text",        ""),
    ("doc_b_text",        ""),
    ("doc_a_name",        ""),
    ("doc_b_name",        ""),
    ("last_focus",        ""),
    ("last_tone",         ""),
]:
    if key not in st.session_state:
        st.session_state[key] = default

# ─────────────────────────────────────────────────────────────
# API KEY
# ─────────────────────────────────────────────────────────────

claude_key = st.secrets.get("ANTHROPIC_API_KEY", "")

# ─────────────────────────────────────────────────────────────
# OTP AUTHENTICATION
# ─────────────────────────────────────────────────────────────

ALLOWED_DOMAIN   = "@segaamerica.com"
OTP_EXPIRY_SECS  = 600   # 10 minutes
COOKIE_EXPIRY_DAYS = 1
COOKIE_NAME      = "sega_doc_auth"

def _send_otp(email: str, code: str) -> bool:
    """Send OTP via AWS SES. Returns True on success."""
    try:
        import boto3
        from botocore.exceptions import ClientError

        ses = boto3.client(
            "ses",
            region_name=st.secrets.get("AWS_SES_REGION", "us-east-1"),
            aws_access_key_id=st.secrets.get("AWS_ACCESS_KEY_ID", ""),
            aws_secret_access_key=st.secrets.get("AWS_SECRET_ACCESS_KEY", ""),
        )

        ses.send_email(
            Source=st.secrets.get("EMAIL_FROM", "noreply@segaamerica.com"),
            Destination={"ToAddresses": [email]},
            Message={
                "Subject": {
                    "Data": "SEGA Doc Comparator — Your verification code",
                    "Charset": "UTF-8",
                },
                "Body": {
                    "Text": {
                        "Data": f"Your SEGA Doc Comparator verification code is: {code}\n\nThis code expires in 10 minutes.\nIf you didn't request this, you can safely ignore this email.",
                        "Charset": "UTF-8",
                    },
                    "Html": {
                        "Data": f"""
                        <div style="font-family:Arial,sans-serif;max-width:480px;margin:0 auto;padding:32px 24px;">
                          <div style="font-size:22px;font-weight:900;letter-spacing:0.1em;color:#1A6BFF;margin-bottom:4px;">SEGA</div>
                          <div style="font-size:14px;color:#444;margin-bottom:28px;">Doc Comparator</div>
                          <div style="font-size:14px;color:#222;margin-bottom:16px;">Your verification code is:</div>
                          <div style="font-size:42px;font-weight:900;letter-spacing:0.18em;color:#1a1a2e;
                                      background:#f0f4ff;border-radius:8px;padding:18px 24px;
                                      display:inline-block;margin-bottom:24px;">{code}</div>
                          <div style="font-size:12px;color:#888;">
                            This code expires in 10 minutes.<br>
                            If you didn't request this, you can safely ignore this email.
                          </div>
                        </div>
                        """,
                        "Charset": "UTF-8",
                    },
                },
            },
        )
        return True
    except ClientError as e:
        st.error(f"SES error: {e.response['Error']['Message']}")
        return False
    except Exception as e:
        st.error(f"Failed to send email: {e}")
        return False

def _sign_cookie(email: str) -> str:
    """Create an HMAC-signed token: email|expiry|signature."""
    secret = st.secrets.get("COOKIE_SIGNING_KEY", "fallback-change-this")
    expiry = int(time.time()) + (COOKIE_EXPIRY_DAYS * 86400)
    payload = f"{email}|{expiry}"
    sig = hmac.new(secret.encode(), payload.encode(), hashlib.sha256).hexdigest()
    return base64.urlsafe_b64encode(f"{payload}|{sig}".encode()).decode()

def _verify_cookie(token: str) -> str | None:
    """Verify signed cookie. Returns email if valid, None otherwise."""
    try:
        secret = st.secrets.get("COOKIE_SIGNING_KEY", "fallback-change-this")
        decoded = base64.urlsafe_b64decode(token.encode()).decode()
        email, expiry_str, sig = decoded.rsplit("|", 2)
        payload = f"{email}|{expiry_str}"
        expected = hmac.new(secret.encode(), payload.encode(), hashlib.sha256).hexdigest()
        if not hmac.compare_digest(sig, expected):
            return None
        if int(time.time()) > int(expiry_str):
            return None
        return email
    except Exception:
        return None

# ── Check for existing valid cookie ──────────────────────────
try:
    from streamlit_cookies_controller import CookieController
    _cookie_manager = CookieController(key="auth_cookie_ctrl")
    _existing_cookie = _cookie_manager.get(COOKIE_NAME)
except Exception:
    _cookie_manager = None
    _existing_cookie = None

_cookie_email = _verify_cookie(_existing_cookie) if _existing_cookie else None

# ── Auth state init ───────────────────────────────────────────
for _k, _v in [
    ("auth_verified",   False),
    ("auth_email",      ""),
    ("otp_code",        ""),
    ("otp_email",       ""),
    ("otp_expiry",      0),
    ("otp_sent",        False),
    ("otp_attempts",    0),
]:
    if _k not in st.session_state:
        st.session_state[_k] = _v

# If valid cookie found, mark as verified
if _cookie_email and not st.session_state.auth_verified:
    st.session_state.auth_verified = True
    st.session_state.auth_email    = _cookie_email

# ── Render login gate if not verified ────────────────────────
if not st.session_state.auth_verified:
    st.markdown("""
    <style>
    .auth-wrap {
        max-width: 420px; margin: 5rem auto; padding: 2.5rem 2.5rem 2rem;
        background: var(--surface); border: 1px solid var(--border);
        border-top: 3px solid var(--blue); border-radius: 0 0 10px 10px;
    }
    .auth-logo { font-family:'Inter Tight',sans-serif; font-size:1.6rem; font-weight:900;
                 letter-spacing:0.12em; color:var(--blue) !important; margin-bottom:0.2rem; }
    .auth-title { font-family:'Inter Tight',sans-serif; font-size:1rem; font-weight:700;
                  color:var(--text) !important; margin-bottom:0.25rem; }
    .auth-sub { font-size:0.8rem; color:var(--muted) !important; margin-bottom:1.5rem; }
    .auth-note { font-size:0.72rem; color:var(--muted) !important; margin-top:1rem;
                 text-align:center; line-height:1.5; }
    </style>
    """, unsafe_allow_html=True)

    _lc, _mc, _rc = st.columns([1, 2, 1])
    with _mc:
        st.markdown("""
        <div class="auth-wrap">
          <div class="auth-logo">SEGA</div>
          <div class="auth-title">Doc Comparator</div>
          <div class="auth-sub">Sign in with your SEGA America email</div>
        </div>
        """, unsafe_allow_html=True)

        if not st.session_state.otp_sent:
            _email_input = st.text_input(
                "Email address",
                placeholder="you@segaamerica.com",
                label_visibility="collapsed",
                key="auth_email_input",
            )
            _send_btn = st.button("Send verification code", width="stretch")

            if _send_btn and _email_input:
                if not _email_input.strip().lower().endswith(ALLOWED_DOMAIN):
                    st.error(f"Access restricted to {ALLOWED_DOMAIN} addresses.")
                else:
                    _code = str(random.randint(100000, 999999))
                    if _send_otp(_email_input.strip().lower(), _code):
                        st.session_state.otp_code     = _code
                        st.session_state.otp_email    = _email_input.strip().lower()
                        st.session_state.otp_expiry   = time.time() + OTP_EXPIRY_SECS
                        st.session_state.otp_sent     = True
                        st.session_state.otp_attempts = 0
                        st.rerun()

        else:
            st.info(f"Code sent to **{st.session_state.otp_email}** — check your inbox.")
            _code_input = st.text_input(
                "6-digit code",
                placeholder="123456",
                label_visibility="collapsed",
                max_chars=6,
                key="auth_code_input",
            )
            _verify_btn = st.button("Verify code", width="stretch")

            if _verify_btn and _code_input:
                if st.session_state.otp_attempts >= 5:
                    st.error("Too many attempts. Please request a new code.")
                    st.session_state.otp_sent = False
                elif time.time() > st.session_state.otp_expiry:
                    st.error("Code has expired. Please request a new one.")
                    st.session_state.otp_sent = False
                elif _code_input.strip() != st.session_state.otp_code:
                    st.session_state.otp_attempts += 1
                    _remaining = 5 - st.session_state.otp_attempts
                    st.error(f"Incorrect code. {_remaining} attempt{'s' if _remaining != 1 else ''} remaining.")
                else:
                    # Success
                    st.session_state.auth_verified = True
                    st.session_state.auth_email    = st.session_state.otp_email
                    st.session_state.otp_code      = ""  # clear code from state
                    # Set persistent cookie
                    if _cookie_manager:
                        _token = _sign_cookie(st.session_state.auth_email)
                        _cookie_manager.set(COOKIE_NAME, _token)
                    st.rerun()

            _resend_col, _ = st.columns([1, 1])
            with _resend_col:
                if st.button("← Use a different email", key="auth_back"):
                    st.session_state.otp_sent  = False
                    st.session_state.otp_code  = ""
                    st.rerun()

        st.markdown(
            f'<div class="auth-note">Restricted to {ALLOWED_DOMAIN} addresses only.<br>'
            f'Codes expire after 10 minutes.</div>',
            unsafe_allow_html=True
        )

    st.stop()  # Nothing below renders until authenticated

# ─────────────────────────────────────────────────────────────
# SIGNED-IN USER + SIGN OUT
# ─────────────────────────────────────────────────────────────

with st.sidebar:
    st.markdown(
        f'<div style="font-size:.7rem;font-weight:600;color:var(--muted);margin-bottom:.5rem;">'
        f'Signed in as<br>'
        f'<span style="color:var(--text);font-weight:700;">{st.session_state.auth_email}</span>'
        f'</div>',
        unsafe_allow_html=True
    )
    if st.button("Sign out", key="sign_out_btn"):
        if _cookie_manager:
            _cookie_manager.remove(COOKIE_NAME)
        for _k in ["auth_verified", "auth_email", "otp_sent", "otp_code",
                   "otp_email", "otp_expiry", "otp_attempts"]:
            st.session_state[_k] = False if _k == "auth_verified" else ""
        st.rerun()

st.markdown("""
<div class="topbar">
  <div class="topbar-logo"><span class="seg">SEGA</span> DOC COMPARATOR</div>
  <div class="topbar-divider"></div>
  <div class="topbar-label">Document Intelligence Platform</div>
  <div class="topbar-pill">Beta</div>
</div>
""", unsafe_allow_html=True)

# ─────────────────────────────────────────────────────────────
# HERO
# ─────────────────────────────────────────────────────────────

st.markdown("""
<div class="hero">
  <div class="hero-title">DOCUMENT <span class="accent">COMPARE</span></div>
  <div class="hero-sub">Upload two documents — PDFs, Word files, or spreadsheets — and get a detailed AI-powered comparison across any dimension you choose.</div>
</div>
""", unsafe_allow_html=True)

# ─────────────────────────────────────────────────────────────
# API KEY STATUS
# ─────────────────────────────────────────────────────────────

if not claude_key:
    st.markdown("""
<div style="background:var(--surface);border:1px solid var(--border);border-top:2px solid var(--neg);
border-radius:0 0 8px 8px;padding:1rem 1.5rem;margin-bottom:1rem;">
<div style="font-size:.62rem;font-weight:700;letter-spacing:.18em;text-transform:uppercase;
color:var(--neg);margin-bottom:.5rem;">Missing API Key</div>
<div style="font-size:.82rem;color:var(--text-dim);line-height:1.7;">
Add the following to <code>.streamlit/secrets.toml</code> in your project folder:
<pre style="background:var(--bg);border:1px solid var(--border);border-radius:6px;
padding:.75rem 1rem;margin:.6rem 0 0;font-size:.8rem;color:var(--blue);">ANTHROPIC_API_KEY = "sk-ant-your-key-here"</pre>
Then restart the app.
</div>
</div>
""", unsafe_allow_html=True)
elif not ANTHROPIC_AVAILABLE:
    st.markdown("""
<div style="background:var(--surface);border:1px solid var(--border);border-top:2px solid var(--neg);
border-radius:0 0 8px 8px;padding:1rem 1.5rem;margin-bottom:1rem;">
<div style="font-size:.62rem;font-weight:700;letter-spacing:.18em;text-transform:uppercase;
color:var(--neg);margin-bottom:.5rem;">Missing Dependency</div>
<div style="font-size:.82rem;color:var(--text-dim);">Run: <code>pip install anthropic</code></div>
</div>
""", unsafe_allow_html=True)
else:
    st.markdown("""
<div style="display:flex;gap:1rem;margin-bottom:.75rem;flex-wrap:wrap;">
  <div style="background:var(--pos-dim);border:1px solid rgba(32,198,90,.25);border-radius:6px;
  padding:.4rem 1rem;font-size:.68rem;font-weight:600;color:var(--pos);">✓ Anthropic API connected</div>
</div>
""", unsafe_allow_html=True)

# ─────────────────────────────────────────────────────────────
# UPLOAD BLOCK
# ─────────────────────────────────────────────────────────────

st.markdown('<div class="upload-block">', unsafe_allow_html=True)

col_a, col_sep, col_b = st.columns([5, 0.3, 5])

with col_a:
    st.markdown('<div class="field-label">Document A</div>', unsafe_allow_html=True)
    file_a = st.file_uploader(
        "doc_a", label_visibility="collapsed",
        type=["pdf", "docx", "xlsx", "xlsm", "csv", "pptx"],
    )
    if file_a:
        ext_a = file_a.name.rsplit(".", 1)[-1].upper()
        st.markdown(
            f'<div class="file-badge"><span class="ext">{ext_a}</span>{file_a.name}</div>',
            unsafe_allow_html=True,
        )

with col_sep:
    st.markdown('<div style="display:flex;align-items:center;justify-content:center;height:100%;padding-top:1.5rem;font-size:1.4rem;color:var(--muted);">⇄</div>', unsafe_allow_html=True)

with col_b:
    st.markdown('<div class="field-label">Document B</div>', unsafe_allow_html=True)
    file_b = st.file_uploader(
        "doc_b", label_visibility="collapsed",
        type=["pdf", "docx", "xlsx", "xlsm", "csv", "pptx"],
    )
    if file_b:
        ext_b = file_b.name.rsplit(".", 1)[-1].upper()
        st.markdown(
            f'<div class="file-badge"><span class="ext">{ext_b}</span>{file_b.name}</div>',
            unsafe_allow_html=True,
        )

st.markdown("<br>", unsafe_allow_html=True)

# ── Options layout ────────────────────────────────────────────
opt1, opt2, opt3 = st.columns([3, 2, 1.5])

with opt1:
    st.markdown('<div class="field-label">Comparison Focus</div>', unsafe_allow_html=True)
    focus_options = [
        "Comprehensive Assessment",
        "Key Differences & Similarities",
        "Risk & Compliance",
        "Data & Numbers",
        "Custom (describe below)",
    ]
    _focus_descs = {
        "Comprehensive Assessment":       "Comprehensive Assessment — Full structured analysis across all dimensions",
        "Key Differences & Similarities": "Key Differences & Similarities — What changed, what stayed the same",
        "Risk & Compliance":              "Risk & Compliance — Obligations, gaps, and legal exposure",
        "Data & Numbers":                 "Data & Numbers — Every figure, metric, and statistic compared",
        "Custom (describe below)":        "Custom — Define your own comparison focus",
    }
    focus = st.selectbox(
        "focus", focus_options, index=0,
        format_func=lambda x: _focus_descs[x],
        label_visibility="collapsed",
    )

with opt2:
    st.markdown('<div class="field-label">Output Tone</div>', unsafe_allow_html=True)
    tone_options = [
        "Analytical & detailed",
        "Executive brief",
        "Plain language",
        "Legal / formal",
    ]
    _tone_descs = {
        "Analytical & detailed": "Analytical & detailed — Thorough, evidence-backed, fully structured",
        "Executive brief":       "Executive brief — Lead with the headline, 400–600 words max",
        "Plain language":        "Plain language — Clear and jargon-free, for any audience",
        "Legal / formal":        "Legal / formal — Precise, numbered sections, flags obligations",
    }
    tone = st.selectbox(
        "tone", tone_options, index=0,
        format_func=lambda x: _tone_descs[x],
        label_visibility="collapsed",
    )

with opt3:
    st.markdown('<div class="field-label">&nbsp;</div>', unsafe_allow_html=True)
    compare_clicked = st.button(
        "COMPARE DOCUMENTS",
        disabled=not (file_a and file_b and claude_key and ANTHROPIC_AVAILABLE),
        width="stretch",
    )

model = "claude-sonnet-4-6"

# Custom focus text area (conditional)
custom_focus = ""
if focus == "Custom (describe below)":
    st.markdown('<div class="field-label" style="margin-top:.75rem;">Describe your comparison focus</div>', unsafe_allow_html=True)
    custom_focus = st.text_area(
        "custom_focus", label_visibility="collapsed",
        placeholder="e.g. Compare the pricing structures and payment terms in both documents...",
        height=80,
    )

st.markdown("</div>", unsafe_allow_html=True)  # close upload-block

# ─────────────────────────────────────────────────────────────
# MARKDOWN RENDER HELPER
# ─────────────────────────────────────────────────────────────

import re as _re

def _escape_dollars(text: str) -> str:
    """Escape bare $ signs so Streamlit doesn't treat them as LaTeX delimiters."""
    # Replace $ not already escaped, not inside code blocks
    return _re.sub(r'(?<!\\)\$', r'\\$', text)

# ─────────────────────────────────────────────────────────────
# COMPARISON LOGIC
# ─────────────────────────────────────────────────────────────

if compare_clicked and file_a and file_b:
    # Reset state for fresh comparison
    st.session_state.comparison_result = ""
    st.session_state.chat_history      = []
    st.session_state.chat_pending      = False

    with st.spinner("Extracting document contents…"):
        file_a.seek(0)
        text_a, type_a = extract_text(file_a)
        file_b.seek(0)
        text_b, type_b = extract_text(file_b)

    st.session_state.doc_a_text = text_a
    st.session_state.doc_b_text = text_b
    st.session_state.doc_a_name = file_a.name
    st.session_state.doc_b_name = file_b.name
    st.session_state.last_focus = focus
    st.session_state.last_tone  = tone

    # Build prompt
    actual_focus = custom_focus.strip() if focus == "Custom (describe below)" and custom_focus.strip() else focus

    focus_instructions = {
        "Comprehensive Assessment": """Deliver a full, structured comparison across every meaningful dimension:

## 1. Most Significant Differences
Identify and explain the 4-6 most impactful differences between the documents. For each: state what changed, quote the specific language from both documents, and explain the significance.

## 2. Numerical & Data Comparison
List every figure, metric, date, percentage, and statistic that appears in both documents. Present them in a table (Metric | Document A | Document B | Delta). Flag any discrepancies or contradictions.

## 3. Comprehensiveness Assessment
Which document is more complete? For each major topic area, state whether Document A, Document B, or both cover it — and note anything present in one that is entirely absent from the other.

## 4. Tone & Framing
How does the language and framing differ? Note any tonal shifts, changes in confidence or hedging, or differences in how key facts are presented.

## 5. Agreements & Alignment
Where do the documents agree, reinforce, or directly echo each other? Quote matching language where relevant.

## 6. Summary Verdict
One paragraph: which document is more authoritative / complete / reliable, and why? What is the single most important thing a reader should know about the difference between these two documents?

Be exhaustive. Quote directly. Use tables for any comparative data.""",

        "Key Differences & Similarities": """Compare the two documents thoroughly:
1. List the 5-7 most significant differences — be specific, cite sections or values
2. List 3-5 meaningful similarities or areas of alignment
3. Which document is more comprehensive, and in what respects?
4. Are there any contradictions or conflicts between the documents?
5. What is the most critical difference a decision-maker should know about?""",

        "Risk & Compliance": """Analyse both documents from a risk and compliance perspective:
1. Identify all risk-related clauses, obligations, or statements in each document
2. Compare the risk profiles — which document carries more risk and for whom?
3. Highlight any compliance gaps, missing protections, or ambiguous language
4. Flag any terms that could be problematic legally or operationally
5. What are the top 3 risk differences between the documents?""",

        "Data & Numbers": """Extract and compare all numerical content:
1. List all key figures, statistics, dates, percentages and metrics from each document
2. Where the same metric appears in both, compare the values — are they consistent?
3. Identify discrepancies or contradictions in numerical data
4. Which document provides more quantitative evidence, and is it more credible?
5. Summarise the financial or quantitative picture each document paints""",
    }

    focus_text = focus_instructions.get(actual_focus, f"""Compare the documents with this specific focus: {actual_focus}

Provide a thorough, structured analysis addressing the user's stated focus. Use concrete evidence from the documents — quote specific language, cite figures, and name sections. Be direct and analytical.""")

    tone_map = {
        "Analytical & detailed":
            "Write in a precise, analytical tone. Use structured headers and sub-points. Support every claim with specific evidence from the documents. Be thorough.",
        "Executive brief":
            "Write as a concise executive briefing. Lead with the single most important finding. Use headers and bullets. Total length: 400-600 words max.",
        "Plain language":
            "Write in plain, accessible language. Avoid jargon. Explain any technical terms. Aim for clarity over comprehensiveness.",
        "Legal / formal":
            "Write in formal, precise language suitable for legal or contractual review. Use structured numbered sections. Flag ambiguities and obligations explicitly.",
    }

    _xlsx_rule = ""
    if type_a == "xlsx" or type_b == "xlsx":
        _xlsx_rule += "\n- EXCEL FILES: For every difference found in a spreadsheet, state the exact sheet name and cell reference (e.g. Sheet1!B4, Financial Model!C12). Do not describe a change without citing its precise location."
    if type_a == "csv" or type_b == "csv":
        _xlsx_rule += "\n- CSV FILES: For every difference found, cite the exact row number and column header (e.g. Row12,Revenue or Row3,Date). Do not describe a change without citing its precise location."
    if type_a == "pptx" or type_b == "pptx":
        _xlsx_rule += "\n- POWERPOINT FILES: Compare text content only — slide titles, body text, and speaker notes. Cite the slide number for every difference (e.g. Slide 4). Ignore all styling, formatting, colours, fonts, and layout differences entirely."

    prompt = f"""You are a senior document analyst. You will compare two documents in depth.

═══════════════════════════════════
DOCUMENT A: {file_a.name}
═══════════════════════════════════
{truncate(text_a)}

═══════════════════════════════════
DOCUMENT B: {file_b.name}
═══════════════════════════════════
{truncate(text_b)}

═══════════════════════════════════
YOUR TASK
═══════════════════════════════════
{focus_text}

OUTPUT TONE: {tone_map.get(tone, tone_map["Analytical & detailed"])}

HARD RULES:
- Every observation must be grounded in specific content from the documents
- Quote directly from the documents when making claims about language or tone
- Use clear markdown headers (##) for sections
- Do not add generic preamble or sign-off — go straight into the analysis
- Do NOT offer suggestions, recommendations, or possible explanations for why a change was made — only state what the documents say and how they differ{_xlsx_rule}"""

    # Stream the response
    st.markdown('<div class="section-header"><span class="dot"></span>COMPARISON RESULT</div>', unsafe_allow_html=True)
    result_placeholder = st.empty()
    full_text = ""

    try:
        client = _anthropic.Anthropic(api_key=claude_key)
        with client.messages.stream(
            model=model,
            max_tokens=4096,
            system=(
                "You are a senior document analyst. "
                "Respond only with your analysis in well-structured markdown. "
                "No preamble, no sign-off. "
                "IMPORTANT FORMATTING RULES: "
                "1. Use **bold** only for standalone words or short phrases — never wrap text containing $, numbers, or punctuation in * or ** as this breaks rendering. "
                "2. When quoting values or figures (e.g. $49.6M), write them as plain text, not inside asterisks. "
                "3. For inline quotes from documents use double quotes (\") not asterisk-italic. "
                "4. Use ## for section headers and ### for sub-headers. "
                "5. Use markdown tables for any comparative data. "
                "6. ALWAYS begin your response with a summary differences table before any other content. "
                "   The table must have columns: Item | Document A | Document B | Delta. "
                "   Include every meaningful difference found. This table must be the very first thing in your response, with no heading or preamble before it. "
                "CRITICAL: Do NOT offer suggestions, recommendations, or possible explanations for why a change was made. Only report what each document states and how they differ. Never say 'this may be because', 'you might consider', 'this could indicate', or similar speculative or advisory language."
            ),
            messages=[{"role": "user", "content": prompt}],
        ) as stream:
            for delta in stream.text_stream:
                full_text += delta
                result_placeholder.markdown(_escape_dollars(full_text) + "▌")

        result_placeholder.markdown(_escape_dollars(full_text))
        st.session_state.comparison_result = full_text

    except _anthropic.AuthenticationError:
        st.error("Invalid Anthropic API key — check your secrets.toml.")
    except _anthropic.RateLimitError:
        st.error("Rate limit reached. Wait a moment and try again.")
    except _anthropic.APIConnectionError as e:
        st.error(f"Could not reach the Anthropic API. Check your internet connection.\nDetail: {e}")
    except Exception as e:
        st.error(f"Unexpected error: {type(e).__name__}: {e}")

# ─────────────────────────────────────────────────────────────
# RESULTS PANEL (persisted from previous run)
# ─────────────────────────────────────────────────────────────

elif st.session_state.comparison_result:
    st.markdown('<div class="section-header"><span class="dot"></span>COMPARISON RESULT</div>', unsafe_allow_html=True)
    st.markdown(_escape_dollars(st.session_state.comparison_result))

# ─────────────────────────────────────────────────────────────
# DOWNLOAD + CHAT (shown when result exists)
# ─────────────────────────────────────────────────────────────

if st.session_state.comparison_result:
    st.markdown("<br>", unsafe_allow_html=True)

    # ── Downloads ─────────────────────────────────────────────
    st.markdown(
        '<div style="font-size:.62rem;font-weight:700;letter-spacing:.18em;'
        'text-transform:uppercase;color:var(--muted);margin-bottom:.5rem;">'
        'DOWNLOAD REPORT</div>',
        unsafe_allow_html=True,
    )
    _fname = f"comparison_{st.session_state.doc_a_name[:20]}_{st.session_state.doc_b_name[:20]}".replace(" ", "_")

    def _make_pdf(md_text, doc_a, doc_b):
        try:
            from reportlab.lib.pagesizes import letter
            from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
            from reportlab.lib.units import cm
            from reportlab.lib import colors
            from reportlab.platypus import (SimpleDocTemplate, Paragraph, Spacer,
                                            HRFlowable, Table, TableStyle)
            from reportlab.lib.enums import TA_LEFT, TA_CENTER
            import io as _io, re as _re

            buf = _io.BytesIO()
            doc = SimpleDocTemplate(buf, pagesize=letter,
                leftMargin=2*cm, rightMargin=2*cm,
                topMargin=2*cm, bottomMargin=2*cm)

            styles = getSampleStyleSheet()
            s_h1   = ParagraphStyle("h1", parent=styles["Normal"],
                fontSize=14, fontName="Helvetica-Bold", spaceBefore=14,
                spaceAfter=4, textColor=colors.HexColor("#1e3a8a"))
            s_h2   = ParagraphStyle("h2", parent=styles["Normal"],
                fontSize=12, fontName="Helvetica-Bold", spaceBefore=10,
                spaceAfter=3, textColor=colors.HexColor("#1e40af"))
            s_h3   = ParagraphStyle("h3", parent=styles["Normal"],
                fontSize=11, fontName="Helvetica-Bold", spaceBefore=8,
                spaceAfter=2, textColor=colors.HexColor("#374151"))
            s_body = ParagraphStyle("body", parent=styles["Normal"],
                fontSize=9.5, fontName="Helvetica", leading=14, spaceAfter=5)
            s_bull = ParagraphStyle("bull", parent=s_body,
                leftIndent=14, bulletIndent=4)
            s_note = ParagraphStyle("note", parent=styles["Normal"],
                fontSize=8, fontName="Helvetica-Oblique",
                textColor=colors.HexColor("#6b7280"))

            def inline(text):
                text = _re.sub(r'\*\*(.+?)\*\*', r'<b>\1</b>', text)
                text = _re.sub(r'\*(.+?)\*',     r'<i>\1</i>', text)
                text = _re.sub(r'`(.+?)`', r'<font name="Courier">\1</font>', text)
                return text

            story = []

            # Header banner
            hdr_data = [[
                Paragraph('<b><font color="white" size="14">Document Comparison Report</font></b>', styles["Normal"]),
            ],[
                Paragraph(f'<font color="#cbd5e1" size="9">{doc_a}  \u21c4  {doc_b}</font>', styles["Normal"]),
            ]]
            hdr_tbl = Table(hdr_data, colWidths=["100%"])
            hdr_tbl.setStyle(TableStyle([
                ("BACKGROUND", (0,0), (-1,-1), colors.HexColor("#1e3a8a")),
                ("TOPPADDING",    (0,0), (-1,-1), 10),
                ("BOTTOMPADDING",(0,0), (-1,-1), 10),
                ("LEFTPADDING",  (0,0), (-1,-1), 14),
                ("RIGHTPADDING", (0,0), (-1,-1), 14),
            ]))
            story.append(hdr_tbl)
            story.append(Spacer(1, 14))

            lines = md_text.split("\n")
            i = 0
            while i < len(lines):
                line = lines[i]

                # Table detection
                if "|" in line and i + 1 < len(lines) and "|" in lines[i+1] and "---" in lines[i+1]:
                    tbl_lines = []
                    while i < len(lines) and "|" in lines[i]:
                        tbl_lines.append(lines[i])
                        i += 1
                    rows = []
                    for tl in tbl_lines:
                        if _re.match(r"^\s*\|[-: |]+\|\s*$", tl):
                            continue
                        cells = [c.strip() for c in tl.strip().strip("|").split("|")]
                        rows.append(cells)
                    if rows:
                        max_cols = max(len(r) for r in rows)
                        for r in rows:
                            while len(r) < max_cols: r.append("")
                        col_w = (17 * cm) / max_cols
                        para_rows = []
                        for ri, row in enumerate(rows):
                            para_row = []
                            for ci, cell in enumerate(row):
                                txt = inline(cell)
                                sty = ParagraphStyle("tc",
                                    parent=styles["Normal"],
                                    fontSize=8.5, fontName="Helvetica-Bold" if ri==0 else "Helvetica",
                                    textColor=colors.white if ri==0 else colors.HexColor("#1a1a2e"),
                                    leading=12)
                                para_row.append(Paragraph(txt, sty))
                            para_rows.append(para_row)
                        t = Table(para_rows, colWidths=[col_w]*max_cols)
                        t.setStyle(TableStyle([
                            ("BACKGROUND",    (0,0), (-1,0),  colors.HexColor("#1e3a8a")),
                            ("ROWBACKGROUNDS",(0,1), (-1,-1),
                             [colors.HexColor("#f0f4ff"), colors.white]),
                            ("GRID",          (0,0), (-1,-1), 0.4, colors.HexColor("#c7d2fe")),
                            ("TOPPADDING",    (0,0), (-1,-1), 4),
                            ("BOTTOMPADDING", (0,0), (-1,-1), 4),
                            ("LEFTPADDING",   (0,0), (-1,-1), 6),
                            ("RIGHTPADDING",  (0,0), (-1,-1), 6),
                        ]))
                        story.append(t)
                        story.append(Spacer(1, 6))
                    continue

                if line.startswith("### "):
                    story.append(Paragraph(inline(line[4:]), s_h3))
                elif line.startswith("## "):
                    story.append(HRFlowable(width="100%", thickness=0.5,
                        color=colors.HexColor("#c7d2fe"), spaceAfter=2))
                    story.append(Paragraph(inline(line[3:]), s_h2))
                elif line.startswith("# "):
                    story.append(Paragraph(inline(line[2:]), s_h1))
                elif line.startswith("- ") or line.startswith("* "):
                    story.append(Paragraph("• " + inline(line[2:]), s_bull))
                elif _re.match(r"^\d+\.\s", line):
                    txt = _re.sub(r"^\d+\.\s*", "", line)
                    story.append(Paragraph("• " + inline(txt), s_bull))
                elif line.strip() == "" or line.strip() == "---":
                    story.append(Spacer(1, 4))
                else:
                    story.append(Paragraph(inline(line), s_body))
                i += 1

            doc.build(story)
            return buf.getvalue()
        except Exception as e:
            return None

    _pdf_bytes = _make_pdf(
        st.session_state.comparison_result,
        st.session_state.doc_a_name,
        st.session_state.doc_b_name,
    )
    _dl_col, _ = st.columns([1, 3])
    with _dl_col:
        if _pdf_bytes:
            st.download_button(
                "⬇ Download PDF",
                data=_pdf_bytes,
                file_name=f"{_fname}.pdf",
                mime="application/pdf",
                width="stretch",
                key="dl_pdf",
            )
        else:
            st.warning("PDF generation failed — install `reportlab`")

    # ── Follow-up Chat ────────────────────────────────────────
    st.markdown("<br>", unsafe_allow_html=True)
    st.markdown(
        '<div class="section-header"><span class="dot"></span>FOLLOW-UP CHAT'
        '<span style="color:var(--muted);font-size:.7rem;font-weight:400;"> '
        '— ask Claude anything about these documents</span></div>',
        unsafe_allow_html=True,
    )

    def _build_chat_system():
        return f"""You are a senior document analyst. You have already compared two documents and produced a report. The user wants to ask follow-up questions.

DOCUMENT A ({st.session_state.doc_a_name}):
{truncate(st.session_state.doc_a_text, 15000)}

DOCUMENT B ({st.session_state.doc_b_name}):
{truncate(st.session_state.doc_b_text, 15000)}

COMPARISON REPORT ALREADY PRODUCED:
{st.session_state.comparison_result[:3000]}{"…[truncated]" if len(st.session_state.comparison_result) > 3000 else ""}

Answer follow-up questions concisely and specifically. Reference the actual document content. Use markdown where it adds clarity. Do not wrap text containing $, numbers, or punctuation inside * or ** — write figures as plain text and use double quotes for inline document quotes."""

    # Render chat history
    for _msg in st.session_state.chat_history:
        with st.chat_message(_msg["role"]):
            st.markdown(_escape_dollars(_msg["content"]))

    # Stream pending reply
    if st.session_state.chat_pending and claude_key:
        st.session_state.chat_pending = False
        _api_messages = [
            {"role": m["role"], "content": m["content"]}
            for m in st.session_state.chat_history
        ]
        try:
            _chat_client = _anthropic.Anthropic(api_key=claude_key)
            with st.chat_message("assistant"):
                _reply = ""
                _ph    = st.empty()
                with _chat_client.messages.stream(
                    model=model,
                    max_tokens=2048,
                    system=_build_chat_system(),
                    messages=_api_messages,
                ) as _stream:
                    for _delta in _stream.text_stream:
                        _reply += _delta
                        _ph.markdown(_escape_dollars(_reply) + "▌")
                _ph.markdown(_escape_dollars(_reply))
            st.session_state.chat_history.append({"role": "assistant", "content": _reply})
        except _anthropic.AuthenticationError:
            st.error("Invalid API key.")
        except _anthropic.RateLimitError:
            st.error("Rate limit reached. Wait a moment and try again.")
        except Exception as _e:
            st.error(f"Chat error: {type(_e).__name__}: {_e}")

    # Chat input
    _user_msg = st.chat_input("Ask a question about the documents…", key="chat_input")
    if _user_msg:
        st.session_state.chat_history.append({"role": "user", "content": _user_msg})
        st.session_state.chat_pending = True
        st.rerun()

    # Clear chat
    if st.session_state.chat_history:
        if st.button("Clear chat history", key="clear_chat"):
            st.session_state.chat_history = []
            st.session_state.chat_pending = False
            st.rerun()

# ─────────────────────────────────────────────────────────────
# EMPTY STATE
# ─────────────────────────────────────────────────────────────

elif not (file_a and file_b):
    st.markdown("""
<div class="empty-state">
  <div class="empty-title">UPLOAD TWO DOCUMENTS</div>
  <div class="empty-sub">
    Drop a PDF, DOCX, or XLSX into each slot above, choose your comparison focus,
    then click <strong style="color:var(--blue);">Compare Documents</strong>.
  </div>
</div>
""", unsafe_allow_html=True)

# ─────────────────────────────────────────────────────────────
# FOOTER
# ─────────────────────────────────────────────────────────────

st.markdown("""
<div class="footer">
  <div class="footer-brand">SEGA DOC COMPARATOR</div>
  <div class="footer-note">Powered by Claude · Documents processed locally · Internal use only</div>
</div>
""", unsafe_allow_html=True)